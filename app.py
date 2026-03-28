from pathlib import Path
import base64
import json
import contextlib
from io import BytesIO, StringIO
from html.parser import HTMLParser
import re
import zipfile

import altair as alt
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import requests
import streamlit as st
import yfinance as yf
from matplotlib import font_manager
from matplotlib.ticker import AutoMinorLocator
from pandas.tseries.offsets import MonthEnd
from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# =====================================
# CONFIG
# =====================================
BASE_DIR = Path(__file__).parent
DATA_FILE = BASE_DIR / "index_database.xlsx"
ALBION_LOGO_FILE = BASE_DIR / "albion_logo.png"
POWERED_BY_FILE = BASE_DIR / "Powered by SA.png"

TIME_SERIES_SHEET = "time_series"
MAPPING_SHEET = "mapping"
REGIONS_SHEET = "regions"
SECTORS_SHEET = "sectors"
FACTORS_SHEET = "factors"

BRAND_ORANGE = "#f36f21"
BRAND_ORANGE_DARK = "#d65f17"
TEXT_GREY = "#555555"
LIGHT_GREY = "#f3f3f3"
MID_GREY = "#d9d9d9"
CHART_BG_GREY = "#f5f5f5"
WHITE = "#ffffff"
APP_TITLE = "smarteranalytics™ Hub"
PAGE_LABELS = {
    "Dashboard": "Market snapshot",
    "Charts": "Performance analysis",
    "Risk": "Risk analysis",
    "Factors": "Risk factors",
    "Geo": "Geographic analysis",
    "Sector": "Sector analysis",
    "Yield": "Rates & FX",
}
DISPLAY_NAME_OVERRIDES = {
    "Global stocks": "Global market",
    "UK stocks": "UK market",
    "Developed stocks": "Developed market",
    "Emerging stocks": "Emerging market",
    "UK value stocks": "UK value",
    "UK small stocks": "UK small",
    "Developed value stocks": "Developed value",
    "Developed small stocks": "Developed small",
    "Emerging value stocks": "Emerging value",
    "Emerging small stocks": "Emerging small",
    "Developed REITs": "REITs",
    "Global GBP hedged bonds (0-5)": "Global bonds (0-5, GBP)",
}

DASHBOARD_HORIZONS = {
    "20Y": "20 Year",
    "10Y": "10 Year",
    "5Y": "5 Year",
    "YTD": "YTD",
}

CHART_PERIODS = {
    "YTD": "YTD",
    "1Y": "1 Year",
    "3Y": "3 Year",
    "5Y": "5 Year",
    "10Y": "10 Year",
    "20Y": "20 Year",
}

RISK_PERIODS = {
    "1Y": "1 Year",
    "3Y": "3 Year",
    "5Y": "5 Year",
    "10Y": "10 Year",
    "20Y": "20 Year",
    "Max": "Maximum",
}

RETURNS_TABLE_PERIODS = ["YTD", "1Y", "3Y", "5Y", "7Y", "10Y", "15Y", "20Y", "25Y", "Period"]

DISPLAY_GROUPS_ABSOLUTE = [
    {
        "title": "Global Equity",
        "items": ["Global stocks"],
        "labels": {"Global stocks": ""},
    },
    {
        "title": "UK Equity",
        "items": ["UK stocks", "UK value stocks", "UK small stocks"],
        "labels": {
            "UK stocks": "Market",
            "UK value stocks": "Value",
            "UK small stocks": "Small",
        },
    },
    {
        "title": "Developed Equity",
        "items": ["Developed stocks", "Developed value stocks", "Developed small stocks"],
        "labels": {
            "Developed stocks": "Market",
            "Developed value stocks": "Value",
            "Developed small stocks": "Small",
        },
    },
    {
        "title": "Emerging Equity",
        "items": ["Emerging stocks", "Emerging value stocks", "Emerging small stocks"],
        "labels": {
            "Emerging stocks": "Market",
            "Emerging value stocks": "Value",
            "Emerging small stocks": "Small",
        },
    },
    {
        "title": "REITs",
        "items": ["Developed REITs"],
        "labels": {"Developed REITs": ""},
    },
    {
        "title": "",
        "items": [
            "Cash (GBP)",
            "UK Gilts (0-5)",
            "UK IL Gilts (0-5)",
            "Global GBP hedged bonds (0-5)",
        ],
        "labels": {
            "Cash (GBP)": "Cash",
            "UK Gilts (0-5)": "UK Gilts (0-5)",
            "UK IL Gilts (0-5)": "UK IL Gilts (0-5)",
            "Global GBP hedged bonds (0-5)": "Global bonds (0-5, GBP)",
        },
    },
]

DISPLAY_GROUPS_RELATIVE_MINOR = [
    {
        "title": "Relative to DM",
        "items": [
            "UK stocks",
            "Emerging stocks",
            "Developed REITs",
            "Developed value stocks",
            "Developed small stocks",
        ],
        "labels": {
            "UK stocks": "UK",
            "Emerging stocks": "EM",
            "Developed REITs": "REIT",
            "Developed value stocks": "DM Value",
            "Developed small stocks": "DM Small",
        },
    },
    {
        "title": "Relative to EM",
        "items": ["Emerging value stocks", "Emerging small stocks"],
        "labels": {
            "Emerging value stocks": "EM Value",
            "Emerging small stocks": "EM Small",
        },
    },
    {
        "title": "Relative to UK",
        "items": ["UK value stocks", "UK small stocks"],
        "labels": {
            "UK value stocks": "UK Value",
            "UK small stocks": "UK Small",
        },
    },
    {
        "title": "",
        "items": [
            "Cash (GBP)",
            "UK Gilts (0-5)",
            "UK IL Gilts (0-5)",
            "Global GBP hedged bonds (0-5)",
        ],
        "labels": {
            "Cash (GBP)": "Cash",
            "UK Gilts (0-5)": "UK Gilts (0-5)",
            "UK IL Gilts (0-5)": "UK IL Gilts (0-5)",
            "Global GBP hedged bonds (0-5)": "Global bonds (0-5, GBP)",
        },
    },
]

REPORT_DISPLAY_GROUPS_ABSOLUTE = DISPLAY_GROUPS_ABSOLUTE[1:]
REPORT_LABEL_OVERRIDES = {
    "Cash": "Cash",
    "UK Gilts (0-5)": "UK Gilts",
    "UK IL Gilts (0-5)": "UK ILG (0-5)",
    "Global bonds (0-5, GBP)": "GSDB (0-5)",
}

ASSET_CLASS_ALIASES = {
    "Global equity": "Global stocks",
    "Global stocks": "Global stocks",
    "Global market": "Global stocks",
    "World equity": "Global stocks",
    "World stocks": "Global stocks",
    "UK equity": "UK stocks",
    "UK stocks": "UK stocks",
    "UK market": "UK stocks",
    "UK value": "UK value stocks",
    "UK value stocks": "UK value stocks",
    "UK small": "UK small stocks",
    "UK small stocks": "UK small stocks",
    "Developed equity": "Developed stocks",
    "Developed stocks": "Developed stocks",
    "Developed market": "Developed stocks",
    "Developed value": "Developed value stocks",
    "Developed value stocks": "Developed value stocks",
    "Developed small": "Developed small stocks",
    "Developed small stocks": "Developed small stocks",
    "Emerging market equity": "Emerging stocks",
    "Emerging equity": "Emerging stocks",
    "Emerging stocks": "Emerging stocks",
    "Emerging market": "Emerging stocks",
    "Emerging value": "Emerging value stocks",
    "Emerging value stocks": "Emerging value stocks",
    "Emerging small": "Emerging small stocks",
    "Emerging small stocks": "Emerging small stocks",
    "Developed REITs": "Developed REITs",
    "REITs": "Developed REITs",
    "Cash GBP": "Cash (GBP)",
    "Cash (GBP)": "Cash (GBP)",
    "UK Gilts": "UK Gilts (0-5)",
    "Short Gilt": "UK Gilts (0-5)",
    "UK Gilts (0-5)": "UK Gilts (0-5)",
    "UK IL Gilts": "UK IL Gilts (0-5)",
    "Short IL Gilt": "UK IL Gilts (0-5)",
    "UK IL Gilts (0-5)": "UK IL Gilts (0-5)",
    "Global bonds (0-5, GBP)": "Global GBP hedged bonds (0-5)",
    "Global GBP Hedged bonds (0-5)": "Global GBP hedged bonds (0-5)",
    "Global GBP hedged bonds (0-5)": "Global GBP hedged bonds (0-5)",
    "Global Short Bond GBP": "Global GBP hedged bonds (0-5)",
    "GSDB (GBP)": "Global GBP hedged bonds (0-5)",
}

ETF_DOWNLOAD_START = "2000-01-01"

MAJOR_GROWTH_BASE_MAP = {
    "Global stocks": None,
    "UK stocks": "Global stocks",
    "UK value stocks": "Global stocks",
    "UK small stocks": "Global stocks",
    "Developed stocks": "Global stocks",
    "Developed value stocks": "Global stocks",
    "Developed small stocks": "Global stocks",
    "Emerging stocks": "Global stocks",
    "Emerging value stocks": "Global stocks",
    "Emerging small stocks": "Global stocks",
    "Developed REITs": "Global stocks",
}

MINOR_GROWTH_BASE_MAP = {
    "Global stocks": None,
    "UK stocks": "Developed stocks",
    "UK value stocks": "UK stocks",
    "UK small stocks": "UK stocks",
    "Developed stocks": None,
    "Developed value stocks": "Developed stocks",
    "Developed small stocks": "Developed stocks",
    "Emerging stocks": "Developed stocks",
    "Emerging value stocks": "Emerging stocks",
    "Emerging small stocks": "Emerging stocks",
    "Developed REITs": "Developed stocks",
}

DEFENSIVE_BASE_MAP = {
    "Cash (GBP)": None,
    "UK Gilts (0-5)": "Cash (GBP)",
    "UK IL Gilts (0-5)": "Cash (GBP)",
    "Global GBP hedged bonds (0-5)": "Cash (GBP)",
}

INFLATION_SERIES_ALIASES = [
    "UK inflation",
    "UK CPIH",
    "CPIH",
    "UK CPI",
    "CPI",
    "Inflation",
    "UK RPI",
    "RPI",
]

ONS_CPI_INDEX_CSV_URL = (
    "https://www.ons.gov.uk/generator?format=csv&uri="
    "/economy/inflationandpriceindices/timeseries/d7bt/mm23"
)
BOE_YIELD_CURVE_ZIP_URL = "https://www.bankofengland.co.uk/-/media/boe/files/statistics/yield-curves/latest-yield-curve-data.zip"
BOE_NOMINAL_MONTH_END_ZIP_URL = "https://www.bankofengland.co.uk/-/media/boe/files/statistics/yield-curves/glcnominalmonthedata.zip"
BOE_REAL_MONTH_END_ZIP_URL = "https://www.bankofengland.co.uk/-/media/boe/files/statistics/yield-curves/glcrealmonthedata.zip"
DIVIDENDDATA_INDEX_LINKED_GILTS_URL = "https://www.dividenddata.co.uk/index-linked-gilts-prices-yields.py"
DIVIDENDDATA_GILT_DETAIL_URL = "https://www.dividenddata.co.uk/gilts.py?ticker={ticker}"
DIVIDENDDATA_SHORT_END_TICKERS = ["TR26", "T27", "T28", "T29", "TR31"]
WORLD_GOVERNMENT_BONDS_BASE_URL = "https://www.worldgovernmentbonds.com"
WORLD_GOVERNMENT_BONDS_COUNTRIES = [
    ("United States", "/country/united-states/"),
    ("China", "/country/china/"),
    ("Japan", "/country/japan/"),
    ("France", "/country/france/"),
    ("United Kingdom", "/country/united-kingdom/"),
    ("Germany", "/country/germany/"),
    ("Canada", "/country/canada/"),
    ("Italy", "/country/italy/"),
    ("Spain", "/country/spain/"),
    ("Netherlands", "/country/netherlands/"),
]
FX_DOWNLOAD_START = "1996-01-01"
GEO_MAX_START = pd.Timestamp("1996-03-18")
GEO_NEUTRAL_CURRENCIES = ["USD", "GBP", "AUD", "EUR", "NOK"]
MSCI_ACWI_IMI_REGION_WEIGHTS = {
    "North America": 65.9,
    "Pacific": 8.0,
    "Europe": 14.9,
    "Emerging": 11.3,
}
PATCHWORK_COUNTRY_SETS = {
    "All countries": None,
    "Largest 20": "largest_20",
    "Largest 10": "largest_10",
    "Regional": ["North America", "Pacific", "Europe", "Emerging"],
}
COMPANIESMARKETCAP_COUNTRIES_URL = "https://companiesmarketcap.com/gbp/all-countries/"
MSCI_ACWI_IMI_REGION_SOURCE_URL = "https://www.msci.com/research-and-insights/video/acwi-imi-complete-geographic-breakdown"
MSCI_ACWI_IMI_REGION_SOURCE_DATE = pd.Timestamp("2025-12-31")
FX_CURRENCIES = ["USD", "EUR", "GBP", "JPY", "CHF", "AUD", "NZD", "CAD", "SEK", "NOK"]
FX_TICKER_SPECS = {
    "EUR": ("EURUSD=X", "direct"),
    "GBP": ("GBPUSD=X", "direct"),
    "JPY": ("JPY=X", "inverse"),
    "CHF": ("CHF=X", "inverse"),
    "AUD": ("AUDUSD=X", "direct"),
    "NZD": ("NZDUSD=X", "direct"),
    "CAD": ("CAD=X", "inverse"),
    "SEK": ("SEK=X", "inverse"),
    "NOK": ("NOK=X", "inverse"),
}
FX_PERIODS = {
    "YTD": "YTD",
    "1Y": "1Y",
    "3Y": "3Y",
    "5Y": "5Y",
    "10Y": "10Y",
    "20Y": "20Y",
    "MAX": "MAX",
}
UK_HISTORICAL_YIELD_PERIODS = ["YTD", "1Y", "3Y", "5Y", "10Y", "MAX"]

DEFAULT_ASSET_ORDER = [
    "Global stocks",
    "UK stocks",
    "UK value stocks",
    "UK small stocks",
    "Developed stocks",
    "Developed value stocks",
    "Developed small stocks",
    "Emerging stocks",
    "Emerging value stocks",
    "Emerging small stocks",
    "Developed REITs",
    "Cash (GBP)",
    "UK Gilts (0-5)",
    "UK IL Gilts (0-5)",
    "Global GBP hedged bonds (0-5)",
]

DEFAULT_CHART_ASSETS = [
    "Cash (GBP)",
    "UK stocks",
    "Developed stocks",
    "Developed value stocks",
    "Developed small stocks",
    "Emerging stocks",
    "Developed REITs",
    "Global GBP hedged bonds (0-5)",
    "UK IL Gilts (0-5)",
]

ASSET_COLOURS = {
    "Global stocks": "#c95b2b",
    "UK stocks": "#d71921",
    "UK value stocks": "#ef5350",
    "UK small stocks": "#8f1015",
    "Developed stocks": "#f36f21",
    "Developed value stocks": "#ff9a4d",
    "Developed small stocks": "#c95f18",
    "Emerging stocks": "#f4c542",
    "Emerging value stocks": "#ffd965",
    "Emerging small stocks": "#cfa400",
    "Developed REITs": "#2e8b57",
    "Cash (GBP)": "#4b5563",
    "UK Gilts (0-5)": "#1f77b4",
    "UK IL Gilts (0-5)": "#5dade2",
    "Global GBP hedged bonds (0-5)": "#0b5394",
}


# =====================================
# HELPERS
# =====================================
@st.cache_data(show_spinner=False)
def img_to_base64(path: str, file_mtime: float) -> str:
    return base64.b64encode(Path(path).read_bytes()).decode("utf-8")


def standardise_series(series: pd.Series) -> pd.Series:
    s = pd.to_numeric(series, errors="coerce")
    max_abs = s.dropna().abs().max()
    if pd.notna(max_abs) and max_abs > 1.5:
        s = s / 100.0
    return s


def normalise_ticker(value: object) -> str:
    if pd.isna(value):
        return ""
    ticker = str(value).strip().upper()
    return "" if ticker.lower() in {"", "nan", "none"} else ticker


def normalise_name(value: object) -> str:
    return str(value).strip().lower()


def display_name(value: object) -> str:
    text = str(value).strip()
    return DISPLAY_NAME_OVERRIDES.get(text, text)


def format_pct(x: float) -> str:
    return "-" if pd.isna(x) else f"{x:.1%}"


def annualised_return_from_growth(growth: float, years: float) -> float:
    if pd.isna(growth) or growth <= 0 or years <= 0:
        return np.nan
    return growth ** (1 / years) - 1


def build_lookup_table(returns_df: pd.DataFrame) -> dict:
    return {} if returns_df.empty else returns_df.set_index("asset_class").to_dict(orient="index")


def heat_colour(value: float, vmin: float, vmax: float) -> str:
    if pd.isna(value):
        return "#E9E9E9"

    if value == 0:
        return "#f6f6f6"

    if value > 0:
        positive_max = max(float(vmax), 0.0)
        if positive_max <= 0:
            return "#f6f6f6"
        norm = min(max(value / positive_max, 0), 1)
        light = np.array([130, 130, 130])
        dark = np.array([0, 170, 95])
        rgb = (light * (1 - norm) + dark * norm).astype(int)
        return f"rgb({rgb[0]},{rgb[1]},{rgb[2]})"

    negative_min = min(float(vmin), 0.0)
    if negative_min >= 0:
        return "#f6f6f6"
    norm = min(max(abs(value) / abs(negative_min), 0), 1)
    light = np.array([130, 130, 130])
    dark = np.array([190, 30, 55])
    rgb = (light * (1 - norm) + dark * norm).astype(int)
    return f"rgb({rgb[0]},{rgb[1]},{rgb[2]})"


def text_colour_for_background(background: str) -> str:
    if not isinstance(background, str):
        return "#111111"

    rgb_match = re.match(r"rgb\((\d+),(\d+),(\d+)\)", background.replace(" ", ""))
    if rgb_match:
        r, g, b = [int(rgb_match.group(i)) for i in range(1, 4)]
    elif background.startswith("#") and len(background) == 7:
        r = int(background[1:3], 16)
        g = int(background[3:5], 16)
        b = int(background[5:7], 16)
    else:
        return "#111111"

    luminance = (0.299 * r) + (0.587 * g) + (0.114 * b)
    return "#111111" if luminance >= 165 else "#ffffff"


def get_series_map_bounds(series_map: dict[str, pd.Series]) -> tuple[pd.Timestamp | None, pd.Timestamp | None]:
    starts = []
    ends = []
    for series in series_map.values():
        s = series.dropna().sort_index()
        if s.empty:
            continue
        starts.append(pd.Timestamp(s.index.min()))
        ends.append(pd.Timestamp(s.index.max()))
    if not starts or not ends:
        return None, None
    return min(starts), max(ends)


def filter_series_by_date(series: pd.Series, start_date: pd.Timestamp | None, end_date: pd.Timestamp | None) -> pd.Series:
    s = series.dropna().sort_index()
    if start_date is not None:
        s = s[s.index >= pd.Timestamp(start_date)]
    if end_date is not None:
        s = s[s.index <= pd.Timestamp(end_date)]
    return s


def filter_series_map_by_date(
    series_map: dict[str, pd.Series],
    start_date: pd.Timestamp | None,
    end_date: pd.Timestamp | None,
) -> dict[str, pd.Series]:
    return {asset: filter_series_by_date(series, start_date, end_date) for asset, series in series_map.items()}


def rank_heat_colour(value: float, vmin: float, vmax: float, low_is_good: bool = False) -> str:
    if pd.isna(value):
        return "#E9E9E9"
    if vmax <= vmin:
        return "#f6f6f6"

    norm = (float(value) - float(vmin)) / (float(vmax) - float(vmin))
    norm = min(max(norm, 0), 1)
    if low_is_good:
        norm = 1 - norm

    start = np.array([190, 30, 55])
    end = np.array([0, 170, 95])
    rgb = (start * (1 - norm) + end * norm).astype(int)
    return f"rgb({rgb[0]},{rgb[1]},{rgb[2]})"


def get_period_start_anchor(end_date: pd.Timestamp, period_key: str) -> pd.Timestamp:
    end_date = pd.Timestamp(end_date)
    if period_key == "YTD":
        return pd.Timestamp(end_date.year - 1, 12, 31)
    if period_key == "MAX":
        return GEO_MAX_START
    years_map = {"1Y": 1, "3Y": 3, "5Y": 5, "10Y": 10, "20Y": 20}
    years = years_map.get(period_key)
    return end_date - pd.DateOffset(years=years) if years is not None else end_date


def get_fx_max_start_anchor(fx_values_df: pd.DataFrame) -> pd.Timestamp:
    if fx_values_df.empty:
        return pd.Timestamp("2000-01-01")

    start_dates = []
    for currency in FX_CURRENCIES:
        series = fx_values_df.get(currency, pd.Series(dtype=float)).dropna().sort_index()
        if series.empty:
            continue
        start_dates.append(pd.Timestamp(series.index.min()).normalize())

    if not start_dates:
        return pd.Timestamp("2000-01-01")

    return max(start_dates)


def get_fx_period_start_anchor(end_date: pd.Timestamp, period_key: str, fx_values_df: pd.DataFrame) -> pd.Timestamp:
    if period_key == "MAX":
        return get_fx_max_start_anchor(fx_values_df)
    return get_period_start_anchor(end_date, period_key)


def get_fx_period_options(end_date: pd.Timestamp | None, fx_values_df: pd.DataFrame) -> list[str]:
    if end_date is None:
        return ["YTD", "1Y", "3Y", "5Y", "10Y", "MAX"]

    options = ["YTD", "1Y", "3Y", "5Y", "10Y"]
    max_anchor = get_fx_max_start_anchor(fx_values_df)
    if max_anchor <= pd.Timestamp(end_date) - pd.DateOffset(years=20):
        options.append("20Y")
    options.append("MAX")
    return options


def get_geo_max_start_anchor(fx_values_df: pd.DataFrame, neutral_currency: str) -> pd.Timestamp:
    anchor = GEO_MAX_START
    if neutral_currency == "USD":
        return anchor

    neutral_series = fx_values_df.get(neutral_currency, pd.Series(dtype=float)).dropna().sort_index()
    if neutral_series.empty:
        return anchor
    return max(anchor, pd.Timestamp(neutral_series.index.min()).normalize())


def get_geo_period_start_anchor(
    end_date: pd.Timestamp,
    period_key: str,
    fx_values_df: pd.DataFrame,
    neutral_currency: str,
) -> pd.Timestamp:
    if period_key == "MAX":
        return get_geo_max_start_anchor(fx_values_df, neutral_currency)
    return get_period_start_anchor(end_date, period_key)


def get_geo_period_options(end_date: pd.Timestamp | None, fx_values_df: pd.DataFrame, neutral_currency: str) -> list[str]:
    if end_date is None:
        return list(FX_PERIODS.keys())

    options = ["YTD", "1Y", "3Y", "5Y", "10Y"]
    max_anchor = get_geo_max_start_anchor(fx_values_df, neutral_currency)
    if max_anchor <= pd.Timestamp(end_date) - pd.DateOffset(years=20):
        options.append("20Y")
    options.append("MAX")
    return options


def get_common_series_inception_anchor(
    series_map: dict[str, pd.Series],
    include_labels: set[str] | None = None,
) -> pd.Timestamp | None:
    start_dates: list[pd.Timestamp] = []
    for label, series in series_map.items():
        if include_labels is not None and str(label) not in include_labels:
            continue
        s = series.dropna().sort_index()
        if s.empty:
            continue
        start_dates.append(pd.Timestamp(s.index.min()).normalize())
    if not start_dates:
        return None
    return max(start_dates)


@st.cache_data(show_spinner=False, ttl=43200)
def fetch_fx_value_series() -> pd.DataFrame:
    tickers = tuple(sorted(spec[0] for spec in FX_TICKER_SPECS.values()))
    raw_prices = fetch_yf_prices(tickers, FX_DOWNLOAD_START)
    if raw_prices.empty:
        return pd.DataFrame()

    fx_frames = []
    for currency, (ticker, mode) in FX_TICKER_SPECS.items():
        if ticker not in raw_prices.columns:
            continue
        series = pd.to_numeric(raw_prices[ticker], errors="coerce").dropna().sort_index()
        if series.empty:
            continue
        if mode == "inverse":
            series = (1.0 / series).replace([np.inf, -np.inf], np.nan).dropna()
        series.name = currency
        fx_frames.append(series)

    if not fx_frames:
        return pd.DataFrame()

    fx_df = pd.concat(fx_frames, axis=1).sort_index().ffill()
    fx_df["USD"] = 1.0
    fx_df = fx_df.reindex(columns=FX_CURRENCIES)
    fx_df["USD"] = 1.0
    return fx_df.sort_index().ffill()


def build_currency_performance_matrix(fx_values_df: pd.DataFrame, period_key: str) -> tuple[pd.DataFrame, pd.Timestamp | None]:
    if fx_values_df.empty:
        return pd.DataFrame(), None

    last_dates = [
        pd.Timestamp(fx_values_df[c].dropna().index.max())
        for c in fx_values_df.columns
        if c in fx_values_df.columns and not fx_values_df[c].dropna().empty
    ]
    if not last_dates:
        return pd.DataFrame(), None

    end_date = min(last_dates)
    start_anchor = get_fx_period_start_anchor(end_date, period_key, fx_values_df)
    rows = []

    for row_currency in FX_CURRENCIES:
        row = {"currency": row_currency}
        row_series = fx_values_df.get(row_currency, pd.Series(dtype=float)).dropna().sort_index()
        for col_currency in FX_CURRENCIES:
            if row_currency == col_currency:
                row[col_currency] = np.nan
                continue
            col_series = fx_values_df.get(col_currency, pd.Series(dtype=float)).dropna().sort_index()
            if row_series.empty or col_series.empty:
                row[col_currency] = np.nan
                continue
            cross_df = pd.concat([row_series.rename("row"), col_series.rename("col")], axis=1).sort_index().ffill().dropna()
            if cross_df.empty:
                row[col_currency] = np.nan
                continue
            cross_series = (cross_df["row"] / cross_df["col"]).dropna()
            if period_key == "MAX":
                row[col_currency] = calc_period_return(cross_series, end_date, "Period", whole_period_start=start_anchor)
            else:
                row[col_currency] = calc_period_return(cross_series, end_date, period_key, whole_period_start=start_anchor)
        rows.append(row)

    return pd.DataFrame(rows), end_date


def build_country_performance_df(
    regions_df: pd.DataFrame,
    fx_values_df: pd.DataFrame,
    period_key: str,
    neutral_currency: str = "GBP",
    preferred_series_map: dict[str, pd.Series] | None = None,
) -> tuple[pd.DataFrame, pd.Timestamp | None]:
    if regions_df.empty:
        return pd.DataFrame(), None

    countries = regions_df[
        (regions_df["country_flag"] == 1) & regions_df["available"].fillna(False)
    ].copy()
    if countries.empty:
        return pd.DataFrame(), None

    tickers = tuple(sorted({normalise_ticker(t) for t in countries["ticker"].tolist() if normalise_ticker(t)}))
    # Geographic country analysis needs pre-2000 ETF history for the MAX view.
    prices_df = fetch_yf_prices(tickers, FX_DOWNLOAD_START)
    if prices_df.empty:
        return pd.DataFrame(), None

    quote_currency_map = fetch_yf_quote_currencies(tickers)

    last_dates = [
        pd.Timestamp(prices_df[c].dropna().index.max())
        for c in prices_df.columns
        if not prices_df[c].dropna().empty
    ]
    if not last_dates:
        return pd.DataFrame(), None

    end_date = min(last_dates)
    neutral_series = (
        fx_values_df.get(neutral_currency, pd.Series(dtype=float)).dropna().sort_index()
        if neutral_currency
        else pd.Series(dtype=float)
    )
    start_anchor = get_geo_period_start_anchor(end_date, period_key, fx_values_df, neutral_currency)
    rows = []

    for row in countries.itertuples():
        ticker = normalise_ticker(row.ticker)
        series = get_price_series(prices_df, ticker)
        if series.empty:
            continue
        series = series[series.index <= end_date]
        if series.empty:
            continue

        quote_currency = quote_currency_map.get(ticker, "")
        converted_series = series.copy()
        if neutral_currency == "USD":
            converted_series = series.copy()
            quote_currency = quote_currency or "USD"
        elif neutral_currency and quote_currency and quote_currency != neutral_currency:
            quote_fx_series = fx_values_df.get(quote_currency, pd.Series(dtype=float)).dropna().sort_index()
            if quote_fx_series.empty or neutral_series.empty:
                continue
            aligned = pd.concat(
                [
                    series.rename("price"),
                    quote_fx_series.rename("quote_fx"),
                    neutral_series.rename("neutral_fx"),
                ],
                axis=1,
            ).sort_index().ffill().dropna()
            if aligned.empty:
                continue
            converted_series = (aligned["price"] * aligned["quote_fx"] / aligned["neutral_fx"]).dropna()
        elif neutral_currency and quote_currency == neutral_currency:
            converted_series = series.copy()
        elif neutral_currency and not quote_currency:
            continue

        if period_key == "MAX":
            performance = calc_period_return(
                converted_series,
                end_date,
                "Period",
                whole_period_start=start_anchor,
            )
        else:
            performance = calc_period_return(
                converted_series,
                end_date,
                period_key,
                whole_period_start=start_anchor,
            )
        if pd.isna(performance):
            continue

        rows.append(
            {
                "ticker": ticker,
                "country": row.investment_area,
                "name": row.name,
                "quote_currency": quote_currency or "-",
                "neutral_currency": neutral_currency or "-",
                "performance": performance,
            }
        )

    if not rows:
        out = pd.DataFrame(columns=["ticker", "country", "name", "quote_currency", "neutral_currency", "performance"])
    else:
        out = pd.DataFrame(rows)

    if preferred_series_map is not None and "Global stocks" in preferred_series_map:
        global_series = preferred_series_map.get("Global stocks", pd.Series(dtype=float)).dropna().sort_index()
        global_series = global_series[global_series.index <= end_date]
        if not global_series.empty:
            if period_key == "MAX":
                global_return = calc_period_return(global_series, end_date, "Period", whole_period_start=start_anchor)
            else:
                global_return = calc_period_return(global_series, end_date, period_key, whole_period_start=start_anchor)
            if pd.notna(global_return):
                global_df = pd.DataFrame(
                    [
                        {
                            "ticker": "Global market",
                            "country": "Global market",
                            "name": "Global market",
                            "quote_currency": neutral_currency,
                            "neutral_currency": neutral_currency,
                            "performance": global_return,
                        }
                    ]
                )
                if out.empty or out.dropna(how="all").empty:
                    out = global_df
                else:
                    out = pd.concat([global_df, out], ignore_index=True)

    if out.empty:
        return pd.DataFrame(), end_date

    out = out.sort_values(["performance", "country"], ascending=[False, True]).reset_index(drop=True)
    return out, end_date


def build_country_series_map(
    regions_df: pd.DataFrame,
    fx_values_df: pd.DataFrame,
    neutral_currency: str = "USD",
    preferred_series_map: dict[str, pd.Series] | None = None,
) -> tuple[dict[str, pd.Series], pd.Timestamp | None]:
    if regions_df.empty:
        return {}, None

    countries = regions_df[
        (regions_df["country_flag"] == 1) & regions_df["available"].fillna(False)
    ].copy()
    if countries.empty:
        return {}, None

    tickers = tuple(sorted({normalise_ticker(t) for t in countries["ticker"].tolist() if normalise_ticker(t)}))
    prices_df = fetch_yf_prices(tickers, FX_DOWNLOAD_START)
    if prices_df.empty:
        return {}, None

    quote_currency_map = fetch_yf_quote_currencies(tickers)
    last_dates = [
        pd.Timestamp(prices_df[c].dropna().index.max())
        for c in prices_df.columns
        if not prices_df[c].dropna().empty
    ]
    if not last_dates:
        return {}, None

    end_date = min(last_dates)
    out: dict[str, pd.Series] = {}

    for row in countries.itertuples():
        ticker = normalise_ticker(row.ticker)
        series = get_price_series(prices_df, ticker)
        if series.empty:
            continue
        series = series[series.index <= end_date]
        if series.empty:
            continue

        quote_currency = quote_currency_map.get(ticker, "") or ("USD" if neutral_currency == "USD" else "")
        converted_series = (
            series.copy()
            if neutral_currency == "USD"
            else convert_price_series_to_neutral_currency(series, quote_currency, neutral_currency, fx_values_df)
        )
        if converted_series.empty:
            continue
        out[str(row.investment_area).strip()] = converted_series.dropna().sort_index()

    if preferred_series_map is not None and "Global stocks" in preferred_series_map:
        global_series = preferred_series_map.get("Global stocks", pd.Series(dtype=float)).dropna().sort_index()
        global_series = global_series[global_series.index <= end_date]
        if not global_series.empty:
            out["Global market"] = global_series

    return out, end_date


def build_labelled_series_map(
    source_df: pd.DataFrame,
    label_col: str,
    fx_values_df: pd.DataFrame,
    neutral_currency: str = "USD",
    preferred_series_map: dict[str, pd.Series] | None = None,
) -> tuple[dict[str, pd.Series], pd.Timestamp | None]:
    if source_df.empty:
        return {}, None

    working = source_df[source_df["available"].fillna(False)].copy()
    if working.empty:
        return {}, None

    tickers = tuple(sorted({normalise_ticker(t) for t in working["ticker"].tolist() if normalise_ticker(t)}))
    prices_df = fetch_yf_prices(tickers, FX_DOWNLOAD_START)
    if prices_df.empty:
        return {}, None

    quote_currency_map = fetch_yf_quote_currencies(tickers)
    last_dates = [
        pd.Timestamp(prices_df[c].dropna().index.max())
        for c in prices_df.columns
        if not prices_df[c].dropna().empty
    ]
    if not last_dates:
        return {}, None

    end_date = min(last_dates)
    out: dict[str, pd.Series] = {}

    for row in working.itertuples():
        ticker = normalise_ticker(row.ticker)
        series = get_price_series(prices_df, ticker)
        if series.empty:
            continue
        series = series[series.index <= end_date]
        if series.empty:
            continue

        quote_currency = quote_currency_map.get(ticker, "") or ("USD" if neutral_currency == "USD" else "")
        converted_series = (
            series.copy()
            if neutral_currency == "USD"
            else convert_price_series_to_neutral_currency(series, quote_currency, neutral_currency, fx_values_df)
        )
        if converted_series.empty:
            continue
        out[str(getattr(row, label_col)).strip()] = converted_series.dropna().sort_index()

    if preferred_series_map is not None and "Global stocks" in preferred_series_map:
        global_series = preferred_series_map.get("Global stocks", pd.Series(dtype=float)).dropna().sort_index()
        global_series = global_series[global_series.index <= end_date]
        if not global_series.empty:
            out["Global market"] = global_series

    return out, end_date


def build_labelled_performance_df(
    source_df: pd.DataFrame,
    label_col: str,
    fx_values_df: pd.DataFrame,
    period_key: str,
    neutral_currency: str = "USD",
    preferred_series_map: dict[str, pd.Series] | None = None,
    max_start_anchor: pd.Timestamp | None = None,
) -> tuple[pd.DataFrame, pd.Timestamp | None]:
    series_map, end_date = build_labelled_series_map(
        source_df,
        label_col,
        fx_values_df,
        neutral_currency,
        preferred_series_map=preferred_series_map,
    )
    if not series_map or end_date is None:
        return pd.DataFrame(), None

    start_anchor = (
        pd.Timestamp(max_start_anchor).normalize()
        if period_key == "MAX" and max_start_anchor is not None
        else get_geo_period_start_anchor(end_date, period_key, fx_values_df, neutral_currency)
    )
    rows = []
    for label, series in series_map.items():
        performance = calc_period_return(
            series,
            end_date,
            "Period" if period_key == "MAX" else period_key,
            whole_period_start=start_anchor,
        )
        if pd.isna(performance):
            continue
        rows.append({"label": label, "ticker": "", "performance": performance})

    if not rows:
        return pd.DataFrame(), end_date

    out = pd.DataFrame(rows).sort_values(["performance", "label"], ascending=[False, True]).reset_index(drop=True)
    return out, end_date


@st.cache_data(show_spinner=False)
def build_factor_style_box_df(
    factors_df: pd.DataFrame,
    fx_values_df: pd.DataFrame,
    period_key: str,
    neutral_currency: str = "USD",
    region_key: str = "US",
) -> tuple[pd.DataFrame, pd.Timestamp | None, pd.Timestamp | None, pd.Timestamp | None]:
    empty = pd.DataFrame(columns=["label", "name", "ticker", "size_style", "value_style", "performance"])
    if factors_df.empty:
        return empty, None, None, None

    working = factors_df[
        factors_df["available"].fillna(False)
        & factors_df["region"].map(canonical_factor_region).eq(canonical_factor_region(region_key))
    ].copy()
    if working.empty:
        return empty, None, None, None

    working["size_style"] = working["size_style"].map(canonical_factor_size_style)
    working["value_style"] = working["value_style"].map(canonical_factor_value_style)
    working["label"] = working["label"].astype(str).str.strip()
    working = working[
        working["size_style"].isin(["Large", "Mid", "Small"]) & working["value_style"].isin(["Value", "Core", "Growth"])
    ].copy()
    if working.empty:
        return empty, None, None, None

    series_map, end_date = build_labelled_series_map(
        source_df=working,
        label_col="label",
        fx_values_df=fx_values_df,
        neutral_currency=neutral_currency,
    )
    if not series_map or end_date is None:
        return empty, None, None, None

    common_start = get_common_series_inception_anchor(series_map, include_labels=set(working["label"].tolist()))
    if common_start is None:
        return empty, None, end_date, None

    start_anchor = (
        pd.Timestamp(common_start).normalize()
        if period_key == "MAX"
        else get_geo_period_start_anchor(end_date, period_key, fx_values_df, neutral_currency)
    )

    rows = []
    for row in working.itertuples():
        label = str(row.label).strip()
        series = series_map.get(label, pd.Series(dtype=float)).dropna().sort_index()
        if series.empty:
            continue
        performance = calc_period_return(
            series,
            end_date,
            "Period" if period_key == "MAX" else period_key,
            whole_period_start=start_anchor,
        )
        if pd.isna(performance):
            continue
        rows.append(
            {
                "label": label,
                "name": str(row.name).strip(),
                "ticker": normalise_ticker(row.ticker),
                "size_style": str(row.size_style).strip(),
                "value_style": str(row.value_style).strip(),
                "performance": performance,
            }
        )

    if not rows:
        return empty, start_anchor, end_date, common_start

    out = pd.DataFrame(rows).drop_duplicates(subset=["size_style", "value_style"], keep="first")
    return out, start_anchor, end_date, common_start


def canonical_region_name(name: object) -> str:
    raw = str(name).strip()
    normalized = re.sub(r"[^a-z]+", "", raw.lower())
    aliases = {
        "northamerica": "North America",
        "americas": "North America",
        "pacific": "Pacific",
        "asiapacific": "Pacific",
        "apac": "Pacific",
        "europe": "Europe",
        "emea": "Europe",
        "europemiddleeastandafrica": "Europe",
        "emerging": "Emerging",
        "emergingmarkets": "Emerging",
        "em": "Emerging",
    }
    return aliases.get(normalized, raw)


def canonical_factor_region(value: object) -> str:
    raw = str(value).strip()
    normalized = re.sub(r"[^a-z]+", "", raw.lower())
    aliases = {
        "us": "US",
        "usa": "US",
        "unitedstates": "US",
        "unitedstatesofamerica": "US",
    }
    return aliases.get(normalized, raw.upper() if len(raw) <= 3 else raw)


def canonical_factor_size_style(value: object) -> str:
    raw = str(value).strip()
    normalized = re.sub(r"[^a-z]+", "", raw.lower())
    aliases = {
        "large": "Large",
        "largecap": "Large",
        "mid": "Mid",
        "midcap": "Mid",
        "medium": "Mid",
        "small": "Small",
        "smallcap": "Small",
    }
    return aliases.get(normalized, raw.title())


def canonical_factor_value_style(value: object) -> str:
    raw = str(value).strip()
    normalized = re.sub(r"[^a-z]+", "", raw.lower())
    aliases = {
        "value": "Value",
        "core": "Core",
        "blend": "Core",
        "growth": "Growth",
    }
    return aliases.get(normalized, raw.title())


def infer_canonical_region(row: pd.Series | object) -> str:
    def _get(attr: str) -> str:
        if isinstance(row, pd.Series):
            return str(row.get(attr, "")).strip()
        return str(getattr(row, attr, "")).strip()

    candidates = [
        _get("investment_area"),
        _get("name"),
        _get("ticker"),
    ]

    explicit_ticker_map = {
        "VGK": "Europe",
        "VPL": "Pacific",
        "VWO": "Emerging",
        "INAA": "North America",
        "IE0030404903": "North America",
    }
    ticker = normalise_ticker(_get("ticker"))
    if ticker in explicit_ticker_map:
        return explicit_ticker_map[ticker]

    for candidate in candidates:
        mapped = canonical_region_name(candidate)
        if mapped in MSCI_ACWI_IMI_REGION_WEIGHTS:
            return mapped

    combined = " ".join(candidates).strip()
    mapped = canonical_region_name(combined)
    return mapped


def convert_price_series_to_neutral_currency(
    series: pd.Series,
    quote_currency: str,
    neutral_currency: str,
    fx_values_df: pd.DataFrame,
) -> pd.Series:
    series = series.dropna().sort_index()
    if series.empty:
        return pd.Series(dtype=float)

    if neutral_currency == "USD":
        return series.copy()

    if neutral_currency and quote_currency and quote_currency != neutral_currency:
        quote_fx_series = fx_values_df.get(quote_currency, pd.Series(dtype=float)).dropna().sort_index()
        neutral_series = fx_values_df.get(neutral_currency, pd.Series(dtype=float)).dropna().sort_index()
        if quote_fx_series.empty or neutral_series.empty:
            return pd.Series(dtype=float)
        aligned = pd.concat(
            [
                series.rename("price"),
                quote_fx_series.rename("quote_fx"),
                neutral_series.rename("neutral_fx"),
            ],
            axis=1,
        ).sort_index().ffill().dropna()
        if aligned.empty:
            return pd.Series(dtype=float)
        return (aligned["price"] * aligned["quote_fx"] / aligned["neutral_fx"]).dropna()

    if neutral_currency and quote_currency == neutral_currency:
        return series.copy()

    if neutral_currency and not quote_currency:
        return pd.Series(dtype=float)

    return series.copy()


def build_region_performance_df(
    regions_df: pd.DataFrame,
    fx_values_df: pd.DataFrame,
    period_key: str,
    neutral_currency: str = "USD",
    preferred_series_map: dict[str, pd.Series] | None = None,
) -> tuple[pd.DataFrame, pd.Timestamp | None]:
    if regions_df.empty:
        return pd.DataFrame(), None

    regions = regions_df[
        (regions_df["region_flag"] == 1) & regions_df["available"].fillna(False)
    ].copy()
    if regions.empty:
        return pd.DataFrame(), None

    regions["canonical_region"] = regions.apply(infer_canonical_region, axis=1)
    regions = regions[regions["canonical_region"].isin(MSCI_ACWI_IMI_REGION_WEIGHTS.keys())].copy()
    if regions.empty:
        return pd.DataFrame(), None

    tickers = tuple(sorted({normalise_ticker(t) for t in regions["ticker"].tolist() if normalise_ticker(t)}))
    prices_df = fetch_yf_prices(tickers, FX_DOWNLOAD_START)
    if prices_df.empty:
        return pd.DataFrame(), None

    quote_currency_map = fetch_yf_quote_currencies(tickers)
    last_dates = [
        pd.Timestamp(prices_df[c].dropna().index.max())
        for c in prices_df.columns
        if not prices_df[c].dropna().empty
    ]
    if not last_dates:
        return pd.DataFrame(), None

    end_date = min(last_dates)
    start_anchor = get_geo_period_start_anchor(end_date, period_key, fx_values_df, neutral_currency)
    rows = []

    for row in regions.drop_duplicates(subset=["canonical_region"], keep="first").itertuples():
        ticker = normalise_ticker(row.ticker)
        series = get_price_series(prices_df, ticker)
        if series.empty:
            continue
        series = series[series.index <= end_date]
        if series.empty:
            continue

        quote_currency = quote_currency_map.get(ticker, "") or "USD"
        converted_series = convert_price_series_to_neutral_currency(series, quote_currency, neutral_currency, fx_values_df)
        if converted_series.empty:
            continue

        performance = calc_period_return(
            converted_series,
            end_date,
            "Period" if period_key == "MAX" else period_key,
            whole_period_start=start_anchor,
        )
        if pd.isna(performance):
            continue

        rows.append(
            {
                "region": row.canonical_region,
                "ticker": ticker,
                "performance": performance,
                "weight": float(MSCI_ACWI_IMI_REGION_WEIGHTS.get(row.canonical_region, np.nan)),
            }
        )

    out = pd.DataFrame(rows) if rows else pd.DataFrame(columns=["region", "ticker", "performance", "weight"])

    if preferred_series_map is not None and "Global stocks" in preferred_series_map:
        global_series = preferred_series_map.get("Global stocks", pd.Series(dtype=float)).dropna().sort_index()
        global_series = global_series[global_series.index <= end_date]
        if not global_series.empty:
            global_return = calc_period_return(
                global_series,
                end_date,
                "Period" if period_key == "MAX" else period_key,
                whole_period_start=start_anchor,
            )
            if pd.notna(global_return):
                global_df = pd.DataFrame(
                    [{"region": "Global market", "ticker": "Global market", "performance": global_return, "weight": 100.0}]
                )
                if out.empty or out.dropna(how="all").empty:
                    out = global_df
                else:
                    out = pd.concat([global_df, out], ignore_index=True)

    if out.empty:
        return pd.DataFrame(), end_date

    region_order = ["Global market"] + list(MSCI_ACWI_IMI_REGION_WEIGHTS.keys())
    out["region_order"] = out["region"].map({name: idx for idx, name in enumerate(region_order)}).fillna(999)
    out = out.sort_values(["region_order", "region"]).drop(columns="region_order").reset_index(drop=True)
    return out, end_date


def build_region_series_map(
    regions_df: pd.DataFrame,
    fx_values_df: pd.DataFrame,
    neutral_currency: str = "USD",
) -> tuple[dict[str, pd.Series], pd.Timestamp | None]:
    if regions_df.empty:
        return {}, None

    regions = regions_df[
        (regions_df["region_flag"] == 1) & regions_df["available"].fillna(False)
    ].copy()
    if regions.empty:
        return {}, None

    regions["canonical_region"] = regions.apply(infer_canonical_region, axis=1)
    regions = regions[regions["canonical_region"].isin(MSCI_ACWI_IMI_REGION_WEIGHTS.keys())].copy()
    if regions.empty:
        return {}, None

    tickers = tuple(sorted({normalise_ticker(t) for t in regions["ticker"].tolist() if normalise_ticker(t)}))
    prices_df = fetch_yf_prices(tickers, FX_DOWNLOAD_START)
    if prices_df.empty:
        return {}, None

    quote_currency_map = fetch_yf_quote_currencies(tickers)
    last_dates = [
        pd.Timestamp(prices_df[c].dropna().index.max())
        for c in prices_df.columns
        if not prices_df[c].dropna().empty
    ]
    if not last_dates:
        return {}, None

    end_date = min(last_dates)
    out: dict[str, pd.Series] = {}

    for row in regions.drop_duplicates(subset=["canonical_region"], keep="first").itertuples():
        ticker = normalise_ticker(row.ticker)
        series = get_price_series(prices_df, ticker)
        if series.empty:
            continue
        series = series[series.index <= end_date]
        if series.empty:
            continue

        quote_currency = quote_currency_map.get(ticker, "") or ("USD" if neutral_currency == "USD" else "")
        converted_series = (
            series.copy()
            if neutral_currency == "USD"
            else convert_price_series_to_neutral_currency(series, quote_currency, neutral_currency, fx_values_df)
        )
        if converted_series.empty:
            continue
        out[str(row.canonical_region).strip()] = converted_series.dropna().sort_index()

    return out, end_date


def build_country_tiles_html(country_df: pd.DataFrame) -> str:
    if country_df.empty:
        return '<div class="table-shell"><div class="table-empty">No country data available.</div></div>'

    vmin = float(country_df["performance"].min()) if not country_df.empty else -0.1
    vmax = float(country_df["performance"].max()) if not country_df.empty else 0.1
    cards = []
    for row in country_df.itertuples():
        colour = "#3b3b3b" if str(row.country) == "Global market" else heat_colour(float(row.performance), vmin, vmax)
        cards.append(
            (
                f'<div class="country-card" style="background:{colour};">'
                f'<div class="country-card-title">{row.country}</div>'
                f'<div class="country-card-meta">{row.ticker}</div>'
                f'<div class="country-card-value">{format_pct(float(row.performance))}</div>'
                f'</div>'
            )
        )

    return f'<div class="country-card-grid">{"".join(cards)}</div>'


def build_label_tiles_html(label_df: pd.DataFrame) -> str:
    if label_df.empty:
        return '<div class="table-shell"><div class="table-empty">No data available.</div></div>'

    vmin = float(label_df["performance"].min()) if not label_df.empty else -0.1
    vmax = float(label_df["performance"].max()) if not label_df.empty else 0.1
    cards = []
    for row in label_df.itertuples():
        colour = "#3b3b3b" if str(row.label) == "Global market" else heat_colour(float(row.performance), vmin, vmax)
        cards.append(
            (
                f'<div class="country-card" style="background:{colour};">'
                f'<div class="country-card-title">{row.label}</div>'
                f'<div class="country-card-value">{format_pct(float(row.performance))}</div>'
                f'</div>'
            )
        )

    return f'<div class="country-card-grid">{"".join(cards)}</div>'


def build_factor_style_box_html(style_box_df: pd.DataFrame) -> str:
    if style_box_df.empty:
        return '<div class="table-shell"><div class="table-empty">No factor style-box data available.</div></div>'

    row_order = ["Large", "Mid", "Small"]
    col_order = ["Value", "Core", "Growth"]
    grid = (
        style_box_df.pivot(index="size_style", columns="value_style", values="performance")
        .reindex(index=row_order, columns=col_order)
    )
    values = pd.to_numeric(style_box_df["performance"], errors="coerce").dropna()
    vmin = float(values.min()) if not values.empty else -0.1
    vmax = float(values.max()) if not values.empty else 0.1

    header_cells = ['<div class="factor-style-header factor-style-corner">Size \\ Style</div>']
    header_cells.extend(f'<div class="factor-style-header">{column}</div>' for column in col_order)

    body_cells = []
    for row_name in row_order:
        body_cells.append(f'<div class="factor-style-row-label">{row_name}</div>')
        for col_name in col_order:
            value = grid.loc[row_name, col_name]
            if pd.isna(value):
                body_cells.append('<div class="factor-style-cell factor-style-cell-empty">-</div>')
                continue
            background = heat_colour(float(value), vmin, vmax)
            text_colour = text_colour_for_background(background)
            body_cells.append(
                f'<div class="factor-style-cell" style="background:{background};color:{text_colour};">{format_pct(float(value))}</div>'
            )

    return (
        '<div class="factor-style-shell">'
        '<div class="factor-style-grid">'
        f'{"".join(header_cells)}'
        f'{"".join(body_cells)}'
        "</div>"
        "</div>"
    )


def build_region_tiles_html(region_df: pd.DataFrame) -> str:
    if region_df.empty:
        return '<div class="table-shell"><div class="table-empty">No regional data available.</div></div>'

    vmin = float(region_df["performance"].min()) if not region_df.empty else -0.1
    vmax = float(region_df["performance"].max()) if not region_df.empty else 0.1

    global_row = region_df[region_df["region"] == "Global market"]
    regional_rows = region_df[region_df["region"] != "Global market"].copy()
    regional_rows = regional_rows[regional_rows["weight"].notna()].copy()
    total_weight = float(regional_rows["weight"].sum()) if not regional_rows.empty else 0.0

    global_html = ""
    if not global_row.empty:
        row = global_row.iloc[0]
        global_html = (
            f'<div class="region-global-card" style="background:#3b3b3b;">'
            f'<div class="region-card-title">{row["region"]}</div>'
            f'<div class="region-card-meta">{row["ticker"]}</div>'
            f'<div class="region-card-value">{format_pct(float(row["performance"]))}</div>'
            f'</div>'
        )

    region_cards = []
    for _, row in regional_rows.iterrows():
        width_pct = (float(row["weight"]) / total_weight * 100.0) if total_weight > 0 else 25.0
        region_cards.append(
            f'<div class="region-card" style="flex:{width_pct:.4f} 1 0;background:{heat_colour(float(row["performance"]), vmin, vmax)};">'
            f'<div class="region-card-title">{row["region"]}</div>'
            f'<div class="region-card-meta">{row["ticker"]} · {row["weight"]:.1f}%</div>'
            f'<div class="region-card-value">{format_pct(float(row["performance"]))}</div>'
            f'</div>'
        )

    return (
        '<div class="region-card-shell">'
        f"{global_html}"
        f'<div class="region-card-row">{"".join(region_cards)}</div>'
        "</div>"
    )


def build_distinct_colour_map(labels: list[str]) -> dict[str, str]:
    labels = [str(label) for label in labels]
    palette = [
        "#1f77b4", "#d62728", "#2ca02c", "#ff7f0e", "#9467bd",
        "#17becf", "#e377c2", "#8c564b", "#bcbd22", "#7f7f7f",
        "#3366cc", "#dc3912", "#109618", "#990099", "#0099c6",
        "#dd4477", "#66aa00", "#b82e2e", "#316395", "#994499",
        "#22aa99", "#aaaa11", "#6633cc", "#e67300", "#8b0707",
        "#651067", "#329262", "#5574a6", "#3b3eac", "#b77322",
    ]
    out: dict[str, str] = {}
    for idx, label in enumerate(sorted(labels)):
        out[label] = palette[idx % len(palette)]
    if "Global market" in labels:
        out["Global market"] = "#3b3b3b"
    return out


def build_country_patchwork_quilt(
    country_series_map: dict[str, pd.Series],
    end_date: pd.Timestamp | None,
    years_back: int = 10,
    include_labels: list[str] | None = None,
) -> tuple[pd.DataFrame, pd.DataFrame, list[int]]:
    if not country_series_map or end_date is None:
        return pd.DataFrame(), pd.DataFrame(), []

    if include_labels is not None:
        allowed = {str(label) for label in include_labels}
        country_series_map = {label: series for label, series in country_series_map.items() if str(label) in allowed}
        if not country_series_map:
            return pd.DataFrame(), pd.DataFrame(), []

    end_ts = pd.Timestamp(end_date)
    last_complete_year = end_ts.year if (end_ts.month == 12 and end_ts.day == 31) else end_ts.year - 1
    years = list(range(last_complete_year - years_back + 1, last_complete_year + 1))
    if not years:
        return pd.DataFrame(), pd.DataFrame(), []

    returns_by_country: dict[str, dict[int, float]] = {}
    summary_rows: list[dict[str, object]] = []
    quilt_start = pd.Timestamp(years[0] - 1, 12, 31)
    quilt_end = pd.Timestamp(years[-1], 12, 31)

    for country, series in country_series_map.items():
        s = series.dropna().sort_index()
        if s.empty:
            continue

        yearly_returns: dict[int, float] = {}
        valid_all_years = True
        for year in years:
            _, base_level = nearest_level_on_or_before(s, pd.Timestamp(year - 1, 12, 31))
            _, end_level = nearest_level_on_or_before(s, pd.Timestamp(year, 12, 31))
            if base_level is None or end_level is None or base_level <= 0:
                valid_all_years = False
                break
            yearly_returns[year] = (float(end_level) / float(base_level)) - 1

        if not valid_all_years:
            continue

        period_return = calc_period_return(s, quilt_end, "Period", whole_period_start=quilt_start)
        ytd_return = calc_period_return(s, end_ts, "YTD")
        returns_by_country[country] = yearly_returns
        summary_rows.append(
            {
                "country": country,
                "period_return": period_return,
                "ytd_return": ytd_return,
            }
        )

    if not returns_by_country:
        return pd.DataFrame(), pd.DataFrame(), years

    quilt_rows: list[dict[str, object]] = []
    max_rank = len(returns_by_country)
    for year in years:
        ranked = sorted(
            (
                {"country": country, "return_value": country_returns[year]}
                for country, country_returns in returns_by_country.items()
            ),
            key=lambda item: (item["return_value"], item["country"]),
            reverse=True,
        )
        for rank, item in enumerate(ranked, start=1):
            quilt_rows.append(
                {
                    "year": year,
                    "rank": rank,
                    "country": item["country"],
                    "return_value": item["return_value"],
                    "max_rank": max_rank,
                }
            )

    quilt_df = pd.DataFrame(quilt_rows)
    summary_df = pd.DataFrame(summary_rows)
    if summary_df.empty:
        return quilt_df, summary_df, years

    summary_df["ytd_rank"] = summary_df["ytd_return"].rank(method="first", ascending=False, na_option="bottom")
    summary_df["period_rank"] = summary_df["period_return"].rank(method="first", ascending=False, na_option="bottom")
    summary_df = summary_df.sort_values(["country"]).reset_index(drop=True)
    return quilt_df, summary_df, years


def resolve_patchwork_labels(
    patchwork_view: str,
    available_labels: list[str],
    country_rankings_df: pd.DataFrame,
) -> list[str] | None:
    selector = PATCHWORK_COUNTRY_SETS.get(patchwork_view)
    if selector is None:
        return selector
    if isinstance(selector, list):
        labels = list(selector)
        if "Global market" in available_labels and "Global market" not in labels:
            labels = ["Global market"] + labels
        return labels

    if country_rankings_df.empty:
        return None

    limit = 20 if selector == "largest_20" else 10 if selector == "largest_10" else None
    if limit is None:
        return None

    available_set = {str(label) for label in available_labels}
    ranked = country_rankings_df[country_rankings_df["country"].isin(available_set)]["country"].astype(str).tolist()
    labels = ranked[:limit]
    if "Global market" in available_labels and "Global market" not in labels:
        labels = ["Global market"] + labels
    return labels


def build_country_patchwork_html(quilt_df: pd.DataFrame, summary_df: pd.DataFrame, years: list[int]) -> str:
    if quilt_df.empty or summary_df.empty or not years:
        return '<div class="table-shell"><div class="table-empty">No patchwork-quilt data available.</div></div>'

    countries = summary_df["country"].astype(str).tolist()
    colour_map = build_distinct_colour_map(countries)

    header_html = "".join(f'<div class="patchwork-year-header">{year}</div>' for year in years)
    body_cols = []
    for year in years:
        year_rows = quilt_df[quilt_df["year"] == year].sort_values("rank")
        cell_html = []
        for _, row in year_rows.iterrows():
            country = str(row["country"])
            cell_html.append(
                f'<div class="patchwork-cell" style="background:{colour_map[country]};">'
                f'<div class="patchwork-cell-country">{country}</div>'
                f'<div class="patchwork-cell-return">{format_pct(float(row["return_value"]))}</div>'
                f'</div>'
            )
        body_cols.append(f'<div class="patchwork-year-col">{"".join(cell_html)}</div>')

    metric_defs = [
        ("YTD", "ytd_return", "ytd_rank"),
        ("Whole period", "period_return", "period_rank"),
    ]
    metric_headers = "".join(
        f'<div class="patchwork-year-header patchwork-metric-header">{label}</div>' for label, _, _ in metric_defs
    )
    metric_cols = []
    for _, metric_key, rank_key in metric_defs:
        metric_rows = summary_df.sort_values([rank_key, "country"], ascending=[True, True])
        cell_html = []
        for _, row in metric_rows.iterrows():
            country = str(row["country"])
            value = pd.to_numeric(pd.Series([row.get(metric_key)]), errors="coerce").iloc[0]
            cell_html.append(
                f'<div class="patchwork-cell patchwork-summary-cell" style="background:{colour_map[country]};">'
                f'<div class="patchwork-cell-country">{country}</div>'
                f'<div class="patchwork-cell-return">{format_pct(float(value))}</div>'
                f'</div>'
            )
        metric_cols.append(
            f'<div class="patchwork-year-col patchwork-metric-col">{"".join(cell_html)}</div>'
        )

    return (
        '<div class="patchwork-shell">'
        f'<div class="patchwork-grid" style="grid-template-columns: repeat({len(years)}, minmax(92px, 1fr)) repeat({len(metric_defs)}, minmax(100px, 1fr));">'
        f"{header_html}{metric_headers}{''.join(body_cols)}{''.join(metric_cols)}"
        "</div>"
        "</div>"
    )


def build_currency_matrix_html(matrix_df: pd.DataFrame) -> str:
    if matrix_df.empty:
        return '<div class="table-shell"><div class="table-empty">No currency data available.</div></div>'

    display_df = matrix_df.copy().rename(columns={"currency": ""})
    numeric_values = pd.to_numeric(matrix_df.drop(columns="currency").stack(), errors="coerce").dropna()
    vmin = float(numeric_values.min()) if not numeric_values.empty else -0.05
    vmax = float(numeric_values.max()) if not numeric_values.empty else 0.05
    cols = list(display_df.columns)
    first_col_width = 10.5
    other_width = (100 - first_col_width) / (len(cols) - 1) if len(cols) > 1 else 100

    colgroup = "".join(
        [f'<col style="width:{first_col_width if idx == 0 else other_width:.4f}%;">' for idx, _ in enumerate(cols)]
    )
    thead = "".join([f"<th>{col}</th>" for col in cols])
    body_rows = []

    for _, row in display_df.iterrows():
        cells = []
        for idx, col in enumerate(cols):
            value = row[col]
            if idx == 0:
                cells.append(f"<td><b>{value}</b></td>")
                continue
            if pd.isna(value):
                cells.append('<td style="background:#9a9a9a;color:#ffffff;font-weight:600;"></td>')
            else:
                cells.append(
                    f'<td style="background:{currency_heat_colour(float(value), vmin, vmax)};color:#111111;">{format_pct(float(value))}</td>'
                )
        body_rows.append(f"<tr>{''.join(cells)}</tr>")

    return f"""
    <div class="table-shell">
        <table class="custom-data-table">
            <colgroup>{colgroup}</colgroup>
            <thead><tr>{thead}</tr></thead>
            <tbody>{''.join(body_rows)}</tbody>
        </table>
    </div>
    """


def currency_heat_colour(value: float, vmin: float, vmax: float) -> str:
    if pd.isna(value):
        return "#FFFFFF"

    value = float(value)
    white = np.array([255, 255, 255])
    green = np.array([0, 170, 95])
    red = np.array([190, 30, 55])

    if value >= 0:
        norm = 0.0 if vmax <= 0 else min(max(value / vmax, 0), 1)
        rgb = (white * (1 - norm) + green * norm).astype(int)
    else:
        norm = 0.0 if vmin >= 0 else min(max(value / vmin, 0), 1)
        rgb = (white * (1 - norm) + red * norm).astype(int)

    return f"rgb({rgb[0]},{rgb[1]},{rgb[2]})"


def correlation_heat_colour(value: float) -> str:
    if pd.isna(value):
        return "#E9E9E9"
    value = float(value)
    grey = np.array([138, 138, 138])
    green = np.array([0, 170, 95])
    red = np.array([190, 30, 55])

    if np.isclose(value, 1.0):
        rgb = grey
    elif value >= 0.5:
        norm = min(max((value - 0.5) / 0.49, 0), 1)
        rgb = (grey * (1 - norm) + red * norm).astype(int)
    else:
        norm = min(max(value / 0.5, 0), 1)
        rgb = (green * (1 - norm) + grey * norm).astype(int)

    return f"rgb({rgb[0]},{rgb[1]},{rgb[2]})"


def safe_relative_return(asset_return: float, base_return: float) -> float:
    if pd.isna(asset_return) or pd.isna(base_return) or (1 + base_return) <= 0:
        return np.nan
    return ((1 + asset_return) / (1 + base_return)) - 1


def get_display_groups(is_relative_mode: bool, relative_detail_mode: str) -> list[dict]:
    if is_relative_mode and relative_detail_mode == "Minor":
        return DISPLAY_GROUPS_RELATIVE_MINOR
    return DISPLAY_GROUPS_ABSOLUTE


def convert_to_relative_returns(absolute_returns_df: pd.DataFrame, relative_detail_mode: str) -> pd.DataFrame:
    relative_df = absolute_returns_df.copy()
    growth_base_map = MAJOR_GROWTH_BASE_MAP if relative_detail_mode == "Major" else MINOR_GROWTH_BASE_MAP
    all_base_maps = {**growth_base_map, **DEFENSIVE_BASE_MAP}
    asset_to_row = relative_df.set_index("asset_class")

    period_cols = [c for c in relative_df.columns if c != "asset_class"]
    for period in period_cols:
        for idx, row in relative_df.iterrows():
            asset = row["asset_class"]
            base_asset = all_base_maps.get(asset)
            if base_asset is None or base_asset not in asset_to_row.index:
                relative_df.at[idx, period] = np.nan
            else:
                relative_df.at[idx, period] = safe_relative_return(row[period], asset_to_row.at[base_asset, period])

    return relative_df


def find_inflation_column(ts: pd.DataFrame) -> str:
    lookup = {normalise_name(c): str(c).strip() for c in ts.columns}
    for alias in INFLATION_SERIES_ALIASES:
        key = normalise_name(alias)
        if key in lookup:
            return lookup[key]
    raise KeyError(
        "Could not find a UK inflation series in time_series sheet. "
        f"Tried aliases: {', '.join(INFLATION_SERIES_ALIASES)}"
    )


def build_inflation_levels_from_timeseries(ts: pd.DataFrame) -> pd.Series:
    inflation_col = find_inflation_column(ts)
    returns = standardise_series(ts[inflation_col])
    series = pd.Series(returns.values, index=ts["Date"], name="UK inflation").dropna().sort_index()
    if series.empty:
        raise ValueError("Inflation series was found but contains no valid data.")
    levels = (1 + series).cumprod()
    levels.name = "UK inflation"
    return levels


@st.cache_data(show_spinner=False, ttl=86400)
def fetch_ons_raw_csv(csv_url: str) -> pd.DataFrame:
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/146.0.0.0 Safari/537.36"
        ),
        "Accept": "text/csv,application/csv,text/plain,*/*",
        "Referer": "https://www.ons.gov.uk/",
    }
    response = requests.get(csv_url, headers=headers, timeout=30)
    response.raise_for_status()
    return pd.read_csv(StringIO(response.text))


@st.cache_data(show_spinner=False, ttl=86400)
def fetch_ons_cpi_index_series(csv_url: str) -> pd.Series:
    raw = fetch_ons_raw_csv(csv_url)
    out, _ = parse_ons_cpi_index_frame(raw)
    return out


def extend_inflation_levels_with_ons_index(existing_inflation_levels: pd.Series, ons_index_series: pd.Series) -> pd.Series:
    workbook = existing_inflation_levels.dropna().sort_index()
    ons = ons_index_series.dropna().sort_index()

    if workbook.empty:
        raise ValueError("Existing inflation level history is required.")

    common_dates = workbook.index.intersection(ons.index)
    if common_dates.empty:
        raise ValueError("No overlapping dates between workbook inflation and ONS CPI index.")

    anchor_date = common_dates.max()
    workbook_anchor = float(workbook.loc[anchor_date])
    ons_anchor = float(ons.loc[anchor_date])

    if pd.isna(ons_anchor) or ons_anchor == 0:
        raise ValueError("Invalid ONS anchor value.")

    scaled_ons = ons * (workbook_anchor / ons_anchor)
    combined = pd.concat([workbook[workbook.index <= anchor_date], scaled_ons[scaled_ons.index > anchor_date]])
    combined = combined[~combined.index.duplicated(keep="last")].sort_index()
    combined.name = "UK inflation"
    return combined


def build_best_available_inflation_levels(ts: pd.DataFrame) -> tuple[pd.Series | None, str, str | None]:
    workbook_levels = build_inflation_levels_from_timeseries(ts)
    try:
        ons_index = fetch_ons_cpi_index_series(ONS_CPI_INDEX_CSV_URL)
        extended = extend_inflation_levels_with_ons_index(workbook_levels, ons_index)
        if extended.index.max() > workbook_levels.index.max():
            return extended, "Workbook time_series + ONS CPI index extension", None
        return workbook_levels, "Workbook time_series", "ONS fetched, but no extension beyond workbook end date."
    except Exception as exc:
        return workbook_levels, "Workbook time_series", f"ONS extension failed: {type(exc).__name__}: {exc}"


def build_monthly_returns_from_levels(levels: pd.Series) -> pd.Series:
    series = levels.dropna().sort_index()
    if series.empty:
        return pd.Series(dtype=float, name="UK inflation")
    returns = series.pct_change()
    returns.name = "UK inflation"
    return returns


def get_live_consistent_end_date(stitched_series_map: dict[str, pd.Series], live_diag: pd.DataFrame) -> pd.Timestamp:
    if live_diag is not None and not live_diag.empty:
        live_rows = live_diag[live_diag["series_type"].isin(["stitched", "live_only"])].copy()
        if not live_rows.empty:
            live_last_dates = pd.to_datetime(live_rows["live_last_date"], errors="coerce").dropna()
            if not live_last_dates.empty:
                return pd.Timestamp(live_last_dates.min())

    last_dates = [pd.Timestamp(s.dropna().index.max()) for s in stitched_series_map.values() if not s.dropna().empty]
    return max(last_dates) if last_dates else pd.Timestamp.today().normalize()


def get_dashboard_end_date(
    stitched_series_map: dict[str, pd.Series],
    live_diag: pd.DataFrame,
    inflation_levels: pd.Series | None,
    is_real_mode: bool,
) -> pd.Timestamp:
    end_date = get_live_consistent_end_date(stitched_series_map, live_diag)
    if is_real_mode and inflation_levels is not None and not inflation_levels.dropna().empty:
        end_date = min(end_date, pd.Timestamp(inflation_levels.dropna().index.max()))
    return end_date


def get_series_map_date_bounds(series_map: dict[str, pd.Series]) -> tuple[pd.Timestamp | None, pd.Timestamp | None]:
    min_dates = []
    max_dates = []
    for series in series_map.values():
        s = series.dropna().sort_index()
        if s.empty:
            continue
        min_dates.append(pd.Timestamp(s.index.min()))
        max_dates.append(pd.Timestamp(s.index.max()))
    if not min_dates or not max_dates:
        return None, None
    return min(min_dates), max(max_dates)


def filter_series_to_window(series: pd.Series, start_date: pd.Timestamp | None, end_date: pd.Timestamp | None) -> pd.Series:
    out = series.dropna().sort_index()
    if start_date is not None:
        out = out[out.index >= start_date]
    if end_date is not None:
        out = out[out.index <= end_date]
    return out


def filter_series_map_to_window(
    series_map: dict[str, pd.Series],
    start_date: pd.Timestamp | None,
    end_date: pd.Timestamp | None,
) -> dict[str, pd.Series]:
    return {
        asset_class: filter_series_to_window(series, start_date, end_date)
        for asset_class, series in series_map.items()
    }


def get_chart_period_start_date(period_key: str, end_date: pd.Timestamp, min_start_date: pd.Timestamp) -> pd.Timestamp:
    end_date = pd.Timestamp(end_date).normalize()
    min_start_date = pd.Timestamp(min_start_date).normalize()
    if period_key == "YTD":
        start_date = pd.Timestamp(end_date.year - 1, 12, 31)
    else:
        years_map = {"1Y": 1, "3Y": 3, "5Y": 5, "10Y": 10, "20Y": 20}
        years = years_map.get(period_key)
        start_date = end_date - pd.DateOffset(years=years) if years is not None else min_start_date
    return max(min_start_date, pd.Timestamp(start_date).normalize())


def match_chart_period_from_dates(
    start_date: pd.Timestamp,
    end_date: pd.Timestamp,
    min_start_date: pd.Timestamp,
) -> str:
    start_date = pd.Timestamp(start_date).normalize()
    end_date = pd.Timestamp(end_date).normalize()
    min_start_date = pd.Timestamp(min_start_date).normalize()
    for period_key in CHART_PERIODS.keys():
        if get_chart_period_start_date(period_key, end_date, min_start_date) == start_date:
            return period_key
    return "Custom"


def convert_to_real_returns(nominal_returns_df: pd.DataFrame, inflation_returns_df: pd.DataFrame) -> pd.DataFrame:
    real_df = nominal_returns_df.copy()
    if inflation_returns_df.empty:
        return real_df

    infl_row = inflation_returns_df[inflation_returns_df["asset_class"] == "UK inflation"]
    if infl_row.empty:
        return real_df

    infl_row = infl_row.iloc[0]
    period_cols = [c for c in real_df.columns if c != "asset_class"]
    for period in period_cols:
        infl_val = infl_row.get(period, np.nan)
        for idx, row in real_df.iterrows():
            real_df.at[idx, period] = safe_relative_return(row[period], infl_val)

    return real_df


def get_methodology_paragraph(page_name: str, is_relative_mode: bool, relative_detail_mode: str, is_real_mode: bool, inflation_source_note: str) -> str:
    basis = "real" if is_real_mode else "nominal"

    if page_name == "Dashboard":
        prefix = f"This tab shows annualised {basis} GBP returns across the displayed horizons, with YTD shown cumulatively."
        if is_relative_mode:
            if relative_detail_mode == "Major":
                relative_text = (
                    " In relative mode, growth assets are shown relative to Global stocks and defensive assets are shown relative to Cash (GBP)."
                )
            else:
                relative_text = (
                    " In relative mode (Minor), UK, EM, REITs, Developed Value and Developed Small are shown relative to Developed stocks; "
                    "EM Value and EM Small are shown relative to Emerging stocks; UK Value and UK Small are shown relative to UK stocks; "
                    "and defensive assets are shown relative to Cash (GBP)."
                )
        else:
            relative_text = " In absolute mode, headline returns are shown directly for each asset class."
    else:
        prefix = (
            "This tab shows return in GBP from the common inception date to the selected end date, "
            "with the chart normalised to a starting value of £1.00. "
            "For the growth of wealth chart, daily live-fund history is preferred where available and monthly index history is stitched in before it to extend the chart further back."
        )
        relative_text = ""

    inflation_text = ""
    if is_real_mode:
        inflation_text = (
            f" Real returns are calculated using UK inflation via (1+asset return)/(1+inflation return)-1. "
            f"Current inflation source: {inflation_source_note}."
        )

    source_text = (
        " Albion index series history is used as the preferred source where available. "
        "For the dashboard and tables, live yfinance mappings are only used after index history ends. "
        "More information on the Albion indices can be found at "
        '<a href="https://smartersuccess.net/indices" target="_blank">smartersuccess.net/indices</a>.'
    )
    return prefix + relative_text + inflation_text + source_text


def dataframe_to_csv_download(df: pd.DataFrame) -> bytes:
    return prepare_dataframe_for_display(df).to_csv(index=False).encode("utf-8")


def prepare_dataframe_for_display(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    if out.empty:
        return out

    for col in out.columns:
        series = out[col]

        if pd.api.types.is_datetime64_any_dtype(series):
            out[col] = pd.to_datetime(series, errors="coerce").dt.strftime("%d/%m/%Y").fillna("")
            continue

        if pd.api.types.is_timedelta64_dtype(series):
            out[col] = series.astype("string").fillna("")
            continue

        if pd.api.types.is_object_dtype(series) or pd.api.types.is_string_dtype(series):
            non_null = series.dropna()
            if non_null.empty:
                out[col] = series.astype("string").fillna("")
                continue

            non_null_str = non_null.astype(str).str.strip()
            date_like_mask = (
                non_null_str.str.match(r"^\d{1,2}/\d{1,2}/\d{4}$", na=False)
                | non_null_str.str.match(r"^\d{4}-\d{2}-\d{2}$", na=False)
                | non_null_str.str.match(r"^\d{4}/\d{2}/\d{2}$", na=False)
            )
            if len(non_null_str) > 0 and date_like_mask.all():
                full_str = series.astype("string").fillna("").str.strip()
                parsed_dates = pd.Series(pd.NaT, index=out.index, dtype="datetime64[ns]")

                slash_mask = full_str.str.match(r"^\d{1,2}/\d{1,2}/\d{4}$", na=False)
                if slash_mask.any():
                    parsed_dates.loc[slash_mask] = pd.to_datetime(
                        full_str.loc[slash_mask],
                        format="%d/%m/%Y",
                        errors="coerce",
                    )

                iso_dash_mask = full_str.str.match(r"^\d{4}-\d{2}-\d{2}$", na=False)
                if iso_dash_mask.any():
                    parsed_dates.loc[iso_dash_mask] = pd.to_datetime(
                        full_str.loc[iso_dash_mask],
                        format="%Y-%m-%d",
                        errors="coerce",
                    )

                iso_slash_mask = full_str.str.match(r"^\d{4}/\d{2}/\d{2}$", na=False)
                if iso_slash_mask.any():
                    parsed_dates.loc[iso_slash_mask] = pd.to_datetime(
                        full_str.loc[iso_slash_mask],
                        format="%Y/%m/%d",
                        errors="coerce",
                    )

                out[col] = parsed_dates.dt.strftime("%d/%m/%Y").fillna("")
                continue

            numeric = pd.to_numeric(non_null, errors="coerce")
            if len(non_null) > 0 and numeric.notna().sum() == len(non_null):
                out[col] = pd.to_numeric(series, errors="coerce")
                continue

            out[col] = series.astype("string").fillna("")

    return out


def format_diag_date(value: object) -> str:
    ts = pd.to_datetime(value, errors="coerce")
    return "" if pd.isna(ts) else ts.strftime("%d/%m/%Y")


def format_diagnostic_table(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    for col in ["index_last_date", "live_first_date", "live_last_date", "stitch_anchor_date"]:
        if col in out.columns:
            out[col] = pd.to_datetime(out[col], errors="coerce").dt.strftime("%d/%m/%Y").fillna("")
    return prepare_dataframe_for_display(out)


def build_series_summary(series: pd.Series | None) -> dict[str, object]:
    s = pd.Series(dtype=float) if series is None else series.dropna().sort_index()
    if s.empty:
        return {
            "points": 0,
            "first_date": pd.NaT,
            "last_date": pd.NaT,
            "last_value": np.nan,
        }
    return {
        "points": int(len(s)),
        "first_date": pd.Timestamp(s.index.min()),
        "last_date": pd.Timestamp(s.index.max()),
        "last_value": float(s.iloc[-1]),
    }


def parse_ons_cpi_index_frame(raw: pd.DataFrame) -> tuple[pd.Series, pd.DataFrame]:
    raw = raw.copy()
    raw.columns = [str(c).strip() for c in raw.columns]

    if raw.shape[1] < 2:
        raise ValueError("ONS CPI index CSV did not contain at least two columns.")

    date_col, value_col = raw.columns[:2]
    df = raw[[date_col, value_col]].copy()
    df.columns = ["date_raw", "value_raw"]

    df["date_raw"] = df["date_raw"].astype(str).str.strip().str.upper()
    df["value"] = pd.to_numeric(df["value_raw"], errors="coerce")
    df["Date"] = pd.NaT
    df["match_type"] = ""

    mask_yyyy_mon = df["date_raw"].str.match(r"^\d{4}\s+[A-Z]{3}$", na=False)
    if mask_yyyy_mon.any():
        df.loc[mask_yyyy_mon, "Date"] = pd.to_datetime(
            "01 " + df.loc[mask_yyyy_mon, "date_raw"],
            format="%d %Y %b",
            errors="coerce",
        )
        df.loc[mask_yyyy_mon, "match_type"] = "YYYY MON"

    mask_mon_yyyy = df["date_raw"].str.match(r"^[A-Z]{3}\s+\d{4}$", na=False)
    if mask_mon_yyyy.any():
        df.loc[mask_mon_yyyy, "Date"] = pd.to_datetime(
            "01 " + df.loc[mask_mon_yyyy, "date_raw"],
            format="%d %b %Y",
            errors="coerce",
        )
        df.loc[mask_mon_yyyy, "match_type"] = "MON YYYY"

    mask_year = df["date_raw"].str.match(r"^\d{4}$", na=False)
    if mask_year.any():
        df.loc[mask_year, "Date"] = pd.to_datetime(
            df.loc[mask_year, "date_raw"] + "-01-01",
            format="%Y-%m-%d",
            errors="coerce",
        )
        df.loc[mask_year, "match_type"] = "YYYY"

    parsed = df.dropna(subset=["value", "Date"]).copy()
    parsed["Date"] = pd.to_datetime(parsed["Date"]) + MonthEnd(0)
    parsed = parsed.sort_values("Date")

    out = pd.Series(parsed["value"].values, index=parsed["Date"], name="ONS CPI index")
    out = out[~out.index.duplicated(keep="last")].sort_index()

    if out.empty:
        raise ValueError("ONS CPI index CSV produced no valid rows.")
    return out, df


def build_ons_fetch_diagnostics(csv_url: str) -> tuple[pd.DataFrame, pd.DataFrame]:
    try:
        raw = fetch_ons_raw_csv(csv_url)
        parsed_series, parsed_df = parse_ons_cpi_index_frame(raw)
        summary = pd.DataFrame(
            [
                {"metric": "Fetch status", "value": "OK"},
                {"metric": "Source URL", "value": csv_url},
                {"metric": "Raw rows", "value": int(len(raw))},
                {"metric": "Raw columns", "value": int(raw.shape[1])},
                {"metric": "Parsed numeric rows", "value": int(parsed_df["value"].notna().sum())},
                {"metric": "Parsed dated rows", "value": int(parsed_df["Date"].notna().sum())},
                {"metric": "Output rows", "value": int(len(parsed_series))},
                {"metric": "Output first date", "value": format_diag_date(parsed_series.index.min())},
                {"metric": "Output last date", "value": format_diag_date(parsed_series.index.max())},
                {"metric": "Output last value", "value": round(float(parsed_series.iloc[-1]), 4)},
            ]
        )

        preview = parsed_df.copy()
        preview["Date"] = pd.to_datetime(preview["Date"], errors="coerce").dt.strftime("%d/%m/%Y").fillna("")
        return summary, preview
    except Exception as exc:
        summary = pd.DataFrame(
            [
                {"metric": "Fetch status", "value": "Failed"},
                {"metric": "Source URL", "value": csv_url},
                {"metric": "Error type", "value": type(exc).__name__},
                {"metric": "Error", "value": str(exc)},
            ]
        )
        return summary, pd.DataFrame()


@st.cache_data(show_spinner=False, ttl=43200)
def fetch_boe_yield_curve_zip(zip_url: str) -> bytes:
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/146.0.0.0 Safari/537.36"
        )
    }
    response = requests.get(zip_url, headers=headers, timeout=45)
    response.raise_for_status()
    return response.content


@st.cache_data(show_spinner=False, ttl=43200)
def fetch_dividenddata_html(page_url: str) -> str:
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/146.0.0.0 Safari/537.36"
        )
    }
    response = requests.get(page_url, headers=headers, timeout=45)
    response.raise_for_status()
    return response.text


def fetch_dividenddata_gilt_detail_html(ticker: str) -> str:
    return fetch_dividenddata_html(DIVIDENDDATA_GILT_DETAIL_URL.format(ticker=ticker))


def fetch_worldgovernmentbonds_html(page_url: str) -> str:
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/124.0 Safari/537.36"
        )
    }
    response = requests.get(page_url, headers=headers, timeout=45)
    response.raise_for_status()
    return response.text


@st.cache_data(show_spinner=False, ttl=43200)
def fetch_companiesmarketcap_country_rankings(page_url: str) -> pd.DataFrame:
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/146.0.0.0 Safari/537.36"
        )
    }
    response = requests.get(page_url, headers=headers, timeout=45)
    response.raise_for_status()

    tables = pd.read_html(StringIO(response.text))
    if not tables:
        return pd.DataFrame(columns=["rank", "country"])

    table = tables[0].copy()
    table.columns = [str(col).strip().lower() for col in table.columns]
    rank_col = next((col for col in table.columns if "rank" in col), None)
    country_col = next((col for col in table.columns if "name" in col), None)
    if rank_col is None or country_col is None:
        return pd.DataFrame(columns=["rank", "country"])

    out = table[[rank_col, country_col]].copy()
    out.columns = ["rank", "country"]
    out["rank"] = pd.to_numeric(out["rank"], errors="coerce")
    out["country"] = out["country"].astype(str).str.strip()
    out = out.dropna(subset=["rank", "country"])

    name_map = {
        "United States": "USA",
        "United Kingdom": "UK",
        "South Korea": "Korea",
        "United Arab Emirates": "UAE",
    }
    out["country"] = out["country"].replace(name_map)
    return out.sort_values("rank").reset_index(drop=True)


def fetch_worldgovernmentbonds_country_payload(page_url: str, global_vars: dict) -> dict:
    headers = {
        "Origin": WORLD_GOVERNMENT_BONDS_BASE_URL,
        "Referer": page_url,
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/124.0 Safari/537.36"
        ),
        "X-Requested-With": "XMLHttpRequest",
    }
    payload = {"GLOBALVAR": global_vars}
    response = requests.post(
        f"{WORLD_GOVERNMENT_BONDS_BASE_URL}/wp-json/country/v1/main",
        headers=headers,
        json=payload,
        timeout=45,
    )
    response.raise_for_status()
    return response.json()


class SimpleHTMLTableParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.tables: list[list[list[str]]] = []
        self._in_table = False
        self._in_row = False
        self._in_cell = False
        self._current_table: list[list[str]] = []
        self._current_row: list[str] = []
        self._current_cell: list[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag == "table":
            self._in_table = True
            self._current_table = []
        elif tag == "tr" and self._in_table:
            self._in_row = True
            self._current_row = []
        elif tag in {"td", "th"} and self._in_row:
            self._in_cell = True
            self._current_cell = []

    def handle_data(self, data: str) -> None:
        if self._in_cell:
            self._current_cell.append(data)

    def handle_endtag(self, tag: str) -> None:
        if tag in {"td", "th"} and self._in_cell:
            cell_text = " ".join(part.strip() for part in self._current_cell if part.strip()).strip()
            self._current_row.append(cell_text)
            self._in_cell = False
            self._current_cell = []
        elif tag == "tr" and self._in_row:
            if self._current_row:
                self._current_table.append(self._current_row)
            self._in_row = False
            self._current_row = []
        elif tag == "table" and self._in_table:
            if self._current_table:
                self.tables.append(self._current_table)
            self._in_table = False
            self._current_table = []


class TextExtractingHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = str(data).strip()
        if text:
            self.parts.append(text)


def pick_boe_zip_member(member_names: list[str], target_text: str) -> str:
    target = target_text.strip().lower()
    candidates = [name for name in member_names if target in name.lower()]
    if not candidates:
        raise FileNotFoundError(f"Could not find '{target_text}' in BOE yield-curve zip.")
    return sorted(candidates)[0]


def parse_maturity_label(value: object) -> float:
    if pd.isna(value):
        return np.nan
    text = str(value).strip()
    numeric = pd.to_numeric(text, errors="coerce")
    if pd.notna(numeric):
        return float(numeric)
    match = re.search(r"(\d+(?:\.\d+)?)", text)
    return float(match.group(1)) if match else np.nan


def parse_time_to_maturity(value: object) -> float:
    if pd.isna(value):
        return np.nan
    text = str(value).strip().lower()
    if not text:
        return np.nan

    years = 0.0
    months = 0.0
    days = 0.0

    year_match = re.search(r"(\d+(?:\.\d+)?)\s*year", text)
    month_match = re.search(r"(\d+(?:\.\d+)?)\s*month", text)
    day_match = re.search(r"(\d+(?:\.\d+)?)\s*day", text)

    if year_match:
        years = float(year_match.group(1))
    if month_match:
        months = float(month_match.group(1))
    if day_match:
        days = float(day_match.group(1))

    if years == 0 and months == 0 and days == 0:
        numeric = pd.to_numeric(text, errors="coerce")
        return float(numeric) if pd.notna(numeric) else np.nan

    return years + (months / 12.0) + (days / 365.25)


def parse_residual_maturity(value: object) -> float:
    if pd.isna(value):
        return np.nan
    text = str(value).strip().lower()
    if not text:
        return np.nan

    month_match = re.search(r"(\d+(?:\.\d+)?)\s*month", text)
    year_match = re.search(r"(\d+(?:\.\d+)?)\s*year", text)
    day_match = re.search(r"(\d+(?:\.\d+)?)\s*day", text)

    years = float(year_match.group(1)) if year_match else 0.0
    months = float(month_match.group(1)) if month_match else 0.0
    days = float(day_match.group(1)) if day_match else 0.0
    if years == 0 and months == 0 and days == 0:
        numeric = pd.to_numeric(text, errors="coerce")
        return float(numeric) if pd.notna(numeric) else np.nan
    return years + (months / 12.0) + (days / 365.25)


def extract_worldgovernmentbonds_js_global_vars(html: str) -> dict:
    match = re.search(r"var\s+jsGlobalVars\s*=\s*(\{.*?\});", html, flags=re.DOTALL)
    if match is None:
        raise ValueError("WorldGovernmentBonds jsGlobalVars payload not found in country page HTML.")
    return json.loads(match.group(1))


def parse_worldgovernmentbonds_yield_curve_table(html: str, country_name: str, page_url: str) -> tuple[pd.DataFrame, pd.DataFrame]:
    parser = SimpleHTMLTableParser()
    parser.feed(html)
    rows = []

    for table_rows in parser.tables:
        for row in table_rows:
            if len(row) < 3:
                continue
            maturity_text = str(row[1]).strip()
            if not maturity_text:
                continue
            maturity_years = parse_residual_maturity(maturity_text)
            yield_percent = pd.to_numeric(
                str(row[2]).replace("%", "").replace(",", "").strip(),
                errors="coerce",
            )
            if pd.isna(maturity_years) or pd.isna(yield_percent):
                continue
            rows.append(
                {
                    "country": country_name,
                    "maturity_label": maturity_text,
                    "maturity_years": float(maturity_years),
                    "yield_percent": float(yield_percent),
                    "curve_date": pd.Timestamp.today().normalize(),
                    "source_url": page_url,
                }
            )

    if not rows:
        raise ValueError(f"WorldGovernmentBonds yield-curve rows not found for {country_name}.")

    out = pd.DataFrame(rows).dropna(subset=["maturity_years", "yield_percent"])

    out = out.sort_values("maturity_years").drop_duplicates(subset=["maturity_years"], keep="first")
    preview = out[["country", "maturity_label", "maturity_years", "yield_percent", "source_url"]].copy()
    return out, preview


def parse_dividenddata_text_fallback(html: str) -> tuple[pd.DataFrame, pd.DataFrame]:
    parser = TextExtractingHTMLParser()
    parser.feed(html)
    tokens = [p.strip() for p in parser.parts if str(p).strip()]
    rows = []
    epic_pattern = re.compile(r"^[A-Z0-9]{2,6}$")
    percent_pattern = re.compile(r"^-?\d+(?:\.\d+)?%$")

    def extract_maturity_text(window: list[str]) -> str | None:
        for i, tok in enumerate(window):
            low = tok.lower()
            if re.fullmatch(r"\d+(?:\.\d+)?", tok):
                if i + 1 < len(window) and window[i + 1].lower().startswith("year"):
                    maturity = f"{tok} {window[i + 1]}"
                    if i + 3 < len(window) and re.fullmatch(r"\d+(?:\.\d+)?", window[i + 2]) and window[i + 3].lower().startswith("day"):
                        maturity += f" {window[i + 2]} {window[i + 3]}"
                    elif i + 3 < len(window) and re.fullmatch(r"\d+(?:\.\d+)?", window[i + 2]) and window[i + 3].lower().startswith("month"):
                        maturity += f" {window[i + 2]} {window[i + 3]}"
                    return maturity
            if "year" in low:
                parts = [tok]
                if i + 2 < len(window) and re.fullmatch(r"\d+(?:\.\d+)?", window[i + 1]) and (
                    window[i + 2].lower().startswith("day") or window[i + 2].lower().startswith("month")
                ):
                    parts.extend([window[i + 1], window[i + 2]])
                return " ".join(parts)
        return None

    for idx, token in enumerate(tokens):
        if not epic_pattern.match(token):
            continue

        lookahead = tokens[idx : idx + 40]
        maturity_text = extract_maturity_text(lookahead)
        if maturity_text is None:
            continue

        maturity_pos = 0
        for j, t in enumerate(lookahead):
            if maturity_text.startswith(t) or t in maturity_text:
                maturity_pos = j
                break

        percent_tokens = [t for t in lookahead[maturity_pos + 1 :] if percent_pattern.match(t)]
        if not percent_tokens:
            continue

        real_yield_text = percent_tokens[-1]
        real_yield = pd.to_numeric(real_yield_text.replace("%", ""), errors="coerce")
        maturity_years = parse_time_to_maturity(maturity_text)
        if pd.isna(real_yield) or pd.isna(maturity_years):
            continue

        rows.append(
            {
                "maturity_years": maturity_years,
                "yield_percent": float(real_yield),
                "curve_type": "Real",
                "curve_date": pd.Timestamp.today().normalize(),
                "source": "DividendData",
                "epic": token,
                "time_to_maturity": maturity_text,
                "real_yield_raw": real_yield_text,
            }
        )

    if not rows:
        joined_text = " ".join(tokens)
        regex_rows = []
        row_pattern = re.compile(
            r"\b(?P<epic>[A-Z0-9]{2,6})\b.*?(?P<maturity>\d+(?:\.\d+)?\s+years?(?:\s+\d+(?:\.\d+)?\s+(?:days?|months?))?).*?(?P<real_yield>-?\d+(?:\.\d+)?)%",
            flags=re.IGNORECASE,
        )
        for match in row_pattern.finditer(joined_text):
            maturity_text = match.group("maturity")
            maturity_years = parse_time_to_maturity(maturity_text)
            real_yield = pd.to_numeric(match.group("real_yield"), errors="coerce")
            if pd.isna(maturity_years) or pd.isna(real_yield):
                continue
            regex_rows.append(
                {
                    "maturity_years": maturity_years,
                    "yield_percent": float(real_yield),
                    "curve_type": "Real",
                    "curve_date": pd.Timestamp.today().normalize(),
                    "source": "DividendData",
                    "epic": match.group("epic"),
                    "time_to_maturity": maturity_text,
                    "real_yield_raw": f"{match.group('real_yield')}%",
                }
            )

        if not regex_rows:
            raise ValueError("DividendData fallback text parser could not find index-linked gilt rows.")
        rows = regex_rows
    out = pd.DataFrame(rows).dropna(subset=["maturity_years", "yield_percent"])
    out = out.sort_values("maturity_years").drop_duplicates(subset=["maturity_years"], keep="first")
    preview = out[["epic", "time_to_maturity", "real_yield_raw", "maturity_years", "yield_percent", "source"]].copy()
    return out[["maturity_years", "yield_percent", "curve_type", "curve_date", "source"]], preview


def fetch_dividenddata_short_end_detail_fallback(tickers: list[str]) -> tuple[pd.DataFrame, pd.DataFrame]:
    rows = []
    preview_rows = []
    today = pd.Timestamp.today().normalize()

    for ticker in tickers:
        try:
            html = fetch_dividenddata_gilt_detail_html(ticker)
        except Exception as exc:
            preview_rows.append(
                {
                    "epic": ticker,
                    "time_to_maturity": "",
                    "real_yield_raw": "",
                    "maturity_years": np.nan,
                    "yield_percent": np.nan,
                    "source": "DividendData detail",
                    "status": f"Fetch failed: {type(exc).__name__}",
                }
            )
            continue

        text_parser = TextExtractingHTMLParser()
        text_parser.feed(html)
        text = " ".join(text_parser.parts)

        maturity_match = re.search(r"Maturity Date.*?([0-9]{1,2}-[A-Za-z]{3}-[0-9]{4})", text, flags=re.IGNORECASE)
        real_yield_match = re.search(r"Real Yield.*?(-?\d+(?:\.\d+)?)%", text, flags=re.IGNORECASE)
        if real_yield_match is None:
            real_yield_match = re.search(r"Current Yield.*?(-?\d+(?:\.\d+)?)%", text, flags=re.IGNORECASE)

        if maturity_match is None or real_yield_match is None:
            preview_rows.append(
                {
                    "epic": ticker,
                    "time_to_maturity": "",
                    "real_yield_raw": "",
                    "maturity_years": np.nan,
                    "yield_percent": np.nan,
                    "source": "DividendData detail",
                    "status": "Parse failed",
                }
            )
            continue

        maturity_date = pd.to_datetime(maturity_match.group(1), format="%d-%b-%Y", errors="coerce")
        real_yield = pd.to_numeric(real_yield_match.group(1), errors="coerce")
        if pd.isna(maturity_date) or pd.isna(real_yield):
            preview_rows.append(
                {
                    "epic": ticker,
                    "time_to_maturity": "",
                    "real_yield_raw": "",
                    "maturity_years": np.nan,
                    "yield_percent": np.nan,
                    "source": "DividendData detail",
                    "status": "Numeric/date conversion failed",
                }
            )
            continue

        maturity_years = max((pd.Timestamp(maturity_date) - today).days / 365.25, 0)
        if maturity_years <= 0:
            preview_rows.append(
                {
                    "epic": ticker,
                    "time_to_maturity": "",
                    "real_yield_raw": "",
                    "maturity_years": np.nan,
                    "yield_percent": np.nan,
                    "source": "DividendData detail",
                    "status": "Maturity already passed",
                }
            )
            continue

        row = {
            "maturity_years": maturity_years,
            "yield_percent": float(real_yield),
            "curve_type": "Real",
            "curve_date": today,
            "source": "DividendData detail",
            "epic": ticker,
            "time_to_maturity": f"{maturity_years:.3f} years",
            "real_yield_raw": f"{float(real_yield):.3f}%",
            "status": "OK",
        }
        rows.append(row)
        preview_rows.append(
            {
                "epic": ticker,
                "time_to_maturity": row["time_to_maturity"],
                "real_yield_raw": row["real_yield_raw"],
                "maturity_years": row["maturity_years"],
                "yield_percent": row["yield_percent"],
                "source": row["source"],
                "status": "OK",
            }
        )

    if not rows:
        preview = pd.DataFrame(preview_rows)
        raise ValueError(
            "DividendData detail-page fallback could not parse any short-end index-linked gilts. "
            f"Attempted tickers: {', '.join(tickers)}"
        )

    out = pd.DataFrame(rows).sort_values("maturity_years").drop_duplicates(subset=["maturity_years"], keep="first")
    preview = pd.DataFrame(preview_rows)
    return out[["maturity_years", "yield_percent", "curve_type", "curve_date", "source"]], preview


def parse_boe_spot_curve_workbook(workbook_bytes: bytes, curve_type: str) -> tuple[pd.DataFrame, pd.Timestamp, pd.DataFrame]:
    raw = pd.read_excel(BytesIO(workbook_bytes), sheet_name="4. spot curve", header=None)
    if raw.shape[0] < 6 or raw.shape[1] < 2:
        raise ValueError(f"BOE {curve_type} spot curve sheet is missing expected rows/columns.")

    candidate_rows = []
    for row_idx in range(min(10, len(raw))):
        parsed_row = raw.iloc[row_idx, 1:].map(parse_maturity_label)
        candidate_rows.append((row_idx, int(parsed_row.notna().sum()), parsed_row))

    maturity_row_idx, maturity_count, maturities = max(candidate_rows, key=lambda x: x[1])
    if maturity_count == 0:
        raise ValueError(f"BOE {curve_type} spot curve maturities could not be parsed.")

    data_rows = raw.iloc[maturity_row_idx + 1 :, :].copy()
    data_rows = data_rows.rename(columns={0: "date_raw"})
    data_rows["curve_date"] = pd.to_datetime(data_rows["date_raw"], errors="coerce")
    value_block = data_rows.iloc[:, 1 : 1 + len(maturities)].apply(pd.to_numeric, errors="coerce")

    valid_mask = data_rows["curve_date"].notna() & value_block.notna().any(axis=1)
    valid_rows = data_rows.loc[valid_mask].copy()
    valid_values = value_block.loc[valid_mask].copy()
    if valid_rows.empty:
        raise ValueError(f"BOE {curve_type} spot curve sheet had no valid dated rows.")

    latest_idx = valid_rows.index[-1]
    curve_date = pd.Timestamp(valid_rows.loc[latest_idx, "curve_date"])
    latest_values = valid_values.loc[latest_idx]

    points = pd.DataFrame(
        {
            "maturity_years": pd.to_numeric(maturities, errors="coerce"),
            "yield_percent": pd.to_numeric(latest_values, errors="coerce"),
            "curve_type": curve_type,
            "curve_date": curve_date,
        }
    ).dropna(subset=["maturity_years", "yield_percent"])

    preview = valid_rows[["curve_date"]].copy()
    preview["points_available"] = valid_values.notna().sum(axis=1).values
    preview["curve_type"] = curve_type
    preview["maturity_row_idx"] = maturity_row_idx + 1
    preview["curve_date"] = pd.to_datetime(preview["curve_date"], errors="coerce")
    return points.sort_values("maturity_years"), curve_date, preview


def parse_boe_spot_curve_history_workbook(workbook_bytes: bytes, curve_type: str) -> tuple[pd.DataFrame, pd.DataFrame]:
    raw = pd.read_excel(BytesIO(workbook_bytes), sheet_name="4. spot curve", header=None)
    if raw.shape[0] < 6 or raw.shape[1] < 2:
        raise ValueError(f"BOE {curve_type} spot curve history sheet is missing expected rows/columns.")

    candidate_rows = []
    for row_idx in range(min(10, len(raw))):
        parsed_row = raw.iloc[row_idx, 1:].map(parse_maturity_label)
        candidate_rows.append((row_idx, int(parsed_row.notna().sum()), parsed_row))

    maturity_row_idx, maturity_count, maturities = max(candidate_rows, key=lambda x: x[1])
    if maturity_count == 0:
        raise ValueError(f"BOE {curve_type} spot curve history maturities could not be parsed.")

    data_rows = raw.iloc[maturity_row_idx + 1 :, :].copy()
    data_rows = data_rows.rename(columns={0: "date_raw"})
    data_rows["curve_date"] = pd.to_datetime(data_rows["date_raw"], errors="coerce")
    value_block = data_rows.iloc[:, 1 : 1 + len(maturities)].apply(pd.to_numeric, errors="coerce")

    valid_mask = data_rows["curve_date"].notna() & value_block.notna().any(axis=1)
    valid_rows = data_rows.loc[valid_mask].copy()
    valid_values = value_block.loc[valid_mask].copy()
    if valid_rows.empty:
        raise ValueError(f"BOE {curve_type} spot curve history sheet had no valid dated rows.")

    frames = []
    for idx in valid_rows.index:
        curve_date = pd.Timestamp(valid_rows.loc[idx, "curve_date"]).normalize()
        frame = pd.DataFrame(
            {
                "maturity_years": pd.to_numeric(maturities, errors="coerce"),
                "yield_percent": pd.to_numeric(valid_values.loc[idx], errors="coerce"),
                "curve_type": curve_type,
                "curve_date": curve_date,
            }
        ).dropna(subset=["maturity_years", "yield_percent"])
        if frame.empty:
            continue
        frames.append(frame)

    if not frames:
        raise ValueError(f"BOE {curve_type} spot curve history rows contained no valid points.")

    out = pd.concat(frames, ignore_index=True)
    out = out.sort_values(["curve_date", "maturity_years"]).drop_duplicates(
        subset=["curve_date", "maturity_years", "curve_type"],
        keep="last",
    )

    preview = valid_rows[["curve_date"]].copy()
    preview["points_available"] = valid_values.notna().sum(axis=1).values
    preview["curve_type"] = curve_type
    preview["maturity_row_idx"] = maturity_row_idx + 1
    preview["curve_date"] = pd.to_datetime(preview["curve_date"], errors="coerce")
    return out.reset_index(drop=True), preview


def fetch_dividenddata_real_yield_extension(page_url: str) -> tuple[pd.DataFrame, pd.DataFrame]:
    html = fetch_dividenddata_html(page_url)
    try:
        parser = SimpleHTMLTableParser()
        parser.feed(html)

        target_rows = None
        header = None
        maturity_idx = None
        real_yield_idx = None
        for table_rows in parser.tables:
            if len(table_rows) < 3:
                continue
            max_cols = max(len(r) for r in table_rows[:3])
            candidate_headers = []
            for depth in (1, 2, 3):
                rows_to_merge = table_rows[:depth]
                merged = []
                for col_idx in range(max_cols):
                    parts = []
                    for row in rows_to_merge:
                        if col_idx < len(row):
                            cell = str(row[col_idx]).strip()
                            if cell:
                                parts.append(cell)
                    merged.append(" ".join(parts).strip())
                candidate_headers.append((depth, merged))

            for depth, merged_header in candidate_headers:
                header_lookup = {col.lower(): idx for idx, col in enumerate(merged_header)}
                maturity_idx = next((idx for col, idx in header_lookup.items() if "time to maturity" in col), None)
                real_yield_idx = next((idx for col, idx in header_lookup.items() if "real yield" in col), None)
                if maturity_idx is not None and real_yield_idx is not None:
                    target_rows = table_rows[depth:]
                    header = merged_header
                    break
            if target_rows is not None:
                break

        if not target_rows or header is None or maturity_idx is None or real_yield_idx is None:
            raise ValueError("DividendData table with 'Time to Maturity' and 'Real Yield' columns was not found.")

        body = target_rows
        table = pd.DataFrame(body, columns=header)
        maturity_col = header[maturity_idx]
        real_yield_col = header[real_yield_idx]
        out = pd.DataFrame(
            {
                "maturity_years": table[maturity_col].map(parse_time_to_maturity),
                "yield_percent": pd.to_numeric(
                    table[real_yield_col].astype(str).str.replace("%", "", regex=False).str.replace(",", "", regex=False),
                    errors="coerce",
                ),
                "curve_type": "Real",
                "curve_date": pd.Timestamp.today().normalize(),
                "source": "DividendData",
            }
        ).dropna(subset=["maturity_years", "yield_percent"])
        out = out.sort_values("maturity_years").drop_duplicates(subset=["maturity_years"], keep="first")

        preview = table[[maturity_col, real_yield_col]].copy()
        preview.columns = ["time_to_maturity", "real_yield_raw"]
        preview["maturity_years"] = preview["time_to_maturity"].map(parse_time_to_maturity)
        preview["yield_percent"] = pd.to_numeric(
            preview["real_yield_raw"].astype(str).str.replace("%", "", regex=False).str.replace(",", "", regex=False),
            errors="coerce",
        )
        preview["source"] = "DividendData"
        return out, preview
    except Exception:
        try:
            return parse_dividenddata_text_fallback(html)
        except Exception:
            return fetch_dividenddata_short_end_detail_fallback(DIVIDENDDATA_SHORT_END_TICKERS)


def build_boe_yield_curve_diagnostics(zip_url: str, dividenddata_url: str) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    try:
        zip_bytes = fetch_boe_yield_curve_zip(zip_url)
        with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
            members = zf.namelist()
            nominal_member = pick_boe_zip_member(members, "GLC Nominal daily data current month")
            real_member = pick_boe_zip_member(members, "GLC Real daily data current month")
            nominal_points, nominal_date, nominal_preview = parse_boe_spot_curve_workbook(
                zf.read(nominal_member),
                "Nominal",
            )
            real_points, real_date, real_preview = parse_boe_spot_curve_workbook(
                zf.read(real_member),
                "Real",
            )

        nominal_points["source"] = "BOE"
        real_points["source"] = "BOE"

        extension_error = None
        extension_preview = pd.DataFrame()
        short_real_extension = pd.DataFrame(columns=real_points.columns)
        try:
            extension_points, extension_preview = fetch_dividenddata_real_yield_extension(dividenddata_url)
            if not real_points.empty and not extension_points.empty:
                real_min_maturity = float(real_points["maturity_years"].min())
                short_real_extension = extension_points[extension_points["maturity_years"] < real_min_maturity].copy()
                if not short_real_extension.empty:
                    short_real_extension["curve_date"] = real_date
        except Exception as exc:
            extension_error = f"{type(exc).__name__}: {exc}"

        real_combined = pd.concat([short_real_extension, real_points], ignore_index=True)
        real_combined = real_combined.sort_values("maturity_years").drop_duplicates(subset=["maturity_years"], keep="last")
        combined = pd.concat([nominal_points, real_combined], ignore_index=True)
        summary = pd.DataFrame(
            [
                {"metric": "Fetch status", "value": "OK"},
                {"metric": "Source URL", "value": zip_url},
                {"metric": "Extension URL", "value": dividenddata_url},
                {"metric": "Archive members", "value": int(len(members))},
                {"metric": "Nominal workbook", "value": nominal_member},
                {"metric": "Nominal latest date", "value": format_diag_date(nominal_date)},
                {"metric": "Nominal points", "value": int(len(nominal_points))},
                {"metric": "Real workbook", "value": real_member},
                {"metric": "Real latest date", "value": format_diag_date(real_date)},
                {"metric": "Real points", "value": int(len(real_points))},
                {"metric": "Real extension status", "value": "OK" if extension_error is None else "Failed"},
                {"metric": "Real extension points", "value": int(len(short_real_extension))},
                {
                    "metric": "Real extension max maturity",
                    "value": "-" if short_real_extension.empty else round(float(short_real_extension["maturity_years"].max()), 3),
                },
                {"metric": "Real extension error", "value": "-" if extension_error is None else extension_error},
            ]
        )
        preview = pd.concat([nominal_preview.tail(10), real_preview.tail(10), extension_preview.head(10)], ignore_index=True)
        preview["curve_date"] = pd.to_datetime(preview["curve_date"], errors="coerce").dt.strftime("%d/%m/%Y").fillna("")
        return combined, summary, preview
    except Exception as exc:
        summary = pd.DataFrame(
            [
                {"metric": "Fetch status", "value": "Failed"},
                {"metric": "Source URL", "value": zip_url},
                {"metric": "Extension URL", "value": dividenddata_url},
                {"metric": "Error type", "value": type(exc).__name__},
                {"metric": "Error", "value": str(exc)},
            ]
        )
        return pd.DataFrame(columns=["maturity_years", "yield_percent", "curve_type", "curve_date"]), summary, pd.DataFrame()


@st.cache_data(show_spinner=False, ttl=43200)
def build_boe_month_end_yield_curve_history(
    nominal_zip_url: str,
    real_zip_url: str,
) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    try:
        nominal_zip_bytes = fetch_boe_yield_curve_zip(nominal_zip_url)
        real_zip_bytes = fetch_boe_yield_curve_zip(real_zip_url)

        curve_frames = []
        preview_frames = []
        nominal_members_loaded: list[str] = []
        real_members_loaded: list[str] = []

        with zipfile.ZipFile(BytesIO(nominal_zip_bytes)) as zf:
            nominal_members = [name for name in zf.namelist() if name.lower().endswith(".xlsx")]
            for member in sorted(nominal_members):
                points, preview = parse_boe_spot_curve_history_workbook(zf.read(member), "Nominal")
                points["source"] = "BOE month end"
                points["workbook"] = member
                preview["workbook"] = member
                curve_frames.append(points)
                preview_frames.append(preview)
                nominal_members_loaded.append(member)

        with zipfile.ZipFile(BytesIO(real_zip_bytes)) as zf:
            real_members = [name for name in zf.namelist() if name.lower().endswith(".xlsx")]
            for member in sorted(real_members):
                points, preview = parse_boe_spot_curve_history_workbook(zf.read(member), "Real")
                points["source"] = "BOE month end"
                points["workbook"] = member
                preview["workbook"] = member
                curve_frames.append(points)
                preview_frames.append(preview)
                real_members_loaded.append(member)

        history_df = pd.concat(curve_frames, ignore_index=True) if curve_frames else pd.DataFrame()
        if not history_df.empty:
            history_df["curve_date"] = pd.to_datetime(history_df["curve_date"], errors="coerce").dt.normalize()
            history_df = history_df.dropna(subset=["curve_date", "maturity_years", "yield_percent", "curve_type"])
            history_df = history_df.sort_values(["curve_type", "curve_date", "maturity_years", "workbook"]).drop_duplicates(
                subset=["curve_type", "curve_date", "maturity_years"],
                keep="last",
            )

        preview = pd.concat(preview_frames, ignore_index=True) if preview_frames else pd.DataFrame()
        if not preview.empty:
            preview["curve_date"] = pd.to_datetime(preview["curve_date"], errors="coerce")

        nominal_dates = (
            history_df.loc[history_df["curve_type"] == "Nominal", "curve_date"].dropna().drop_duplicates().sort_values()
            if not history_df.empty
            else pd.Series(dtype="datetime64[ns]")
        )
        real_dates = (
            history_df.loc[history_df["curve_type"] == "Real", "curve_date"].dropna().drop_duplicates().sort_values()
            if not history_df.empty
            else pd.Series(dtype="datetime64[ns]")
        )
        common_dates = nominal_dates[nominal_dates.isin(set(real_dates.tolist()))] if not nominal_dates.empty and not real_dates.empty else pd.Series(dtype="datetime64[ns]")

        summary = pd.DataFrame(
            [
                {"metric": "Fetch status", "value": "OK"},
                {"metric": "Nominal source URL", "value": nominal_zip_url},
                {"metric": "Real source URL", "value": real_zip_url},
                {"metric": "Nominal workbooks", "value": int(len(nominal_members_loaded))},
                {"metric": "Real workbooks", "value": int(len(real_members_loaded))},
                {"metric": "Nominal month-end dates", "value": int(len(nominal_dates))},
                {"metric": "Real month-end dates", "value": int(len(real_dates))},
                {"metric": "Common month-end dates", "value": int(len(common_dates))},
                {
                    "metric": "Latest common month end",
                    "value": "-" if common_dates.empty else pd.Timestamp(common_dates.max()).strftime("%d/%m/%Y"),
                },
                {"metric": "Total points", "value": int(len(history_df))},
            ]
        )

        if not preview.empty:
            preview["curve_date"] = preview["curve_date"].dt.strftime("%d/%m/%Y").fillna("")

        return history_df.reset_index(drop=True), summary, preview
    except Exception as exc:
        summary = pd.DataFrame(
            [
                {"metric": "Fetch status", "value": "Failed"},
                {"metric": "Nominal source URL", "value": nominal_zip_url},
                {"metric": "Real source URL", "value": real_zip_url},
                {"metric": "Error type", "value": type(exc).__name__},
                {"metric": "Error", "value": str(exc)},
            ]
        )
        return (
            pd.DataFrame(columns=["maturity_years", "yield_percent", "curve_type", "curve_date", "source", "workbook"]),
            summary,
            pd.DataFrame(),
        )


@st.cache_data(show_spinner=False, ttl=21600)
def build_global_yield_curve_diagnostics(
    base_url: str,
    country_specs: list[tuple[str, str]],
) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    summary_rows = [
        {"metric": "Fetch status", "value": "OK"},
        {"metric": "Base URL", "value": base_url},
        {"metric": "Countries requested", "value": int(len(country_specs))},
    ]
    preview_frames = []
    curve_frames = []
    success_count = 0

    for country_name, path in country_specs:
        page_url = f"{base_url.rstrip('/')}{path}"
        try:
            html = fetch_worldgovernmentbonds_html(page_url)
            global_vars = extract_worldgovernmentbonds_js_global_vars(html)
            payload = fetch_worldgovernmentbonds_country_payload(page_url, global_vars)
            table_html = str(payload.get("mainTable", "")).strip()
            if not table_html:
                raise ValueError(f"WorldGovernmentBonds API returned no mainTable HTML for {country_name}.")

            country_curve, country_preview = parse_worldgovernmentbonds_yield_curve_table(table_html, country_name, page_url)
            last_data_desc = str(payload.get("lastDataValDesc", "")).strip()
            curve_date = pd.to_datetime(last_data_desc, errors="coerce", dayfirst=True)
            if pd.notna(curve_date):
                country_curve["curve_date"] = pd.Timestamp(curve_date).normalize()
            curve_frames.append(country_curve)
            if pd.notna(curve_date):
                country_preview["curve_date"] = pd.Timestamp(curve_date).strftime("%d/%m/%Y")
            else:
                country_preview["curve_date"] = last_data_desc
            preview_frames.append(country_preview)
            summary_rows.append({"metric": f"{country_name} status", "value": "OK"})
            summary_rows.append({"metric": f"{country_name} points", "value": int(len(country_curve))})
            summary_rows.append({"metric": f"{country_name} latest date", "value": last_data_desc or "-"})
            success_count += 1
        except Exception as exc:
            summary_rows.append({"metric": f"{country_name} status", "value": f"{type(exc).__name__}: {exc}"})
            summary_rows.append({"metric": f"{country_name} points", "value": 0})
            summary_rows.append({"metric": f"{country_name} latest date", "value": "-"})

    if success_count == 0:
        summary_rows[0]["value"] = "Failed"

    combined = (
        pd.concat(curve_frames, ignore_index=True)
        if curve_frames
        else pd.DataFrame(columns=["country", "maturity_label", "maturity_years", "yield_percent", "curve_date", "source_url"])
    )
    preview = (
        pd.concat(preview_frames, ignore_index=True)
        if preview_frames
        else pd.DataFrame(columns=["country", "maturity_label", "maturity_years", "yield_percent", "source_url"])
    )
    return combined, pd.DataFrame(summary_rows), preview


def build_asset_coverage_table(
    mapping: pd.DataFrame,
    monthly_levels: dict[str, pd.Series],
    stitched_series_map: dict[str, pd.Series],
    chart_series_map: dict[str, pd.Series],
    live_diag: pd.DataFrame,
    chart_diag: pd.DataFrame,
) -> pd.DataFrame:
    mapping_view = mapping.drop_duplicates(subset=["asset_class"]).copy()
    rows = []

    for _, row in mapping_view.iterrows():
        asset_class = str(row.get("asset_class", "")).strip()
        if not asset_class:
            continue

        index_summary = build_series_summary(monthly_levels.get(asset_class))
        dashboard_summary = build_series_summary(stitched_series_map.get(asset_class))
        chart_summary = build_series_summary(chart_series_map.get(asset_class))

        live_row = (
            live_diag[live_diag["asset_class"] == asset_class].iloc[0].to_dict()
            if live_diag is not None and not live_diag.empty and asset_class in live_diag["asset_class"].values
            else {}
        )
        chart_row = (
            chart_diag[chart_diag["asset_class"] == asset_class].iloc[0].to_dict()
            if chart_diag is not None and not chart_diag.empty and asset_class in chart_diag["asset_class"].values
            else {}
        )

        rows.append(
            {
                "asset_class": asset_class,
                "index_name": row.get("index_name", ""),
                "live_fund_primary": row.get("live_fund_primary", ""),
                "live_fund_secondary": row.get("live_fund_secondary", ""),
                "selected_ticker": live_row.get("selected_ticker", ""),
                "dashboard_series_type": live_row.get("series_type", ""),
                "chart_series_type": chart_row.get("series_type", ""),
                "index_points": index_summary["points"],
                "index_first_date": index_summary["first_date"],
                "index_last_date": index_summary["last_date"],
                "dashboard_points": dashboard_summary["points"],
                "dashboard_last_date": dashboard_summary["last_date"],
                "chart_points": chart_summary["points"],
                "chart_first_date": chart_summary["first_date"],
                "chart_last_date": chart_summary["last_date"],
            }
        )

    out = pd.DataFrame(rows)
    if out.empty:
        return out

    for col in [
        "index_first_date",
        "index_last_date",
        "dashboard_last_date",
        "chart_first_date",
        "chart_last_date",
    ]:
        out[col] = pd.to_datetime(out[col], errors="coerce").dt.strftime("%d/%m/%Y").fillna("")

    return order_asset_rows(out)


def build_mapping_diagnostics_table(mapping: pd.DataFrame, ts: pd.DataFrame) -> pd.DataFrame:
    rows = []
    ts_columns = {str(c).strip() for c in ts.columns}
    asset_counts = mapping["asset_class"].value_counts(dropna=False).to_dict() if "asset_class" in mapping.columns else {}

    for _, row in mapping.iterrows():
        index_name = str(row.get("index_name", "")).strip()
        asset_class = str(row.get("asset_class", "")).strip()
        primary = normalise_ticker(row.get("live_fund_primary", ""))
        secondary = normalise_ticker(row.get("live_fund_secondary", ""))

        issues = []
        if not asset_class:
            issues.append("Missing asset_class")
        if not index_name:
            issues.append("Missing index_name")
        elif index_name not in ts_columns:
            issues.append("Index column missing from time_series")
        if asset_class and asset_counts.get(asset_class, 0) > 1:
            issues.append("Duplicate asset_class in mapping")
        if not primary and not secondary:
            issues.append("No live ticker")

        rows.append(
            {
                "asset_class": asset_class,
                "index_name": index_name,
                "live_fund_primary": primary,
                "live_fund_secondary": secondary,
                "status": "Issue" if issues else "OK",
                "issues": "; ".join(issues) if issues else "",
            }
        )

    return order_asset_rows(pd.DataFrame(rows))


def build_live_price_diagnostics(prices_df: pd.DataFrame) -> pd.DataFrame:
    if prices_df is None or prices_df.empty:
        return pd.DataFrame(columns=["ticker", "points", "first_date", "last_date", "last_price"])

    rows = []
    for ticker in sorted(prices_df.columns):
        series = pd.to_numeric(prices_df[ticker], errors="coerce").dropna().sort_index()
        if series.empty:
            rows.append(
                {
                    "ticker": ticker,
                    "points": 0,
                    "first_date": "",
                    "last_date": "",
                    "last_price": np.nan,
                }
            )
            continue
        rows.append(
            {
                "ticker": ticker,
                "points": int(len(series)),
                "first_date": format_diag_date(series.index.min()),
                "last_date": format_diag_date(series.index.max()),
                "last_price": round(float(series.iloc[-1]), 4),
            }
        )
    return pd.DataFrame(rows)


def build_return_anchor_table(
    series_map: dict[str, pd.Series],
    end_date: pd.Timestamp,
    period_keys: list[str],
    whole_period_start: pd.Timestamp,
) -> pd.DataFrame:
    rows = []
    years_map = {"1Y": 1, "3Y": 3, "5Y": 5, "7Y": 7, "10Y": 10, "15Y": 15, "20Y": 20, "25Y": 25}

    for asset_class, series in series_map.items():
        s = series.dropna().sort_index()
        if s.empty:
            continue

        end_anchor_date = pd.Timestamp(s.index.max()) if not s[s.index <= end_date].empty else pd.NaT
        for period_key in period_keys:
            if period_key == "YTD":
                base_date, base_level = nearest_level_on_or_before(s, pd.Timestamp(end_date.year - 1, 12, 31))
            elif period_key == "Period":
                base_date, base_level = nearest_level_on_or_before(s, whole_period_start)
            else:
                years = years_map.get(period_key)
                base_date, base_level = nearest_level_on_or_before(s, end_date - pd.DateOffset(years=years))

            rows.append(
                {
                    "asset_class": asset_class,
                    "period": period_key,
                    "base_date_used": format_diag_date(base_date),
                    "end_date_used": format_diag_date(end_anchor_date),
                    "base_level": np.nan if base_level is None else round(float(base_level), 6),
                    "end_level": round(float(s[s.index <= end_date].iloc[-1]), 6) if not s[s.index <= end_date].empty else np.nan,
                    "return_value": calc_period_return(s, end_date, period_key, whole_period_start=whole_period_start),
                }
            )

    out = pd.DataFrame(rows)
    if not out.empty:
        out["return_value"] = out["return_value"].map(lambda x: np.nan if pd.isna(x) else round(x * 100, 4))
    return order_asset_rows(out)


def years_between(start_date: pd.Timestamp, end_date: pd.Timestamp) -> float:
    return max((pd.Timestamp(end_date) - pd.Timestamp(start_date)).days / 365.25, 0)


def nearest_level_on_or_before(series: pd.Series, date: pd.Timestamp) -> tuple[pd.Timestamp | None, float | None]:
    s = series.dropna().sort_index()
    s = s[s.index <= date]
    if s.empty:
        return None, None
    return pd.Timestamp(s.index[-1]), float(s.iloc[-1])


def calc_period_return(series: pd.Series, end_date: pd.Timestamp, period_key: str, whole_period_start: pd.Timestamp | None = None) -> float:
    s = series.dropna().sort_index()
    s = s[s.index <= end_date]
    if s.empty:
        return np.nan

    end_level = float(s.iloc[-1])

    if period_key == "YTD":
        _, base_level = nearest_level_on_or_before(s, pd.Timestamp(end_date.year - 1, 12, 31))
        return (end_level / base_level - 1) if base_level and base_level > 0 else np.nan

    if period_key == "Period":
        if whole_period_start is not None:
            base_date, base_level = nearest_level_on_or_before(s, whole_period_start)
        else:
            base_date, base_level = pd.Timestamp(s.index.min()), float(s.iloc[0])

        if base_level is None or base_level <= 0:
            return np.nan

        years = years_between(base_date, s.index[-1])
        growth = end_level / base_level
        return annualised_return_from_growth(growth, years)

    years_map = {"1Y": 1, "3Y": 3, "5Y": 5, "7Y": 7, "10Y": 10, "15Y": 15, "20Y": 20, "25Y": 25}
    years = years_map.get(period_key)
    if years is None:
        return np.nan

    _, base_level = nearest_level_on_or_before(s, end_date - pd.DateOffset(years=years))
    if base_level is None or base_level <= 0:
        return np.nan

    growth = end_level / base_level
    return annualised_return_from_growth(growth, years)


@st.cache_data(show_spinner=False)
def load_data(file_path: str, file_mtime: float) -> tuple[pd.DataFrame, pd.DataFrame]:
    ts = pd.read_excel(file_path, sheet_name=TIME_SERIES_SHEET)
    mapping = pd.read_excel(file_path, sheet_name=MAPPING_SHEET)

    ts.columns = [str(c).strip() for c in ts.columns]
    mapping.columns = [str(c).strip() for c in mapping.columns]

    ts = ts.copy()
    mapping = mapping.copy()

    ts.iloc[:, 0] = pd.to_datetime(ts.iloc[:, 0], errors="coerce")
    ts = ts.rename(columns={ts.columns[0]: "Date"})
    ts = ts.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)

    normalized_mapping_cols = {str(c).strip().lower(): str(c).strip() for c in mapping.columns}
    rename_map = {}
    explicit_name_map = {
        "index_name": "index_name",
        "asset_class": "asset_class",
        "live_fund_primary": "live_fund_primary",
        "live_fund_secondary": "live_fund_secondary",
        "growth/defensive": "growth_defensive",
        "growth_defensive": "growth_defensive",
    }
    for source_name, target_name in explicit_name_map.items():
        if source_name in normalized_mapping_cols:
            rename_map[normalized_mapping_cols[source_name]] = target_name

    if "index_name" not in rename_map.values() and len(mapping.columns) >= 1:
        rename_map[mapping.columns[0]] = "index_name"
    if "asset_class" not in rename_map.values() and len(mapping.columns) >= 2:
        rename_map[mapping.columns[1]] = "asset_class"
    if "live_fund_primary" not in rename_map.values() and len(mapping.columns) >= 3:
        rename_map[mapping.columns[2]] = "live_fund_primary"
    if "live_fund_secondary" not in rename_map.values() and len(mapping.columns) >= 4:
        rename_map[mapping.columns[3]] = "live_fund_secondary"
    mapping = mapping.rename(columns=rename_map)

    if "index_name" in mapping.columns:
        mapping["index_name"] = mapping["index_name"].astype(str).str.strip()

    if "asset_class" in mapping.columns:
        mapping["asset_class"] = mapping["asset_class"].astype(str).str.strip().replace(ASSET_CLASS_ALIASES)

    if "live_fund_primary" in mapping.columns:
        mapping["live_fund_primary"] = mapping["live_fund_primary"].map(normalise_ticker)

    if "live_fund_secondary" in mapping.columns:
        mapping["live_fund_secondary"] = mapping["live_fund_secondary"].map(normalise_ticker)

    if "growth_defensive" in mapping.columns:
        mapping["growth_defensive"] = pd.to_numeric(mapping["growth_defensive"], errors="coerce")

    valid_rows = mapping[mapping["index_name"].isin(ts.columns)].copy()
    return ts, valid_rows


@st.cache_data(show_spinner=False)
def load_regions_data(file_path: str, file_mtime: float) -> pd.DataFrame:
    try:
        regions = pd.read_excel(file_path, sheet_name=REGIONS_SHEET)
    except Exception:
        return pd.DataFrame(columns=["ticker", "name", "investment_area", "available", "country_flag", "region_flag"])

    regions.columns = [str(c).strip() for c in regions.columns]
    regions = regions.copy()

    normalized_cols = {str(c).strip().lower(): str(c).strip() for c in regions.columns}
    rename_map = {}
    explicit_name_map = {
        "ticker": "ticker",
        "name": "name",
        "investment_area": "investment_area",
        "available": "available",
        "country": "country_flag",
        "region": "region_flag",
    }
    for source_name, target_name in explicit_name_map.items():
        if source_name in normalized_cols:
            rename_map[normalized_cols[source_name]] = target_name
    regions = regions.rename(columns=rename_map)

    if "ticker" not in regions.columns:
        return pd.DataFrame(columns=["ticker", "name", "investment_area", "available", "country_flag", "region_flag"])

    regions["ticker"] = regions["ticker"].map(normalise_ticker)
    if "name" in regions.columns:
        regions["name"] = regions["name"].astype(str).str.strip()
    else:
        regions["name"] = regions["ticker"]
    if "investment_area" in regions.columns:
        regions["investment_area"] = regions["investment_area"].astype(str).str.strip()
    else:
        regions["investment_area"] = regions["name"]

    if "available" in regions.columns:
        regions["available"] = (
            regions["available"].astype(str).str.strip().str.lower().isin({"yes", "y", "true", "1"})
        )
    else:
        regions["available"] = True

    regions["country_flag"] = pd.to_numeric(regions.get("country_flag", 0), errors="coerce").fillna(0).astype(int)
    regions["region_flag"] = pd.to_numeric(regions.get("region_flag", 0), errors="coerce").fillna(0).astype(int)

    regions = regions[regions["ticker"].astype(str).str.len() > 0].copy()
    regions = regions.drop_duplicates(subset=["ticker"], keep="first")
    return regions[["ticker", "name", "investment_area", "available", "country_flag", "region_flag"]]


@st.cache_data(show_spinner=False)
def load_sectors_data(file_path: str, file_mtime: float) -> pd.DataFrame:
    try:
        sectors = pd.read_excel(file_path, sheet_name=SECTORS_SHEET)
    except Exception:
        return pd.DataFrame(columns=["ticker", "name", "sector", "available"])

    sectors.columns = [str(c).strip() for c in sectors.columns]
    sectors = sectors.copy()

    normalized_cols = {str(c).strip().lower(): str(c).strip() for c in sectors.columns}
    rename_map = {}
    explicit_name_map = {
        "ticker": "ticker",
        "name": "name",
        "sector": "sector",
        "available": "available",
    }
    for source_name, target_name in explicit_name_map.items():
        if source_name in normalized_cols:
            rename_map[normalized_cols[source_name]] = target_name
    sectors = sectors.rename(columns=rename_map)

    if "ticker" not in sectors.columns:
        return pd.DataFrame(columns=["ticker", "name", "sector", "available"])

    sectors["ticker"] = sectors["ticker"].map(normalise_ticker)
    sectors["name"] = sectors.get("name", sectors["ticker"]).astype(str).str.strip()
    sectors["sector"] = sectors.get("sector", sectors["name"]).astype(str).str.strip()
    if "available" in sectors.columns:
        sectors["available"] = (
            sectors["available"].astype(str).str.strip().str.lower().isin({"yes", "y", "true", "1"})
        )
    else:
        sectors["available"] = True

    sectors = sectors[sectors["ticker"].astype(str).str.len() > 0].copy()
    sectors = sectors.drop_duplicates(subset=["ticker"], keep="first")
    return sectors[["ticker", "name", "sector", "available"]]


@st.cache_data(show_spinner=False)
def load_factors_data(file_path: str, file_mtime: float) -> pd.DataFrame:
    try:
        factors = pd.read_excel(file_path, sheet_name=FACTORS_SHEET)
    except Exception:
        return pd.DataFrame(columns=["ticker", "name", "region", "size_style", "value_style", "label", "available"])

    factors.columns = [str(c).strip() for c in factors.columns]
    factors = factors.copy()

    normalized_cols = {str(c).strip().lower(): str(c).strip() for c in factors.columns}
    rename_map = {}
    explicit_name_map = {
        "ticker": "ticker",
        "name": "name",
        "region": "region",
        "size": "size_style",
        "size style": "size_style",
        "value": "value_style",
        "value style": "value_style",
        "final": "label",
        "available": "available",
    }
    for source_name, target_name in explicit_name_map.items():
        if source_name in normalized_cols:
            rename_map[normalized_cols[source_name]] = target_name
    factors = factors.rename(columns=rename_map)

    if "ticker" not in factors.columns:
        return pd.DataFrame(columns=["ticker", "name", "region", "size_style", "value_style", "label", "available"])

    factors["ticker"] = factors["ticker"].map(normalise_ticker)
    factors["name"] = factors.get("name", factors["ticker"]).astype(str).str.strip()
    factors["region"] = factors.get("region", "").map(canonical_factor_region)
    factors["size_style"] = factors.get("size_style", "").map(canonical_factor_size_style)
    factors["value_style"] = factors.get("value_style", "").map(canonical_factor_value_style)
    factors["label"] = factors.get("label", factors["name"]).astype(str).str.strip()

    if "available" in factors.columns:
        factors["available"] = (
            factors["available"].astype(str).str.strip().str.lower().isin({"yes", "y", "true", "1"})
        )
    else:
        factors["available"] = True

    factors = factors[factors["ticker"].astype(str).str.len() > 0].copy()
    factors = factors.drop_duplicates(subset=["ticker"], keep="first")
    return factors[["ticker", "name", "region", "size_style", "value_style", "label", "available"]]


@st.cache_data(show_spinner=False, ttl=43200)
def fetch_yf_quote_currencies(tickers: tuple[str, ...]) -> dict[str, str]:
    out: dict[str, str] = {}
    for ticker in tickers:
        norm_ticker = normalise_ticker(ticker)
        if not norm_ticker:
            continue
        currency = ""
        try:
            with contextlib.redirect_stdout(StringIO()), contextlib.redirect_stderr(StringIO()):
                fast = yf.Ticker(norm_ticker).fast_info
            if fast:
                currency = str(fast.get("currency", "") or "").upper().strip()
        except Exception:
            currency = ""
        out[norm_ticker] = currency
    return out


@st.cache_data(show_spinner=False)
def build_monthly_index_levels(ts: pd.DataFrame, mapping: pd.DataFrame) -> dict[str, pd.Series]:
    output = {}
    for _, row in mapping.iterrows():
        asset_class = row["asset_class"]
        index_name = row["index_name"]

        if index_name not in ts.columns:
            continue

        returns = standardise_series(ts[index_name])
        series = pd.Series(returns.values, index=ts["Date"], name=asset_class).dropna().sort_index()
        if series.empty:
            continue

        levels = (1 + series).cumprod()
        levels.name = asset_class
        output[asset_class] = levels

    return output


@st.cache_data(show_spinner=False, ttl=43200)
def fetch_yf_prices(tickers: tuple[str, ...], start_date: str) -> pd.DataFrame:
    tickers = tuple(sorted({normalise_ticker(t) for t in tickers if normalise_ticker(t)}))
    if not tickers:
        return pd.DataFrame()

    with contextlib.redirect_stdout(StringIO()), contextlib.redirect_stderr(StringIO()):
        data = yf.download(
            list(tickers),
            start=start_date,
            progress=False,
            auto_adjust=True,
            actions=False,
            threads=False,
            group_by="ticker",
        )

    if data is None or len(data) == 0:
        return pd.DataFrame()

    close_frames = []
    if isinstance(data.columns, pd.MultiIndex):
        for ticker in tickers:
            if ticker in data.columns.get_level_values(0):
                sub = data[ticker].copy()
                if "Close" in sub.columns:
                    close_frames.append(pd.to_numeric(sub["Close"], errors="coerce").rename(ticker))
    else:
        if "Close" in data.columns and len(tickers) == 1:
            close_frames.append(pd.to_numeric(data["Close"], errors="coerce").rename(tickers[0]))

    if not close_frames:
        return pd.DataFrame()

    out = pd.concat(close_frames, axis=1)
    out.index = pd.to_datetime(out.index).tz_localize(None)
    return out.sort_index()


def get_price_series(prices_df: pd.DataFrame, ticker: str) -> pd.Series:
    if prices_df.empty or ticker not in prices_df.columns:
        return pd.Series(dtype=float)
    return pd.to_numeric(prices_df[ticker], errors="coerce").dropna().sort_index()


def pick_live_ticker_for_asset(row: pd.Series, prices_df: pd.DataFrame) -> tuple[str, str, str, pd.Series]:
    primary = normalise_ticker(row.get("live_fund_primary", ""))
    secondary = normalise_ticker(row.get("live_fund_secondary", ""))

    primary_series = get_price_series(prices_df, primary) if primary else pd.Series(dtype=float)
    secondary_series = get_price_series(prices_df, secondary) if secondary else pd.Series(dtype=float)

    if not primary_series.empty:
        return primary, "primary", "Primary available in yfinance", primary_series
    if not secondary_series.empty:
        return secondary, "secondary", "Primary unavailable; using secondary", secondary_series
    if primary and secondary:
        return "", "none", "Neither primary nor secondary returned price history", pd.Series(dtype=float)
    if primary:
        return "", "none", "Primary returned no price history and no secondary provided", pd.Series(dtype=float)
    if secondary:
        return "", "none", "No primary provided; secondary returned no price history", pd.Series(dtype=float)
    return "", "none", "No live fund tickers provided", pd.Series(dtype=float)


def build_stitched_asset_series(monthly_levels: dict[str, pd.Series], mapping: pd.DataFrame, prices_df: pd.DataFrame) -> tuple[dict[str, pd.Series], pd.DataFrame]:
    stitched = {}
    diag_rows = []
    deduped_mapping = mapping.drop_duplicates(subset=["asset_class"], keep="last").copy()

    for asset_class in sorted(deduped_mapping["asset_class"].dropna().astype(str).unique()):
        row_match = deduped_mapping[deduped_mapping["asset_class"] == asset_class]

        if asset_class not in monthly_levels:
            diag_rows.append(
                {
                    "asset_class": asset_class,
                    "live_fund_primary": "",
                    "live_fund_secondary": "",
                    "selected_ticker": "",
                    "selected_source": "none",
                    "series_type": "missing",
                    "index_last_date": pd.NaT,
                    "live_first_date": pd.NaT,
                    "live_last_date": pd.NaT,
                    "stitch_anchor_date": pd.NaT,
                    "note": "No index series found for asset class",
                }
            )
            continue

        row = row_match.iloc[0] if not row_match.empty else pd.Series(dtype=object)
        primary = normalise_ticker(row.get("live_fund_primary", ""))
        secondary = normalise_ticker(row.get("live_fund_secondary", ""))

        index_levels = monthly_levels[asset_class].dropna().sort_index()
        index_last_date = pd.Timestamp(index_levels.index.max())
        anchor_level = float(index_levels.loc[index_last_date])

        selected_ticker, selected_source, note, live_prices = pick_live_ticker_for_asset(row, prices_df)

        if live_prices.empty:
            stitched[asset_class] = index_levels.copy()
            diag_rows.append(
                {
                    "asset_class": asset_class,
                    "live_fund_primary": primary,
                    "live_fund_secondary": secondary,
                    "selected_ticker": selected_ticker,
                    "selected_source": selected_source,
                    "series_type": "index_only",
                    "index_last_date": index_last_date,
                    "live_first_date": pd.NaT,
                    "live_last_date": pd.NaT,
                    "stitch_anchor_date": pd.NaT,
                    "note": note,
                }
            )
            continue

        live_prices = live_prices.dropna().sort_index()
        live_first_date = pd.Timestamp(live_prices.index.min())
        live_last_date = pd.Timestamp(live_prices.index.max())

        live_anchor_date, live_anchor_price = nearest_level_on_or_before(live_prices, index_last_date)
        if live_anchor_date is None or live_anchor_price is None:
            stitched[asset_class] = index_levels.copy()
            diag_rows.append(
                {
                    "asset_class": asset_class,
                    "live_fund_primary": primary,
                    "live_fund_secondary": secondary,
                    "selected_ticker": selected_ticker,
                    "selected_source": selected_source,
                    "series_type": "index_only",
                    "index_last_date": index_last_date,
                    "live_first_date": live_first_date,
                    "live_last_date": live_last_date,
                    "stitch_anchor_date": index_last_date,
                    "note": f"{note}. Live series has no price on or before the final index date, so no stitch was applied",
                }
            )
            continue

        live_extension = live_prices[live_prices.index > index_last_date].copy()
        if live_extension.empty:
            stitched[asset_class] = index_levels.copy()
            diag_rows.append(
                {
                    "asset_class": asset_class,
                    "live_fund_primary": primary,
                    "live_fund_secondary": secondary,
                    "selected_ticker": selected_ticker,
                    "selected_source": selected_source,
                    "series_type": "index_only",
                    "index_last_date": index_last_date,
                    "live_first_date": live_first_date,
                    "live_last_date": live_last_date,
                    "stitch_anchor_date": index_last_date,
                    "note": f"{note}. Live series exists but has no data after the final index date",
                }
            )
            continue

        live_levels = anchor_level * (live_extension / live_anchor_price)
        stitched_series = pd.concat([index_levels, live_levels])
        stitched_series = stitched_series[~stitched_series.index.duplicated(keep="last")].sort_index()
        stitched[asset_class] = stitched_series

        diag_rows.append(
            {
                "asset_class": asset_class,
                "live_fund_primary": primary,
                "live_fund_secondary": secondary,
                "selected_ticker": selected_ticker,
                "selected_source": selected_source,
                "series_type": "stitched",
                "index_last_date": index_last_date,
                "live_first_date": live_first_date,
                "live_last_date": live_last_date,
                "stitch_anchor_date": index_last_date,
                "note": f"{note}. Index history retained through its final date, then extended with live adjusted-close history",
            }
        )

    return stitched, pd.DataFrame(diag_rows)


def build_chart_preferred_series(monthly_levels: dict[str, pd.Series], mapping: pd.DataFrame, prices_df: pd.DataFrame) -> tuple[dict[str, pd.Series], pd.DataFrame]:
    chart_series_map = {}
    diag_rows = []
    deduped_mapping = mapping.drop_duplicates(subset=["asset_class"], keep="last").copy()

    for asset_class in sorted(deduped_mapping["asset_class"].dropna().astype(str).unique()):
        row_match = deduped_mapping[deduped_mapping["asset_class"] == asset_class]
        row = row_match.iloc[0] if not row_match.empty else pd.Series(dtype=object)

        primary = normalise_ticker(row.get("live_fund_primary", ""))
        secondary = normalise_ticker(row.get("live_fund_secondary", ""))

        index_levels = monthly_levels.get(asset_class, pd.Series(dtype=float)).dropna().sort_index()
        selected_ticker, selected_source, note, live_prices = pick_live_ticker_for_asset(row, prices_df)

        if live_prices.empty and index_levels.empty:
            diag_rows.append(
                {
                    "asset_class": asset_class,
                    "live_fund_primary": primary,
                    "live_fund_secondary": secondary,
                    "selected_ticker": selected_ticker,
                    "selected_source": selected_source,
                    "series_type": "missing",
                    "index_last_date": pd.NaT,
                    "live_first_date": pd.NaT,
                    "live_last_date": pd.NaT,
                    "stitch_anchor_date": pd.NaT,
                    "note": "No index or live series available for chart series",
                }
            )
            continue

        if live_prices.empty and not index_levels.empty:
            chart_series_map[asset_class] = index_levels.copy()
            diag_rows.append(
                {
                    "asset_class": asset_class,
                    "live_fund_primary": primary,
                    "live_fund_secondary": secondary,
                    "selected_ticker": selected_ticker,
                    "selected_source": selected_source,
                    "series_type": "index_fallback",
                    "index_last_date": pd.Timestamp(index_levels.index.max()),
                    "live_first_date": pd.NaT,
                    "live_last_date": pd.NaT,
                    "stitch_anchor_date": pd.NaT,
                    "note": "No usable live history. Monthly index history used for chart series",
                }
            )
            continue

        live_series = live_prices.dropna().sort_index()
        live_first_date = pd.Timestamp(live_series.index.min())
        live_last_date = pd.Timestamp(live_series.index.max())

        if index_levels.empty:
            chart_series_map[asset_class] = live_series.copy()
            diag_rows.append(
                {
                    "asset_class": asset_class,
                    "live_fund_primary": primary,
                    "live_fund_secondary": secondary,
                    "selected_ticker": selected_ticker,
                    "selected_source": selected_source,
                    "series_type": "live_only",
                    "index_last_date": pd.NaT,
                    "live_first_date": live_first_date,
                    "live_last_date": live_last_date,
                    "stitch_anchor_date": live_first_date,
                    "note": f"{note}. No index history available, so daily live history is used on its own",
                }
            )
            continue

        live_anchor_date = live_first_date
        live_anchor_value = float(live_series.iloc[0])

        index_anchor_date, index_anchor_level = nearest_level_on_or_before(index_levels, live_anchor_date)
        if index_anchor_date is None or index_anchor_level is None or index_anchor_level <= 0:
            chart_series_map[asset_class] = live_series.copy()
            diag_rows.append(
                {
                    "asset_class": asset_class,
                    "live_fund_primary": primary,
                    "live_fund_secondary": secondary,
                    "selected_ticker": selected_ticker,
                    "selected_source": selected_source,
                    "series_type": "live_preferred",
                    "index_last_date": pd.Timestamp(index_levels.index.max()),
                    "live_first_date": live_first_date,
                    "live_last_date": live_last_date,
                    "stitch_anchor_date": live_anchor_date,
                    "note": f"{note}. Daily live history used for chart series; no valid index anchor before live start so no backward stitch applied",
                }
            )
            continue

        historical_index = index_levels[index_levels.index < live_anchor_date].copy()
        scaled_historical_index = historical_index * (live_anchor_value / index_anchor_level)

        combined = pd.concat([scaled_historical_index, live_series])
        combined = combined[~combined.index.duplicated(keep="last")].sort_index()
        combined.name = asset_class
        chart_series_map[asset_class] = combined

        diag_rows.append(
            {
                "asset_class": asset_class,
                "live_fund_primary": primary,
                "live_fund_secondary": secondary,
                "selected_ticker": selected_ticker,
                "selected_source": selected_source,
                "series_type": "chart_stitched",
                "index_last_date": pd.Timestamp(index_levels.index.max()),
                "live_first_date": live_first_date,
                "live_last_date": live_last_date,
                "stitch_anchor_date": live_anchor_date,
                "note": (
                    f"{note}. Monthly index history stitched in before the first live daily observation; "
                    "daily live history is prioritised from that point onward"
                ),
            }
        )

    return chart_series_map, pd.DataFrame(diag_rows)


def calc_horizon_returns_from_levels(stitched_series_map: dict[str, pd.Series], end_date: pd.Timestamp, period_keys: list[str]) -> pd.DataFrame:
    rows = []
    for asset_class, series in stitched_series_map.items():
        row = {"asset_class": asset_class}
        for period_key in period_keys:
            row[period_key] = calc_period_return(series, end_date, period_key)
        rows.append(row)
    return pd.DataFrame(rows)


def calc_whole_period_returns(stitched_series_map: dict[str, pd.Series], end_date: pd.Timestamp, whole_period_start: pd.Timestamp) -> pd.DataFrame:
    rows = []
    for asset_class, series in stitched_series_map.items():
        rows.append(
            {
                "asset_class": asset_class,
                "Period": calc_period_return(series, end_date, "Period", whole_period_start=whole_period_start),
            }
        )
    return pd.DataFrame(rows)


def merge_return_tables(*dfs: pd.DataFrame) -> pd.DataFrame:
    out = None
    for df in dfs:
        out = df if out is None else out.merge(df, on="asset_class", how="outer")
    return out if out is not None else pd.DataFrame(columns=["asset_class"])


def order_asset_rows(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty or "asset_class" not in df.columns:
        return df
    order_map = {name: idx for idx, name in enumerate(DEFAULT_ASSET_ORDER)}
    return (
        df.assign(_order=df["asset_class"].map(order_map).fillna(9999))
        .sort_values(["_order", "asset_class"])
        .drop(columns="_order")
    )


def build_real_level_series(nominal_series: pd.Series, inflation_levels: pd.Series) -> pd.Series:
    nominal = nominal_series.dropna().sort_index()
    infl = inflation_levels.dropna().sort_index()
    if nominal.empty or infl.empty:
        return pd.Series(dtype=float)
    aligned_infl = infl.reindex(nominal.index, method="ffill")
    real = nominal / aligned_infl
    return real.dropna().sort_index()


def build_growth_of_wealth_df(
    chart_series_map: dict[str, pd.Series],
    selected_assets: list[str],
    end_date: pd.Timestamp,
    period_key: str,
    is_real_mode: bool,
    inflation_levels: pd.Series | None,
) -> pd.DataFrame:
    rows = []

    for asset in selected_assets:
        if asset not in chart_series_map:
            continue

        series = chart_series_map[asset].dropna().sort_index()
        series = series[series.index <= end_date]
        if series.empty:
            continue

        if is_real_mode:
            if inflation_levels is None or inflation_levels.dropna().empty:
                continue
            series = build_real_level_series(series, inflation_levels)
            series = series[series.index <= end_date]
            if series.empty:
                continue

        if period_key == "YTD":
            base_date, base_level = nearest_level_on_or_before(series, pd.Timestamp(end_date.year - 1, 12, 31))
        elif period_key == "Custom":
            if len(series) < 2:
                continue
            base_date = pd.Timestamp(series.index.min())
            base_level = float(series.iloc[0])
        else:
            years = int(period_key.replace("Y", ""))
            base_date, base_level = nearest_level_on_or_before(series, end_date - pd.DateOffset(years=years))

        if base_date is None or base_level is None or base_level <= 0:
            continue

        chart_series = series[series.index >= base_date].copy() / base_level
        if chart_series.empty:
            continue

        chart_df = chart_series.reset_index()
        chart_df.columns = ["Date", "Growth"]
        chart_df["asset_class"] = asset
        rows.append(chart_df)

    return pd.concat(rows, ignore_index=True) if rows else pd.DataFrame(columns=["Date", "Growth", "asset_class"])


def build_monthly_level_window(series: pd.Series, end_date: pd.Timestamp, years: int) -> pd.Series:
    s = series.dropna().sort_index()
    s = s[s.index <= end_date]
    if s.empty:
        return pd.Series(dtype=float)

    monthly = s.resample("ME").last().dropna()
    if monthly.empty:
        return pd.Series(dtype=float)

    base_date, _ = nearest_level_on_or_before(monthly, end_date - pd.DateOffset(years=years))
    if base_date is None:
        return pd.Series(dtype=float)

    return monthly[monthly.index >= base_date].copy()


def build_monthly_level_window_from_start(series: pd.Series, end_date: pd.Timestamp, start_date: pd.Timestamp) -> pd.Series:
    s = series.dropna().sort_index()
    s = s[s.index <= end_date]
    if s.empty:
        return pd.Series(dtype=float)

    monthly = s.resample("ME").last().dropna()
    if monthly.empty:
        return pd.Series(dtype=float)

    base_date, _ = nearest_level_on_or_before(monthly, start_date)
    if base_date is None:
        return pd.Series(dtype=float)

    return monthly[monthly.index >= base_date].copy()


def calc_annualised_volatility(series: pd.Series, end_date: pd.Timestamp, years: int) -> float:
    window = build_monthly_level_window(series, end_date, years)
    monthly_returns = window.pct_change().dropna()
    if monthly_returns.empty or len(monthly_returns) < 2:
        return np.nan
    return float(monthly_returns.std(ddof=1) * np.sqrt(12))


def build_monthly_return_window(series: pd.Series, end_date: pd.Timestamp, years: int) -> pd.Series:
    window = build_monthly_level_window(series, end_date, years)
    if window.empty:
        return pd.Series(dtype=float)
    return window.pct_change().dropna()


def build_monthly_return_window_from_start(series: pd.Series, end_date: pd.Timestamp, start_date: pd.Timestamp) -> pd.Series:
    window = build_monthly_level_window_from_start(series, end_date, start_date)
    if window.empty:
        return pd.Series(dtype=float)
    return window.pct_change().dropna()


def calc_tracking_error(series: pd.Series, benchmark_series: pd.Series, end_date: pd.Timestamp, years: int) -> float:
    asset_returns = build_monthly_return_window(series, end_date, years)
    benchmark_returns = build_monthly_return_window(benchmark_series, end_date, years)
    aligned = pd.concat([asset_returns.rename("asset"), benchmark_returns.rename("benchmark")], axis=1).dropna()
    if aligned.empty or len(aligned) < 2:
        return np.nan
    return float((aligned["asset"] - aligned["benchmark"]).std(ddof=1) * np.sqrt(12))


def calc_tracking_error_from_start(series: pd.Series, benchmark_series: pd.Series, end_date: pd.Timestamp, start_date: pd.Timestamp) -> float:
    asset_returns = build_monthly_return_window_from_start(series, end_date, start_date)
    benchmark_returns = build_monthly_return_window_from_start(benchmark_series, end_date, start_date)
    aligned = pd.concat([asset_returns.rename("asset"), benchmark_returns.rename("benchmark")], axis=1).dropna()
    if aligned.empty or len(aligned) < 2:
        return np.nan
    excess_returns = aligned["asset"] - aligned["benchmark"]
    return float(excess_returns.std(ddof=1) * np.sqrt(12))


def calc_downside_deviation(series: pd.Series, end_date: pd.Timestamp, years: int) -> float:
    monthly_returns = build_monthly_return_window(series, end_date, years)
    if monthly_returns.empty:
        return np.nan
    downside = monthly_returns.clip(upper=0)
    return float(np.sqrt((downside.pow(2)).mean()) * np.sqrt(12))


def calc_downside_deviation_from_start(series: pd.Series, end_date: pd.Timestamp, start_date: pd.Timestamp) -> float:
    monthly_returns = build_monthly_return_window_from_start(series, end_date, start_date)
    if monthly_returns.empty:
        return np.nan
    downside = monthly_returns.clip(upper=0)
    return float(np.sqrt((downside.pow(2)).mean()) * np.sqrt(12))


def calc_annualised_volatility_from_start(series: pd.Series, end_date: pd.Timestamp, start_date: pd.Timestamp) -> float:
    monthly_returns = build_monthly_return_window_from_start(series, end_date, start_date)
    if monthly_returns.empty or len(monthly_returns) < 2:
        return np.nan
    return float(monthly_returns.std(ddof=1) * np.sqrt(12))


def calc_worst_drawdown_since_inception(series: pd.Series) -> float:
    s = series.dropna().sort_index()
    if s.empty:
        return np.nan
    drawdown = (s / s.cummax()) - 1
    return float(drawdown.min())


def build_risk_summary_table(
    series_map: dict[str, pd.Series],
    asset_style_map: dict[str, float],
    selected_assets: list[str],
    end_date: pd.Timestamp,
    start_date: pd.Timestamp,
) -> pd.DataFrame:
    rows = []
    cash_series = series_map.get("Cash (GBP)", pd.Series(dtype=float))
    global_series = series_map.get("Global stocks", pd.Series(dtype=float))
    developed_series = series_map.get("Developed stocks", pd.Series(dtype=float))
    assets = selected_assets if selected_assets else list(series_map.keys())

    for asset in assets:
        if asset not in series_map:
            continue

        series = series_map[asset]
        period_return = calc_period_return(series, end_date, "Period", whole_period_start=start_date)
        period_vol = calc_annualised_volatility_from_start(series, end_date, start_date)
        downside_dev = calc_downside_deviation_from_start(series, end_date, start_date)
        worst_drawdown = calc_worst_drawdown_since_inception(series)

        growth_flag = asset_style_map.get(asset, np.nan)
        benchmark_asset = "Global stocks" if growth_flag == 1 else "Cash (GBP)"
        benchmark_series = global_series if benchmark_asset == "Global stocks" else cash_series
        benchmark_return = (
            calc_period_return(benchmark_series, end_date, "Period", whole_period_start=start_date)
            if not benchmark_series.empty
            else np.nan
        )
        tracking_error = (
            calc_tracking_error_from_start(series, benchmark_series, end_date, start_date)
            if not benchmark_series.empty and asset != benchmark_asset
            else np.nan
        )

        return_vol_ratio = (
            period_return / period_vol
            if pd.notna(period_return) and pd.notna(period_vol) and period_vol > 0
            else np.nan
        )
        sharpe_ratio = (
            (period_return - calc_period_return(cash_series, end_date, "Period", whole_period_start=start_date)) / period_vol
            if not cash_series.empty and pd.notna(period_vol) and period_vol > 0
            else np.nan
        )
        information_ratio = (
            (period_return - benchmark_return) / tracking_error
            if pd.notna(benchmark_return) and pd.notna(tracking_error) and tracking_error > 0
            else np.nan
        )
        sortino_ratio = (
            period_return / downside_dev
            if pd.notna(period_return) and pd.notna(downside_dev) and downside_dev > 0
            else np.nan
        )
        calmar_ratio = (
            period_return / abs(worst_drawdown)
            if pd.notna(period_return) and pd.notna(worst_drawdown) and worst_drawdown < 0
            else np.nan
        )

        if asset == "Cash (GBP)":
            return_vol_ratio = np.nan
            sharpe_ratio = np.nan
            information_ratio = np.nan
            sortino_ratio = np.nan
            worst_drawdown = np.nan
            calmar_ratio = np.nan
            tracking_error = np.nan

        rows.append(
            {
                "asset_class": asset,
                "Period return": period_return,
                "Period vol": period_vol,
                "Return/vol ratio": return_vol_ratio,
                "Sharpe ratio": sharpe_ratio,
                "Information ratio": information_ratio,
                "Sortino ratio": sortino_ratio,
                "Worst drawdown": worst_drawdown,
                "Calmar ratio": calmar_ratio,
                "Tracking error": tracking_error,
            }
        )

    out = order_asset_rows(pd.DataFrame(rows))
    if not out.empty and "asset_class" in out.columns:
        cash_rows = out[out["asset_class"] == "Cash (GBP)"]
        non_cash_rows = out[out["asset_class"] != "Cash (GBP)"]
        out = pd.concat([non_cash_rows, cash_rows], ignore_index=True)
    return out


def build_correlation_matrix_table(
    series_map: dict[str, pd.Series],
    selected_assets: list[str],
    end_date: pd.Timestamp,
    start_date: pd.Timestamp,
) -> pd.DataFrame:
    assets = selected_assets if selected_assets else list(series_map.keys())
    return_frames = []

    for asset in assets:
        if asset not in series_map:
            continue
        monthly_returns = build_monthly_return_window_from_start(series_map[asset], end_date, start_date)
        if monthly_returns.empty:
            continue
        return_frames.append(monthly_returns.rename(asset))

    if not return_frames:
        return pd.DataFrame()

    returns_df = pd.concat(return_frames, axis=1).dropna(how="all")
    if returns_df.empty:
        return pd.DataFrame()

    corr = returns_df.corr()
    ordered_assets = [c for c in assets if c in corr.index and c != "Cash (GBP)"]
    if "Cash (GBP)" in corr.index:
        ordered_assets.append("Cash (GBP)")
    corr = corr.reindex(index=ordered_assets, columns=ordered_assets)

    for row_idx, row_name in enumerate(corr.index):
        for col_idx, col_name in enumerate(corr.columns):
            if col_idx > row_idx:
                corr.loc[row_name, col_name] = np.nan

    corr = corr.reset_index().rename(columns={"index": "asset_class"})
    return corr


def build_risk_metrics_table(
    series_map: dict[str, pd.Series],
    selected_assets: list[str],
    end_date: pd.Timestamp,
    period_keys: list[str],
    whole_period_start: pd.Timestamp,
) -> pd.DataFrame:
    rows = []
    assets = selected_assets if selected_assets else list(series_map.keys())

    for asset in assets:
        if asset not in series_map:
            continue

        row = {"asset_class": asset}
        series = series_map[asset]
        for period_key in period_keys:
            if period_key == "Max":
                row[f"{period_key} Return"] = calc_period_return(series, end_date, "Period", whole_period_start=whole_period_start)
                row[f"{period_key} Vol"] = calc_annualised_volatility_from_start(series, end_date, whole_period_start)
            else:
                years = int(period_key.replace("Y", ""))
                row[f"{period_key} Return"] = calc_period_return(series, end_date, period_key)
                row[f"{period_key} Vol"] = calc_annualised_volatility(series, end_date, years)
        rows.append(row)

    return order_asset_rows(pd.DataFrame(rows))


def build_risk_scatter_df(
    series_map: dict[str, pd.Series],
    selected_assets: list[str],
    end_date: pd.Timestamp,
    period_key: str,
    whole_period_start: pd.Timestamp,
) -> pd.DataFrame:
    rows = []
    assets = selected_assets if selected_assets else list(series_map.keys())

    for asset in assets:
        if asset not in series_map:
            continue

        series = series_map[asset]
        if period_key == "Max":
            annual_return = calc_period_return(series, end_date, "Period", whole_period_start=whole_period_start)
            annual_vol = calc_annualised_volatility_from_start(series, end_date, whole_period_start)
        else:
            years = int(period_key.replace("Y", ""))
            annual_return = calc_period_return(series, end_date, period_key)
            annual_vol = calc_annualised_volatility(series, end_date, years)
        if pd.isna(annual_return) or pd.isna(annual_vol):
            continue

        rows.append(
            {
                "asset_class": asset,
                "display_asset_class": display_name(asset),
                "annual_return": annual_return,
                "annual_volatility": annual_vol,
            }
        )

    return pd.DataFrame(rows)


def build_calendar_year_returns(stitched_series_map: dict[str, pd.Series], end_date: pd.Timestamp, years_back: int = 10) -> pd.DataFrame:
    last_complete_year = end_date.year - 1
    years = list(reversed(list(range(last_complete_year - years_back + 1, last_complete_year + 1))))
    rows = []

    for asset_class, series in stitched_series_map.items():
        row = {"asset_class": asset_class}
        s = series.dropna().sort_index()

        for year in years:
            _, end_level = nearest_level_on_or_before(s, pd.Timestamp(year, 12, 31))
            _, start_level = nearest_level_on_or_before(s, pd.Timestamp(year - 1, 12, 31))
            if end_level is None or start_level is None or start_level <= 0:
                row[str(year)] = np.nan
            else:
                row[str(year)] = (end_level / start_level) - 1

        rows.append(row)

    return order_asset_rows(pd.DataFrame(rows))


def format_pct_strings(df: pd.DataFrame, exclude_cols: set[str] | None = None) -> pd.DataFrame:
    exclude_cols = exclude_cols or {"asset_class"}
    out = df.copy()
    for col in out.columns:
        if col not in exclude_cols:
            out[col] = out[col].map(format_pct)
    return out


def build_html_table(
    df: pd.DataFrame,
    percent_cols: list[str] | None = None,
    conditional_cols: list[str] | None = None,
    header_wrap_cols: list[str] | None = None,
    invert_conditional_cols: list[str] | None = None,
    rank_conditional_cols: list[str] | None = None,
    decimal_cols: list[str] | None = None,
    correlation_conditional_cols: list[str] | None = None,
) -> str:
    if df.empty:
        return '<div class="table-shell"><div class="table-empty">No data available.</div></div>'

    display_df = df.copy()
    if "asset_class" in display_df.columns:
        display_df["asset_class"] = display_df["asset_class"].map(display_name)
        display_df = display_df.rename(columns={"asset_class": "Asset class"})

    cols = list(display_df.columns)
    percent_cols = set(percent_cols or [])
    conditional_cols = set(conditional_cols or [])
    header_wrap_cols = set(header_wrap_cols or [])
    invert_conditional_cols = set(invert_conditional_cols or [])
    rank_conditional_cols = set(rank_conditional_cols or [])
    decimal_cols = set(decimal_cols or [])
    correlation_conditional_cols = set(correlation_conditional_cols or [])
    if cols:
        asset_col_width = 28 if "Asset class" in cols else 0
        other_count = len(cols) - (1 if "Asset class" in cols else 0)
        other_width = ((100 - asset_col_width) / other_count) if other_count > 0 else 100
    else:
        asset_col_width = 0
        other_width = 100

    colgroup = "".join(
        [
            f'<col style="width:{asset_col_width if col == "Asset class" else other_width:.4f}%;">'
            for col in cols
        ]
    )

    thead = "".join(
        [
            f"<th>{col.replace(' ', '<br>') if col in header_wrap_cols else col}</th>"
            for col in cols
        ]
    )

    heat_bounds = {}
    for col in cols:
        source_col = "asset_class" if col == "Asset class" else col
        if source_col in conditional_cols and source_col in df.columns and pd.api.types.is_numeric_dtype(df[source_col]):
            numeric = pd.to_numeric(df[source_col], errors="coerce").dropna()
            heat_bounds[source_col] = (
                float(numeric.min()) if not numeric.empty else 0.0,
                float(numeric.max()) if not numeric.empty else 0.0,
            )

    body_rows = []
    for row_idx, (_, row) in enumerate(display_df.iterrows()):
        cells = []
        for col in cols:
            value = row[col]
            source_col = "asset_class" if col == "Asset class" else col
            cell_text = "-"
            cell_style = ""

            if pd.notna(value):
                if source_col in percent_cols:
                    cell_text = format_pct(float(value))
                elif source_col in decimal_cols:
                    cell_text = f"{float(value):.2f}"
                else:
                    cell_text = str(value)

            if source_col in heat_bounds and pd.notna(value):
                vmin, vmax = heat_bounds[source_col]
                if source_col in correlation_conditional_cols:
                    bg_colour = correlation_heat_colour(float(value))
                elif source_col in rank_conditional_cols:
                    bg_colour = rank_heat_colour(
                        float(value),
                        vmin,
                        vmax,
                        low_is_good=source_col in invert_conditional_cols,
                    )
                else:
                    heat_value = -float(value) if source_col in invert_conditional_cols else float(value)
                    heat_vmin = -vmax if source_col in invert_conditional_cols else vmin
                    heat_vmax = -vmin if source_col in invert_conditional_cols else vmax
                    bg_colour = heat_colour(heat_value, heat_vmin, heat_vmax)
                cell_style = (
                    f' style="background:{bg_colour};'
                    ' color:#ffffff; font-weight:600;"'
                )

            cells.append(f"<td{cell_style}>{cell_text}</td>")
        body_rows.append(f"<tr>{''.join(cells)}</tr>")

    tbody = "".join(body_rows)

    return f"""
    <div class="table-shell">
        <table class="custom-data-table">
            <colgroup>{colgroup}</colgroup>
            <thead>
                <tr>{thead}</tr>
            </thead>
            <tbody>
                {tbody}
            </tbody>
        </table>
    </div>
    """


def convert_pct_table_for_display(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    for col in out.columns:
        if col != "asset_class" and pd.api.types.is_numeric_dtype(out[col]):
            out[col] = out[col].map(lambda x: np.nan if pd.isna(x) else round(x * 100, 2))
    return prepare_dataframe_for_display(out)


def get_chart_axis_format(period_key: str) -> str:
    if period_key in {"YTD", "1Y"}:
        return "%d/%m/%y"
    if period_key in {"3Y", "5Y"}:
        return "%b-%y"
    return "%Y"


def get_chart_right_padding_days(period_key: str) -> int:
    if period_key == "YTD":
        return 10
    if period_key == "1Y":
        return 18
    if period_key in {"3Y", "5Y"}:
        return 35
    return 120


def build_chart(chart_df: pd.DataFrame, selected_assets: list[str], period_key: str) -> alt.Chart:
    if chart_df.empty:
        return alt.Chart(pd.DataFrame({"x": [], "y": []})).mark_line()

    chart_df = chart_df.copy()
    chart_df["asset_class"] = chart_df["asset_class"].map(display_name)

    colour_domain = [display_name(a) for a in selected_assets if display_name(a) in chart_df["asset_class"].unique()]
    colour_range = [ASSET_COLOURS.get(a, "#1f77b4") for a in selected_assets if display_name(a) in chart_df["asset_class"].unique()]

    ymin = float(chart_df["Growth"].min())
    ymax = float(chart_df["Growth"].max())
    spread = max(ymax - ymin, 0.02)
    pad = spread * 0.10
    domain_min = ymin - pad
    domain_max = ymax + pad
    x_axis_format = get_chart_axis_format(period_key)

    return (
        alt.Chart(chart_df)
        .mark_line(strokeWidth=2.5)
        .encode(
            x=alt.X(
                "Date:T",
                title=None,
                axis=alt.Axis(
                    format=x_axis_format,
                    labelColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=False,
                ),
            ),
            y=alt.Y(
                "Growth:Q",
                title=None,
                scale=alt.Scale(domain=[domain_min, domain_max], zero=False),
                axis=alt.Axis(
                    labelColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=True,
                    gridColor=MID_GREY,
                    gridDash=[2, 4],
                    labelExpr="'£' + format(datum.value, ',.2f')",
                ),
            ),
            color=alt.Color(
                "asset_class:N",
                title="Asset class",
                scale=alt.Scale(domain=colour_domain, range=colour_range),
                legend=alt.Legend(
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    orient="bottom",
                    direction="horizontal",
                    columns=max(1, len(colour_domain)),
                    title=None,
                    symbolLimit=len(colour_domain),
                    labelLimit=180,
                ),
            ),
            tooltip=[
                alt.Tooltip("Date:T", title="Date"),
                alt.Tooltip("asset_class:N", title="Asset class"),
                alt.Tooltip("Growth:Q", title="Growth of wealth", format=".3f"),
            ],
        )
        .properties(
            height=420,
            width="container",
            padding={"left": 14, "top": 8, "right": max(28, int(get_chart_right_padding_days(period_key) * 0.8)), "bottom": 10},
        )
        .configure_view(stroke=None, fill=CHART_BG_GREY)
        .configure_axis(labelFont="Calibri", titleFont="Calibri")
        .configure_legend(
            labelFont="Calibri",
            titleFont="Calibri",
            fillColor=CHART_BG_GREY,
            strokeColor=CHART_BG_GREY,
        )
        .configure(background=CHART_BG_GREY)
    )


def build_risk_scatter_chart(risk_df: pd.DataFrame, selected_assets: list[str]) -> alt.Chart:
    if risk_df.empty:
        return alt.Chart(pd.DataFrame({"x": [], "y": []})).mark_circle()

    colour_domain = [
        display_name(a)
        for a in selected_assets
        if display_name(a) in risk_df["display_asset_class"].unique()
    ]
    colour_range = [
        ASSET_COLOURS.get(a, "#1f77b4")
        for a in selected_assets
        if display_name(a) in risk_df["display_asset_class"].unique()
    ]

    return (
        alt.Chart(risk_df)
        .mark_circle(size=180, opacity=0.95)
        .encode(
            x=alt.X(
                "annual_volatility:Q",
                title="Annualised volatility",
                axis=alt.Axis(
                    format=".0%",
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=True,
                    gridColor=MID_GREY,
                    gridDash=[2, 4],
                ),
            ),
            y=alt.Y(
                "annual_return:Q",
                title="Annualised return",
                axis=alt.Axis(
                    format=".0%",
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=True,
                    gridColor=MID_GREY,
                    gridDash=[2, 4],
                ),
            ),
            color=alt.Color(
                "display_asset_class:N",
                title=None,
                scale=alt.Scale(domain=colour_domain, range=colour_range),
                legend=alt.Legend(
                    labelColor=TEXT_GREY,
                    orient="bottom",
                    direction="horizontal",
                    columns=max(1, len(colour_domain)),
                    symbolLimit=len(colour_domain),
                    labelLimit=180,
                ),
            ),
            tooltip=[
                alt.Tooltip("display_asset_class:N", title="Asset class"),
                alt.Tooltip("annual_return:Q", title="Annualised return", format=".2%"),
                alt.Tooltip("annual_volatility:Q", title="Annualised volatility", format=".2%"),
            ],
        )
        .properties(height=420, width="container", padding={"left": 14, "top": 8, "right": 28, "bottom": 10})
        .configure_view(stroke=None, fill=CHART_BG_GREY)
        .configure_axis(labelFont="Calibri", titleFont="Calibri")
        .configure_legend(
            labelFont="Calibri",
            titleFont="Calibri",
            fillColor=CHART_BG_GREY,
            strokeColor=CHART_BG_GREY,
        )
        .configure(background=CHART_BG_GREY)
    )


def select_curve_date_on_or_before(available_dates: list[pd.Timestamp], target_date: pd.Timestamp) -> pd.Timestamp | None:
    candidates = [pd.Timestamp(date).normalize() for date in available_dates if pd.Timestamp(date).normalize() <= pd.Timestamp(target_date).normalize()]
    return max(candidates) if candidates else None


def build_uk_yield_curve_overlay_options(history_df: pd.DataFrame) -> list[dict[str, object]]:
    if history_df.empty:
        return []

    history = history_df.copy()
    history["curve_date"] = pd.to_datetime(history["curve_date"], errors="coerce").dt.normalize()
    nominal_dates = set(history.loc[history["curve_type"] == "Nominal", "curve_date"].dropna().drop_duplicates().tolist())
    real_dates = set(history.loc[history["curve_type"] == "Real", "curve_date"].dropna().drop_duplicates().tolist())
    common_dates = sorted(pd.Timestamp(date).normalize() for date in nominal_dates.intersection(real_dates))
    if not common_dates:
        return []

    latest_common = max(common_dates)
    target_defs = [
        ("Last month end", latest_common),
        ("LME-1", (latest_common.to_period("M") - 1).end_time.normalize()),
        ("LME-2", (latest_common.to_period("M") - 2).end_time.normalize()),
        ("LME-3", (latest_common.to_period("M") - 3).end_time.normalize()),
    ]

    options: list[dict[str, object]] = []
    for label, target in target_defs:
        resolved = select_curve_date_on_or_before(common_dates, pd.Timestamp(target))
        if resolved is None:
            continue
        options.append({"label": label, "curve_date": resolved, "key": re.sub(r"[^a-z0-9]+", "_", label.lower()).strip("_")})
    return options


def build_breakeven_curve_points(curve_df: pd.DataFrame) -> pd.DataFrame:
    expected = {"maturity_years", "yield_percent", "curve_type", "curve_date", "snapshot_label", "is_current", "snapshot_sort"}
    if curve_df.empty or not expected.issubset(curve_df.columns):
        return pd.DataFrame(columns=list(expected))

    out_frames = []
    for (snapshot_label, is_current, snapshot_sort), group in curve_df.groupby(["snapshot_label", "is_current", "snapshot_sort"], dropna=False):
        nominal = group[group["curve_type"] == "Nominal"][["maturity_years", "yield_percent", "curve_date"]].copy()
        real = group[group["curve_type"] == "Real"][["maturity_years", "yield_percent", "curve_date"]].copy()
        if nominal.empty or real.empty:
            continue

        nominal = nominal.rename(columns={"yield_percent": "nominal_yield", "curve_date": "nominal_date"})
        real = real.rename(columns={"yield_percent": "real_yield", "curve_date": "real_date"})
        nominal_sorted = nominal.sort_values("maturity_years").dropna(subset=["maturity_years", "nominal_yield"])
        real_sorted = real.sort_values("maturity_years").dropna(subset=["maturity_years", "real_yield"])
        if nominal_sorted.empty or real_sorted.empty:
            continue

        interpolated_nominal = np.interp(
            real_sorted["maturity_years"].to_numpy(),
            nominal_sorted["maturity_years"].to_numpy(),
            nominal_sorted["nominal_yield"].to_numpy(),
        )
        bei = real_sorted.copy()
        bei["nominal_yield"] = interpolated_nominal
        bei["nominal_date"] = nominal_sorted["nominal_date"].max()
        bei["curve_type"] = "Breakeven inflation"
        bei["yield_percent"] = (((1 + (bei["nominal_yield"] / 100.0)) / (1 + (bei["real_yield"] / 100.0))) - 1) * 100.0
        bei["curve_date"] = bei[["nominal_date", "real_date"]].max(axis=1)
        bei["snapshot_label"] = snapshot_label
        bei["is_current"] = bool(is_current)
        bei["snapshot_sort"] = int(snapshot_sort)
        out_frames.append(
            bei[["maturity_years", "yield_percent", "curve_type", "curve_date", "snapshot_label", "is_current", "snapshot_sort"]]
        )

    return pd.concat(out_frames, ignore_index=True) if out_frames else pd.DataFrame(columns=list(expected))


def shade_hex_colour(hex_colour: str, blend: float) -> str:
    text = str(hex_colour).strip()
    if not text.startswith("#") or len(text) != 7:
        return text
    blend = min(max(float(blend), 0.0), 1.0)
    r = int(text[1:3], 16)
    g = int(text[3:5], 16)
    b = int(text[5:7], 16)
    r = int(round(r + (255 - r) * blend))
    g = int(round(g + (255 - g) * blend))
    b = int(round(b + (255 - b) * blend))
    return f"#{r:02x}{g:02x}{b:02x}"


def build_yield_curve_chart(yield_curve_df: pd.DataFrame) -> alt.Chart:
    if yield_curve_df.empty:
        return alt.Chart(pd.DataFrame({"x": [], "y": []})).mark_line()

    chart_df = yield_curve_df.copy()
    chart_df["curve_type"] = pd.Categorical(chart_df["curve_type"], categories=["Nominal", "Real", "Breakeven inflation"], ordered=True)
    chart_df["yield_decimal"] = pd.to_numeric(chart_df["yield_percent"], errors="coerce") / 100.0
    chart_df["snapshot_label"] = chart_df.get("snapshot_label", "Latest").astype(str)
    chart_df["is_current"] = chart_df.get("is_current", True).astype(bool)
    chart_df["snapshot_sort"] = pd.to_numeric(chart_df.get("snapshot_sort", 0), errors="coerce").fillna(0).astype(int)
    chart_df["variant_idx"] = chart_df["snapshot_sort"].clip(lower=0, upper=4)

    base_colour_map = {"Nominal": "#c95b2b", "Real": "#1f77b4", "Breakeven inflation": "#2e8b57"}
    historical_blends = {1: 0.12, 2: 0.12, 3: 0.38, 4: 0.38}
    opacity_map = {1: 0.66, 2: 0.66, 3: 0.33, 4: 0.33}
    stroke_dash_map = {1: [1, 0], 2: [9, 3], 3: [1, 0], 4: [9, 3]}
    historical_colour_domain = []
    historical_colour_range = []
    for variant_idx in [1, 2, 3, 4]:
        for curve_type in ["Nominal", "Real", "Breakeven inflation"]:
            historical_colour_domain.append(f"{curve_type}|{variant_idx}")
            historical_colour_range.append(
                shade_hex_colour(base_colour_map.get(curve_type, "#666666"), historical_blends.get(variant_idx, 0.35))
            )
    chart_df["historical_series"] = chart_df.apply(
        lambda row: f"{row['curve_type']}|{int(row['variant_idx'])}" if not bool(row["is_current"]) else "",
        axis=1,
    )

    base = alt.Chart(chart_df).encode(
            x=alt.X(
                "maturity_years:Q",
                title="Maturity (years)",
                axis=alt.Axis(
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=False,
                ),
            ),
            y=alt.Y(
                "yield_decimal:Q",
                title="Yield",
                axis=alt.Axis(
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=True,
                    gridColor=MID_GREY,
                    gridDash=[2, 4],
                    format=".1%",
                ),
            ),
            order=alt.Order("snapshot_sort:Q", sort="ascending"),
            tooltip=[
                alt.Tooltip("curve_type:N", title="Curve"),
                alt.Tooltip("snapshot_label:N", title="Snapshot"),
                alt.Tooltip("curve_date:T", title="Date"),
                alt.Tooltip("maturity_years:Q", title="Maturity (years)", format=".1f"),
                alt.Tooltip("yield_decimal:Q", title="Yield", format=".2%"),
            ],
        )

    latest = (
        base.transform_filter(alt.datum.is_current == True)
        .mark_line(point=True, strokeWidth=2.8, opacity=1.0)
        .encode(
            detail=["curve_type:N", "snapshot_label:N"],
            color=alt.Color(
                "curve_type:N",
                title=None,
                scale=alt.Scale(domain=["Nominal", "Real", "Breakeven inflation"], range=["#c95b2b", "#1f77b4", "#2e8b57"]),
                legend=alt.Legend(
                    labelColor=TEXT_GREY,
                    orient="bottom",
                    direction="horizontal",
                    columns=3,
                    symbolLimit=3,
                ),
            ),
        )
    )

    layers = []
    for variant_idx in [1, 2, 3, 4]:
        layers.append(
            base.transform_filter((alt.datum.is_current == False) & (alt.datum.snapshot_sort == variant_idx))
            .mark_line(
                point=False,
                strokeWidth=2.2,
                opacity=opacity_map[variant_idx],
                strokeDash=stroke_dash_map[variant_idx],
            )
            .encode(
                detail=["curve_type:N", "snapshot_label:N"],
                color=alt.Color(
                    "historical_series:N",
                    scale=alt.Scale(domain=historical_colour_domain, range=historical_colour_range),
                    legend=None,
                ),
            )
        )

    return (
        alt.layer(*layers, latest)
        .properties(height=420, width="container", padding={"left": 14, "top": 8, "right": 28, "bottom": 10})
        .configure_view(stroke=None, fill=CHART_BG_GREY)
        .configure_axis(labelFont="Calibri", titleFont="Calibri")
        .configure_legend(
            labelFont="Calibri",
            titleFont="Calibri",
            fillColor=CHART_BG_GREY,
            strokeColor=CHART_BG_GREY,
        )
        .configure(background=CHART_BG_GREY)
    )


@st.cache_data(show_spinner=False)
def build_uk_historical_yield_df(history_df: pd.DataFrame) -> pd.DataFrame:
    if history_df.empty:
        return pd.DataFrame(columns=["curve_date", "curve_type", "target_maturity", "selected_maturity", "yield_percent", "series_label"])

    working = history_df.copy()
    working["curve_date"] = pd.to_datetime(working["curve_date"], errors="coerce").dt.normalize()
    working["maturity_years"] = pd.to_numeric(working["maturity_years"], errors="coerce")
    working["yield_percent"] = pd.to_numeric(working["yield_percent"], errors="coerce")
    working = working.dropna(subset=["curve_date", "curve_type", "maturity_years", "yield_percent"])
    if working.empty:
        return pd.DataFrame(columns=["curve_date", "curve_type", "target_maturity", "selected_maturity", "yield_percent", "series_label"])

    def maturity_label(value: float) -> str:
        number = float(value)
        if abs(number - round(number)) < 1e-9:
            return f"{int(round(number))}Y"
        return f"{number:.1f}Y"

    target_maturities = [1.0, 5.0, 10.0, 30.0]
    rows = []
    for (curve_date, curve_type), group in working.groupby(["curve_date", "curve_type"], dropna=False):
        group = group.sort_values("maturity_years").drop_duplicates(subset=["maturity_years"], keep="last")
        if group.empty:
            continue
        min_maturity = float(group["maturity_years"].min())
        max_maturity = float(group["maturity_years"].max())
        for target_maturity in target_maturities:
            if curve_type == "Real" and target_maturity == 1.0:
                selected = group.iloc[0]
                label_maturity = float(selected["maturity_years"])
            else:
                if target_maturity < min_maturity or target_maturity > max_maturity:
                    continue
                distance = (group["maturity_years"] - target_maturity).abs()
                selected = group.loc[distance.idxmin()]
                label_maturity = float(target_maturity)
            rows.append(
                {
                    "curve_date": pd.Timestamp(curve_date),
                    "curve_type": str(curve_type),
                    "target_maturity": float(target_maturity),
                    "selected_maturity": float(selected["maturity_years"]),
                    "yield_percent": float(selected["yield_percent"]),
                    "series_label": f"{curve_type} {maturity_label(label_maturity)}",
                }
            )

    out = pd.DataFrame(rows)
    if out.empty:
        return out

    bei_rows = []
    for curve_date, group in out.groupby("curve_date", dropna=False):
        nominal = group[group["curve_type"] == "Nominal"].copy()
        real = group[group["curve_type"] == "Real"].copy()
        if nominal.empty or real.empty:
            continue
        nominal = nominal.sort_values("selected_maturity")
        if nominal.empty:
            continue
        for row in real.itertuples():
            maturity = float(row.selected_maturity)
            interpolated_nominal = np.interp(
                maturity,
                nominal["selected_maturity"].to_numpy(),
                nominal["yield_percent"].to_numpy(),
            )
            bei_rows.append(
                {
                    "curve_date": pd.Timestamp(curve_date),
                    "curve_type": "Breakeven inflation",
                    "target_maturity": float(row.target_maturity),
                    "selected_maturity": maturity,
                    "yield_percent": (((1 + (interpolated_nominal / 100.0)) / (1 + (float(row.yield_percent) / 100.0))) - 1) * 100.0,
                    "series_label": f"BEI {maturity_label(maturity)}",
                }
            )

    if bei_rows:
        out = pd.concat([out, pd.DataFrame(bei_rows)], ignore_index=True)

    return out.sort_values(["curve_date", "curve_type", "selected_maturity"]).reset_index(drop=True)


def get_uk_historical_yield_start_date(period_key: str, end_date: pd.Timestamp, min_start_date: pd.Timestamp) -> pd.Timestamp:
    if period_key == "MAX":
        return pd.Timestamp(min_start_date).normalize()
    return get_chart_period_start_date(period_key, end_date, min_start_date)


def build_uk_historical_yield_chart(history_df: pd.DataFrame, selected_series: list[str], period_key: str) -> alt.Chart:
    if history_df.empty:
        return alt.Chart(pd.DataFrame({"curve_date": [], "yield_decimal": [], "series_label": []})).mark_line()

    chart_df = history_df.copy()
    chart_df["curve_date"] = pd.to_datetime(chart_df["curve_date"], errors="coerce")
    chart_df["yield_decimal"] = pd.to_numeric(chart_df["yield_percent"], errors="coerce") / 100.0
    chart_df = chart_df.dropna(subset=["curve_date", "yield_decimal", "series_label"])
    if selected_series:
        chart_df = chart_df[chart_df["series_label"].isin(selected_series)].copy()
    if chart_df.empty:
        return alt.Chart(pd.DataFrame({"curve_date": [], "yield_decimal": [], "series_label": []})).mark_line()

    series_order = [label for label in selected_series if label in chart_df["series_label"].unique()] if selected_series else sorted(chart_df["series_label"].unique())
    colour_map = {
        "Nominal 1Y": "#dd865d",
        "Nominal 5Y": "#c95b2b",
        "Nominal 10Y": "#a9441c",
        "Nominal 30Y": "#7f2e0f",
        "Real 2.5Y": "#7cb7e5",
        "Real 5Y": "#529ad0",
        "Real 10Y": "#2f7fbc",
        "Real 30Y": "#15527b",
        "BEI 2.5Y": "#82c792",
        "BEI 5Y": "#58ad6e",
        "BEI 10Y": "#2e8b57",
        "BEI 30Y": "#1f6a41",
    }
    colour_range = [colour_map.get(label, "#666666") for label in series_order]
    x_axis_format = get_chart_axis_format(period_key if period_key in CHART_PERIODS else "20Y")

    return (
        alt.Chart(chart_df)
        .mark_line(strokeWidth=2.4)
        .encode(
            x=alt.X(
                "curve_date:T",
                title=None,
                axis=alt.Axis(
                    format=x_axis_format,
                    labelColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=False,
                ),
            ),
            y=alt.Y(
                "yield_decimal:Q",
                title="Yield",
                axis=alt.Axis(
                    format=".1%",
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=True,
                    gridColor=MID_GREY,
                    gridDash=[2, 4],
                ),
            ),
            color=alt.Color(
                "series_label:N",
                title=None,
                scale=alt.Scale(domain=series_order, range=colour_range),
                legend=alt.Legend(
                    labelColor=TEXT_GREY,
                    orient="bottom",
                    direction="horizontal",
                    columns=max(1, len(series_order)),
                    labelLimit=180,
                ),
            ),
            tooltip=[
                alt.Tooltip("curve_date:T", title="Date"),
                alt.Tooltip("series_label:N", title="Series"),
                alt.Tooltip("yield_decimal:Q", title="Yield", format=".2%"),
                alt.Tooltip("selected_maturity:Q", title="Source maturity", format=".1f"),
            ],
        )
        .properties(height=360, width="container", padding={"left": 14, "top": 8, "right": 28, "bottom": 10})
        .configure_view(stroke=None, fill=CHART_BG_GREY)
        .configure_axis(labelFont="Calibri", titleFont="Calibri")
        .configure_legend(
            labelFont="Calibri",
            titleFont="Calibri",
            fillColor=CHART_BG_GREY,
            strokeColor=CHART_BG_GREY,
        )
        .configure(background=CHART_BG_GREY)
    )


@st.cache_data(show_spinner=False)
def build_uk_term_spread_df(history_df: pd.DataFrame) -> pd.DataFrame:
    if history_df.empty:
        return pd.DataFrame(columns=["curve_date", "curve_type", "series_label", "spread_percent", "short_maturity", "long_maturity"])

    working = history_df.copy()
    working["curve_date"] = pd.to_datetime(working["curve_date"], errors="coerce").dt.normalize()
    working["maturity_years"] = pd.to_numeric(working["maturity_years"], errors="coerce")
    working["yield_percent"] = pd.to_numeric(working["yield_percent"], errors="coerce")
    working = working.dropna(subset=["curve_date", "curve_type", "maturity_years", "yield_percent"])
    if working.empty:
        return pd.DataFrame(columns=["curve_date", "curve_type", "series_label", "spread_percent", "short_maturity", "long_maturity"])

    spread_defs = [
        ("Nominal", 30.0, 10.0, "Nominal 30-10Y"),
        ("Nominal", 10.0, 2.0, "Nominal 10-2Y"),
        ("Nominal", 5.0, 1.0, "Nominal 5-1Y"),
        ("Nominal", 10.0, 0.5, "Nominal 10-0.5Y"),
        ("Nominal", 2.0, 0.5, "Nominal 2-0.5Y"),
        ("Real", 30.0, 10.0, "Real 30-10Y"),
        ("Real", 10.0, 2.5, "Real 10-2.5Y"),
        ("Real", 5.0, 2.5, "Real 5-2.5Y"),
    ]

    rows = []
    for (curve_date, curve_type), group in working.groupby(["curve_date", "curve_type"], dropna=False):
        group = group.sort_values("maturity_years").drop_duplicates(subset=["maturity_years"], keep="last")
        if group.empty:
            continue
        for target_curve_type, long_target, short_target, series_label in spread_defs:
            if curve_type != target_curve_type:
                continue
            long_idx = (group["maturity_years"] - long_target).abs().idxmin()
            short_idx = (group["maturity_years"] - short_target).abs().idxmin()
            long_row = group.loc[long_idx]
            short_row = group.loc[short_idx]
            if abs(float(long_row["maturity_years"]) - long_target) > 0.26:
                continue
            if abs(float(short_row["maturity_years"]) - short_target) > 0.26:
                continue
            rows.append(
                {
                    "curve_date": pd.Timestamp(curve_date),
                    "curve_type": str(curve_type),
                    "series_label": series_label,
                    "spread_percent": float(long_row["yield_percent"]) - float(short_row["yield_percent"]),
                    "short_maturity": float(short_row["maturity_years"]),
                    "long_maturity": float(long_row["maturity_years"]),
                }
            )

    out = pd.DataFrame(rows)
    if out.empty:
        return out
    return out.sort_values(["curve_date", "curve_type", "series_label"]).reset_index(drop=True)


def build_uk_term_spread_chart(spread_df: pd.DataFrame, selected_series: list[str], period_key: str) -> alt.Chart:
    if spread_df.empty:
        return alt.Chart(pd.DataFrame({"curve_date": [], "spread_decimal": [], "series_label": []})).mark_line()

    chart_df = spread_df.copy()
    chart_df["curve_date"] = pd.to_datetime(chart_df["curve_date"], errors="coerce")
    chart_df["spread_decimal"] = pd.to_numeric(chart_df["spread_percent"], errors="coerce") / 100.0
    chart_df = chart_df.dropna(subset=["curve_date", "spread_decimal", "series_label"])
    if selected_series:
        chart_df = chart_df[chart_df["series_label"].isin(selected_series)].copy()
    if chart_df.empty:
        return alt.Chart(pd.DataFrame({"curve_date": [], "spread_decimal": [], "series_label": []})).mark_line()

    series_order = [label for label in selected_series if label in chart_df["series_label"].unique()] if selected_series else sorted(chart_df["series_label"].unique())
    colour_map = {
        "Nominal 30-10Y": "#e15759",
        "Nominal 10-2Y": "#f28e2b",
        "Nominal 5-1Y": "#edc948",
        "Nominal 10-0.5Y": "#b07aa1",
        "Nominal 2-0.5Y": "#9c755f",
        "Real 30-10Y": "#76b7b2",
        "Real 10-2.5Y": "#4e79a7",
        "Real 5-2.5Y": "#1f3f75",
    }
    colour_range = [colour_map.get(label, "#666666") for label in series_order]
    x_axis_format = get_chart_axis_format(period_key if period_key in CHART_PERIODS else "20Y")

    return (
        alt.Chart(chart_df)
        .mark_line(strokeWidth=2.4)
        .encode(
            x=alt.X(
                "curve_date:T",
                title=None,
                axis=alt.Axis(
                    format=x_axis_format,
                    labelColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=False,
                ),
            ),
            y=alt.Y(
                "spread_decimal:Q",
                title="Spread",
                axis=alt.Axis(
                    format=".1%",
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=True,
                    gridColor=MID_GREY,
                    gridDash=[2, 4],
                ),
            ),
            color=alt.Color(
                "series_label:N",
                title=None,
                scale=alt.Scale(domain=series_order, range=colour_range),
                legend=alt.Legend(
                    labelColor=TEXT_GREY,
                    orient="bottom",
                    direction="horizontal",
                    columns=max(1, len(series_order)),
                    labelLimit=180,
                ),
            ),
            tooltip=[
                alt.Tooltip("curve_date:T", title="Date"),
                alt.Tooltip("series_label:N", title="Spread"),
                alt.Tooltip("spread_decimal:Q", title="Spread", format=".2%"),
                alt.Tooltip("long_maturity:Q", title="Long maturity", format=".1f"),
                alt.Tooltip("short_maturity:Q", title="Short maturity", format=".1f"),
            ],
        )
        .properties(height=340, width="container", padding={"left": 14, "top": 8, "right": 28, "bottom": 10})
        .configure_view(stroke=None, fill=CHART_BG_GREY)
        .configure_axis(labelFont="Calibri", titleFont="Calibri")
        .configure_legend(
            labelFont="Calibri",
            titleFont="Calibri",
            fillColor=CHART_BG_GREY,
            strokeColor=CHART_BG_GREY,
        )
        .configure(background=CHART_BG_GREY)
    )


def build_global_yield_curve_chart(global_yield_curve_df: pd.DataFrame) -> alt.Chart:
    if global_yield_curve_df.empty:
        return alt.Chart(pd.DataFrame({"maturity_years": [], "yield_decimal": [], "country": []}))

    chart_df = global_yield_curve_df.copy()
    chart_df["yield_decimal"] = pd.to_numeric(chart_df["yield_percent"], errors="coerce") / 100.0
    chart_df["maturity_years"] = pd.to_numeric(chart_df["maturity_years"], errors="coerce")
    chart_df = chart_df.dropna(subset=["maturity_years", "yield_decimal", "country"])
    if chart_df.empty:
        return alt.Chart(pd.DataFrame({"maturity_years": [], "yield_decimal": [], "country": []}))
    legend_columns = max(1, int(chart_df["country"].nunique()))

    return (
        alt.Chart(chart_df)
        .mark_line(point=False, strokeWidth=2.3)
        .encode(
            x=alt.X(
                "maturity_years:Q",
                title="Maturity (years)",
                axis=alt.Axis(
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=False,
                ),
            ),
            y=alt.Y(
                "yield_decimal:Q",
                title="Yield",
                axis=alt.Axis(
                    format=".1%",
                    labelColor=TEXT_GREY,
                    titleColor=TEXT_GREY,
                    labelFontSize=15,
                    tickColor=MID_GREY,
                    domainColor=MID_GREY,
                    grid=True,
                    gridColor=MID_GREY,
                    gridDash=[2, 4],
                ),
            ),
            color=alt.Color(
                "country:N",
                title=None,
                legend=alt.Legend(
                    labelColor=TEXT_GREY,
                    orient="bottom",
                    direction="horizontal",
                    columns=legend_columns,
                    labelLimit=200,
                    symbolType="stroke",
                    symbolLimit=legend_columns,
                ),
            ),
            tooltip=[
                alt.Tooltip("country:N", title="Country"),
                alt.Tooltip("maturity_years:Q", title="Maturity (years)", format=".2f"),
                alt.Tooltip("yield_decimal:Q", title="Yield", format=".2%"),
            ],
        )
        .properties(height=420, width="container", padding={"left": 14, "top": 8, "right": 28, "bottom": 10})
        .configure_view(stroke=None, fill=CHART_BG_GREY)
        .configure_axis(labelFont="Calibri", titleFont="Calibri")
        .configure_legend(
            labelFont="Calibri",
            titleFont="Calibri",
            fillColor=CHART_BG_GREY,
            strokeColor=CHART_BG_GREY,
        )
        .configure(background=CHART_BG_GREY)
    )


def build_yield_curve_display_df(
    yield_curve_df: pd.DataFrame,
    selected_series: list[str],
    historical_curve_df: pd.DataFrame | None = None,
) -> pd.DataFrame:
    empty_cols = ["maturity_years", "yield_percent", "curve_type", "curve_date", "snapshot_label", "is_current", "snapshot_sort"]
    if yield_curve_df.empty and (historical_curve_df is None or historical_curve_df.empty):
        return pd.DataFrame(columns=empty_cols)

    frames = []
    if not yield_curve_df.empty:
        latest = yield_curve_df.copy()
        latest["snapshot_label"] = "Latest"
        latest["is_current"] = True
        latest["snapshot_sort"] = 0
        frames.append(latest)
    if historical_curve_df is not None and not historical_curve_df.empty:
        historical = historical_curve_df.copy()
        historical["snapshot_label"] = historical["snapshot_label"].astype(str)
        historical["is_current"] = False
        historical["snapshot_sort"] = pd.to_numeric(historical["snapshot_sort"], errors="coerce").fillna(1).astype(int)
        frames.append(historical)

    base = pd.concat(frames, ignore_index=True) if frames else pd.DataFrame(columns=empty_cols)
    if base.empty:
        return pd.DataFrame(columns=empty_cols)

    base["curve_date"] = pd.to_datetime(base["curve_date"], errors="coerce")
    base = base.dropna(subset=["maturity_years", "yield_percent", "curve_type", "curve_date"])
    if base.empty:
        return pd.DataFrame(columns=empty_cols)

    bei = build_breakeven_curve_points(base)
    combined = pd.concat([base, bei], ignore_index=True)
    if selected_series:
        combined = combined[combined["curve_type"].isin(selected_series)].copy()
    return combined.sort_values(["snapshot_sort", "curve_type", "maturity_years"]).reset_index(drop=True)


def ppt_rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


@st.cache_data(show_spinner=False)
def get_report_font_name() -> str:
    return "Calibri Light"


@st.cache_data(show_spinner=False)
def get_matplotlib_report_font_name() -> str:
    available_fonts = {f.name for f in font_manager.fontManager.ttflist}
    for candidate in ["Calibri Light", "Calibri", "Arial", "DejaVu Sans"]:
        if candidate in available_fonts:
            return candidate
    return "DejaVu Sans"


def mpl_colour(value: str) -> str:
    text = str(value).strip()
    if text.startswith("rgb(") and text.endswith(")"):
        inner = text[4:-1]
        parts = [p.strip() for p in inner.split(",")]
        if len(parts) == 3:
            try:
                r, g, b = [int(float(p)) for p in parts]
                return f"#{r:02x}{g:02x}{b:02x}"
            except ValueError:
                return "#e9e9e9"
    return text


def add_ppt_title(slide, title: str, subtitle: str | None = None) -> None:
    report_font = get_report_font_name()
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.55))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = report_font
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = ppt_rgb(TEXT_GREY)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12.0), Inches(0.3))
        p2 = subtitle_box.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = report_font
        run2.font.size = Pt(11)
        run2.font.color.rgb = ppt_rgb(TEXT_GREY)


def add_ppt_logo(slide, image_path: Path, left: float, top: float, width: float) -> None:
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def render_table_to_image(
    df: pd.DataFrame,
    title: str | None = None,
    percent_cols: list[str] | None = None,
    first_col_width: float | None = None,
) -> BytesIO:
    percent_cols = set(percent_cols or [])
    report_font = get_matplotlib_report_font_name()
    display_df = df.copy()
    if "asset_class" in display_df.columns:
        display_df["asset_class"] = display_df["asset_class"].map(display_name)
        display_df = display_df.rename(columns={"asset_class": "Asset class"})

    for col in display_df.columns:
        source_col = "asset_class" if col == "Asset class" else col
        if source_col in percent_cols:
            display_df[col] = display_df[col].map(lambda x: format_pct(x) if pd.notna(x) else "-")
        else:
            display_df[col] = display_df[col].map(lambda x: "-" if pd.isna(x) else x)

    rows, cols = display_df.shape
    fig_h = max(2.2, 0.34 * (rows + 2))
    fig_w = max(10.0, 1.15 * cols)
    with plt.rc_context({"font.family": [report_font, "Calibri", "Arial", "DejaVu Sans"]}):
        fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=180)
        ax.axis("off")
        if title:
            ax.set_title(title, fontsize=14, color=TEXT_GREY, loc="left", pad=12, fontweight="bold")

        table = ax.table(
            cellText=display_df.values,
            colLabels=display_df.columns,
            loc="upper left",
            cellLoc="center",
            colLoc="center",
            bbox=[0, 0, 1, 0.92 if title else 1],
        )
        table.auto_set_font_size(False)
        table.set_fontsize(8)
        if first_col_width is not None and len(display_df.columns) > 1:
            other_width = (1.0 - first_col_width) / (len(display_df.columns) - 1)
            for col_idx in range(len(display_df.columns)):
                width = first_col_width if col_idx == 0 else other_width
                for row_idx in range(len(display_df.index) + 1):
                    table[(row_idx, col_idx)].set_width(width)

        for (row, col), cell in table.get_celld().items():
            cell.set_edgecolor(MID_GREY)
            cell.set_linewidth(0.35)
            if row == 0:
                cell.set_facecolor(LIGHT_GREY)
                cell.set_text_props(weight="bold", color=TEXT_GREY)
            else:
                cell.set_facecolor(WHITE if row % 2 == 0 else "#fafafa")
                if col == 0:
                    cell.set_text_props(ha="left", color=TEXT_GREY)
                else:
                    cell.set_text_props(color=TEXT_GREY)

        plt.tight_layout()
        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
    buf.seek(0)
    return buf


def render_growth_chart_to_image(chart_df: pd.DataFrame) -> BytesIO:
    report_font = get_matplotlib_report_font_name()
    with plt.rc_context({"font.family": [report_font, "Calibri", "Arial", "DejaVu Sans"]}):
        fig, ax = plt.subplots(figsize=(11.6, 5.5), dpi=180)
        for asset_class, group in chart_df.groupby("asset_class"):
            group = group.sort_values("Date")
            ax.plot(group["Date"], group["Growth"], label=display_name(asset_class), linewidth=2.1)
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")
        ax.set_ylabel("Growth of wealth (GBP)", color=TEXT_GREY)
        ax.grid(False)
        ax.yaxis.set_minor_locator(AutoMinorLocator(2))
        ax.grid(which="minor", axis="y", linestyle=(0, (2, 4)), color=MID_GREY, linewidth=0.8)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color(MID_GREY)
        ax.spines["bottom"].set_color(MID_GREY)
        ax.tick_params(colors=TEXT_GREY, labelsize=10)
        ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.12), ncol=4, frameon=False, fontsize=9)
        fig.autofmt_xdate()
        plt.tight_layout()
        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
    buf.seek(0)
    return buf


def render_yield_chart_to_image(yield_df: pd.DataFrame) -> BytesIO:
    report_font = get_matplotlib_report_font_name()
    with plt.rc_context({"font.family": [report_font, "Calibri", "Arial", "DejaVu Sans"]}):
        fig, ax = plt.subplots(figsize=(11.6, 5.5), dpi=180)
        for curve_type, group in yield_df.groupby("curve_type"):
            group = group.sort_values("maturity_years")
            ax.plot(group["maturity_years"], pd.to_numeric(group["yield_percent"], errors="coerce") / 100.0, label=curve_type, linewidth=2.2)
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")
        ax.set_xlabel("Maturity (years)", color=TEXT_GREY)
        ax.set_ylabel("Yield", color=TEXT_GREY)
        ax.grid(axis="y", linestyle=(0, (2, 4)), color=MID_GREY, linewidth=0.8)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color(MID_GREY)
        ax.spines["bottom"].set_color(MID_GREY)
        ax.tick_params(colors=TEXT_GREY, labelsize=10)
        ax.yaxis.set_major_formatter(FuncFormatter(lambda y, _: f"{y:.0%}"))
        ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.12), ncol=3, frameon=False, fontsize=9)
        plt.tight_layout()
        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
    buf.seek(0)
    return buf


def build_market_commentary_bullets(ytd_returns_df: pd.DataFrame) -> list[str]:
    if ytd_returns_df.empty or "YTD" not in ytd_returns_df.columns:
        return [
            "Markets were mixed over the period, with equity and defensive returns diverging across regions.",
            "Use this slide for an analyst summary of the most important market developments.",
        ]

    table = ytd_returns_df[["asset_class", "YTD"]].dropna().copy()
    if table.empty:
        return ["YTD performance commentary could not be generated from the selected end date."]

    leaders = table.sort_values("YTD", ascending=False).head(3)
    laggards = table.sort_values("YTD", ascending=True).head(2)
    bullets = [
        f"{display_name(row.asset_class)} returned {format_pct(row.YTD)} YTD." for row in leaders.itertuples()
    ]
    bullets.extend(
        [f"{display_name(row.asset_class)} lagged at {format_pct(row.YTD)} YTD." for row in laggards.itertuples()]
    )

    cash_row = table[table["asset_class"] == "Cash (GBP)"]
    if not cash_row.empty:
        bullets.append(f"Cash (GBP) delivered {format_pct(float(cash_row['YTD'].iloc[0]))} YTD.")

    return bullets[:7]


def get_assets_from_display_groups(display_groups: list[dict]) -> list[str]:
    assets = []
    for group in display_groups:
        for item in group.get("items", []):
            if item not in assets:
                assets.append(item)
    return assets


def wrap_tile_label(label_text: str, max_chars: int = 16) -> str:
    text = str(label_text).strip()
    if not text or len(text) <= max_chars:
        return text
    words = text.split()
    lines = []
    current = ""
    for word in words:
        candidate = word if not current else f"{current} {word}"
        if len(candidate) <= max_chars:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return "\n".join(lines[:2])


def get_report_tile_label(label_text: str) -> str:
    text = str(label_text).strip()
    return REPORT_LABEL_OVERRIDES.get(text, text)


def render_snapshot_tiles_to_image(
    returns_df: pd.DataFrame,
    display_groups: list[dict],
    period_order: list[str] | None = None,
    layout_variant: str = "absolute",
) -> BytesIO:
    period_order = period_order or ["20Y", "10Y", "5Y", "YTD"]
    lookup = build_lookup_table(returns_df)
    report_font = get_matplotlib_report_font_name()
    with plt.rc_context({"font.family": [report_font, "Calibri", "Arial", "DejaVu Sans"]}):
        fig, axes = plt.subplots(1, len(period_order), figsize=(12.4, 6.7), dpi=180)
        if len(period_order) == 1:
            axes = [axes]

        fig.patch.set_facecolor("white")
        global_bottom = 1.0

        for ax, period in zip(axes, period_order):
            ax.set_xlim(0, 1)
            ax.set_ylim(0, 1)
            ax.axis("off")

            period_vals = returns_df[period].dropna() if period in returns_df.columns else pd.Series(dtype=float)
            vmin = float(period_vals.min()) if len(period_vals) else -0.05
            vmax = float(period_vals.max()) if len(period_vals) else 0.15

            ax.text(0.5, 0.985, DASHBOARD_HORIZONS.get(period, period), ha="center", va="top", fontsize=16, color="black", fontweight="medium")

            y = 0.93
            group_gap = 0.02 if layout_variant == "relative_minor" else 0.018
            title_gap = 0.026 if layout_variant == "relative_minor" else 0.028
            tile_h_big = 0.074 if layout_variant == "relative_minor" else 0.072
            tile_h_small = 0.068 if layout_variant == "relative_minor" else 0.067

            for group in display_groups:
                title = group["title"]
                items = group["items"]
                labels = group["labels"]

                if title:
                    ax.text(0.5, y, title, ha="center", va="top", fontsize=10.5, color="black", fontweight="medium")
                    y -= title_gap
                else:
                    y -= 0.008

                def draw_tile(x, y_top, w, h, value, label_text="", plain=False):
                    colour = mpl_colour(heat_colour(value, vmin, vmax))
                    rect = plt.Rectangle((x, y_top - h), w, h, facecolor=colour, edgecolor="none")
                    ax.add_patch(rect)
                    if label_text:
                        display_label = get_report_tile_label(label_text)
                        wrapped_label = wrap_tile_label(display_label, 18 if layout_variant == "relative_minor" else 15)
                        line_count = wrapped_label.count("\n") + 1
                        label_font = (
                            6.9 if layout_variant == "relative_minor" and line_count > 1 else
                            7.4 if layout_variant == "relative_minor" else
                            6.9 if line_count > 1 else 7.4
                        )
                        ax.text(
                            x + w / 2,
                            y_top - (0.007 if layout_variant == "relative_minor" else 0.008),
                            wrapped_label,
                            ha="center",
                            va="top",
                            fontsize=7.8 if plain and layout_variant != "relative_minor" else label_font,
                            color="black",
                            fontweight="bold" if plain else "normal",
                            linespacing=0.9,
                        )
                    else:
                        line_count = 0
                    ax.text(
                        x + w / 2,
                        y_top - h + (0.013 if layout_variant == "relative_minor" else 0.011),
                        format_pct(value),
                        ha="center",
                        va="bottom",
                        fontsize=11.2 if layout_variant == "relative_minor" else 12,
                        color="white",
                        fontweight="bold",
                    )

                if len(items) == 1:
                    item = items[0]
                    val = lookup.get(item, {}).get(period, np.nan)
                    draw_tile(0.25, y, 0.5, tile_h_big, val, labels.get(item, ""), plain=False)
                    y -= tile_h_big + group_gap
                elif len(items) == 2:
                    for idx, item in enumerate(items):
                        val = lookup.get(item, {}).get(period, np.nan)
                        draw_tile(0.03 + idx * 0.48, y, 0.44, tile_h_small, val, labels.get(item, item), plain=True)
                    y -= tile_h_small + group_gap
                elif len(items) == 3:
                    broad = items[0]
                    draw_tile(0.0325, y, 0.935, tile_h_big, lookup.get(broad, {}).get(period, np.nan), labels.get(broad, ""), plain=False)
                    y -= tile_h_big + 0.01
                    for idx, item in enumerate(items[1:]):
                        draw_tile(0.03 + idx * 0.48, y, 0.44, tile_h_small, lookup.get(item, {}).get(period, np.nan), labels.get(item, ""), plain=False)
                    y -= tile_h_small + group_gap
                elif len(items) == 4:
                    if layout_variant == "relative_minor":
                        relative_tile_h = 0.084
                        relative_gap = 0.018
                        x_positions = [0.04, 0.525]
                        widths = [0.405, 0.405]
                        for row in range(2):
                            for col in range(2):
                                idx = row * 2 + col
                                item = items[idx]
                                draw_tile(
                                    x_positions[col],
                                    y,
                                    widths[col],
                                    relative_tile_h,
                                    lookup.get(item, {}).get(period, np.nan),
                                    labels.get(item, item),
                                    plain=True,
                                )
                            y -= relative_tile_h + relative_gap
                        y -= group_gap - relative_gap
                    else:
                        for row in range(2):
                            for col in range(2):
                                idx = row * 2 + col
                                item = items[idx]
                                draw_tile(0.03 + col * 0.48, y, 0.44, tile_h_small, lookup.get(item, {}).get(period, np.nan), labels.get(item, item), plain=True)
                            y -= tile_h_small + 0.008
                        y -= group_gap - 0.008
                elif len(items) == 5:
                    widths = [0.295, 0.295, 0.295]
                    x_positions = [0.02, 0.3525, 0.685]
                    for idx in range(3):
                        item = items[idx]
                        draw_tile(x_positions[idx], y, widths[idx], tile_h_small, lookup.get(item, {}).get(period, np.nan), labels.get(item, item), plain=True)
                    y -= tile_h_small + 0.008
                    for idx in range(2):
                        item = items[idx + 3]
                        draw_tile(0.16 + idx * 0.36, y, 0.31, tile_h_small, lookup.get(item, {}).get(period, np.nan), labels.get(item, item), plain=True)
                    y -= tile_h_small + group_gap

            global_bottom = min(global_bottom, y - 0.01)

        if layout_variant == "relative_minor":
            ylim_bottom = max(global_bottom, 0.42)
            for ax in axes:
                ax.set_ylim(ylim_bottom, 1.0)

        plt.tight_layout(w_pad=0.7)
        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
    buf.seek(0)
    return buf


def render_relative_minor_snapshot_to_image(
    returns_df: pd.DataFrame,
    period_order: list[str] | None = None,
) -> BytesIO:
    period_order = period_order or ["20Y", "10Y", "5Y", "YTD"]
    lookup = build_lookup_table(returns_df)
    report_font = get_matplotlib_report_font_name()

    with plt.rc_context({"font.family": [report_font, "Calibri", "Arial", "DejaVu Sans"]}):
        fig, axes = plt.subplots(1, len(period_order), figsize=(12.4, 5.7), dpi=180)
        if len(period_order) == 1:
            axes = [axes]
        fig.patch.set_facecolor("white")

        for ax, period in zip(axes, period_order):
            ax.set_xlim(0, 1)
            ax.set_ylim(0, 1)
            ax.axis("off")

            period_vals = returns_df[period].dropna() if period in returns_df.columns else pd.Series(dtype=float)
            vmin = float(period_vals.min()) if len(period_vals) else -0.05
            vmax = float(period_vals.max()) if len(period_vals) else 0.15
            ax.text(0.5, 0.985, DASHBOARD_HORIZONS.get(period, period), ha="center", va="top", fontsize=16, color="black", fontweight="medium")

            def draw_tile(x, y_top, w, h, value, label_text):
                rect = plt.Rectangle((x, y_top - h), w, h, facecolor=mpl_colour(heat_colour(value, vmin, vmax)), edgecolor="none")
                ax.add_patch(rect)
                ax.text(
                    x + w / 2,
                    y_top - 0.012,
                    get_report_tile_label(label_text),
                    ha="center",
                    va="top",
                    fontsize=6.7,
                    color="black",
                    fontweight="bold",
                )
                ax.text(
                    x + w / 2,
                    y_top - h + 0.018,
                    format_pct(value),
                    ha="center",
                    va="bottom",
                    fontsize=10.6,
                    color="white",
                    fontweight="bold",
                )

            # Relative to DM
            ax.text(0.5, 0.91, "Relative to DM", ha="center", va="top", fontsize=10.5, color="black", fontweight="medium")
            y = 0.865
            top_items = [
                ("UK stocks", "UK"),
                ("Emerging stocks", "EM"),
                ("Developed REITs", "REIT"),
            ]
            for idx, (asset, label) in enumerate(top_items):
                draw_tile(0.02 + idx * 0.325, y, 0.29, 0.08, lookup.get(asset, {}).get(period, np.nan), label)
            y -= 0.096
            bottom_items = [
                ("Developed value stocks", "DM Value"),
                ("Developed small stocks", "DM Small"),
            ]
            for idx, (asset, label) in enumerate(bottom_items):
                draw_tile(0.18 + idx * 0.34, y, 0.29, 0.08, lookup.get(asset, {}).get(period, np.nan), label)

            # Relative to EM
            ax.text(0.5, 0.63, "Relative to EM", ha="center", va="top", fontsize=10.5, color="black", fontweight="medium")
            y = 0.585
            for idx, (asset, label) in enumerate([("Emerging value stocks", "EM Value"), ("Emerging small stocks", "EM Small")]):
                draw_tile(0.06 + idx * 0.46, y, 0.38, 0.08, lookup.get(asset, {}).get(period, np.nan), label)

            # Relative to UK
            ax.text(0.5, 0.46, "Relative to UK", ha="center", va="top", fontsize=10.5, color="black", fontweight="medium")
            y = 0.415
            for idx, (asset, label) in enumerate([("UK value stocks", "UK Value"), ("UK small stocks", "UK Small")]):
                draw_tile(0.06 + idx * 0.46, y, 0.38, 0.08, lookup.get(asset, {}).get(period, np.nan), label)

            # Defensive block
            y = 0.285
            defensive_items = [
                ("Cash (GBP)", "Cash"),
                ("UK Gilts (0-5)", "UK Gilts"),
                ("UK IL Gilts (0-5)", "UK ILG (0-5)"),
                ("Global GBP hedged bonds (0-5)", "GSDB (0-5)"),
            ]
            for row in range(2):
                for col in range(2):
                    idx = row * 2 + col
                    asset, label = defensive_items[idx]
                    draw_tile(0.04 + col * 0.48, y, 0.40, 0.095, lookup.get(asset, {}).get(period, np.nan), label)
                y -= 0.115

        plt.tight_layout(w_pad=0.7)
        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
    buf.seek(0)
    return buf


@st.cache_data(show_spinner=False)
def build_quarterly_market_metrics_report(
    report_end_date_text: str,
    report_end_date_long_text: str,
    report_month_text: str,
    nominal_abs_df: pd.DataFrame,
    real_abs_df: pd.DataFrame,
    relative_minor_df: pd.DataFrame,
    returns_nominal_df: pd.DataFrame,
    returns_real_df: pd.DataFrame,
    growth_df: pd.DataFrame,
    yield_df: pd.DataFrame,
    bullet_points: list[str],
    methodology_notes: list[str],
    index_notes: list[str],
    further_notes: list[str],
    snapshot_assets: list[str],
) -> bytes:
    report_font = get_report_font_name()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    add_ppt_logo(slide, ALBION_LOGO_FILE, 0.65, 0.45, 2.0)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Quarterly Market Metrics"
    run.font.name = report_font
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ppt_rgb(TEXT_GREY)
    p.alignment = PP_ALIGN.CENTER
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(11.7), Inches(0.5))
    p2 = subtitle.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = f"Data to {report_end_date_long_text}"
    run2.font.name = report_font
    run2.font.size = Pt(16)
    run2.font.color.rgb = ppt_rgb(TEXT_GREY)
    p2.alignment = PP_ALIGN.CENTER
    add_ppt_logo(slide, POWERED_BY_FILE, 4.8, 6.45, 3.6)

    # 2. Commentary
    slide = prs.slides.add_slide(blank)
    add_ppt_title(slide, "Market commentary")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.3))
    tf = box.text_frame
    for idx, bullet in enumerate(bullet_points or ["Add analyst commentary here."]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.name = report_font
        p.font.size = Pt(16)
        p.font.color.rgb = ppt_rgb(TEXT_GREY)

    snapshot_slide_specs = [
        ("Absolute annualised nominal returns", nominal_abs_df, REPORT_DISPLAY_GROUPS_ABSOLUTE, "absolute", 0.90),
        ("Absolute annualised real returns", real_abs_df, REPORT_DISPLAY_GROUPS_ABSOLUTE, "absolute", 0.90),
    ]

    for slide_title, df, groups, layout_variant, top_pos in snapshot_slide_specs:
        slide = prs.slides.add_slide(blank)
        add_ppt_title(slide, slide_title)
        slide.shapes.add_picture(
            render_snapshot_tiles_to_image(df, groups, layout_variant=layout_variant),
            Inches(0.45),
            Inches(top_pos),
            width=Inches(12.35),
        )

    slide = prs.slides.add_slide(blank)
    add_ppt_title(slide, "Relative annualised returns")
    slide.shapes.add_picture(
        render_relative_minor_snapshot_to_image(relative_minor_df),
        Inches(0.45),
        Inches(1.05),
        width=Inches(12.35),
    )

    filtered_nominal = returns_nominal_df[returns_nominal_df["asset_class"].isin(snapshot_assets)].copy()
    filtered_real = returns_real_df[returns_real_df["asset_class"].isin(snapshot_assets)].copy()
    slide_specs = [
        ("Annualised returns (nominal)", filtered_nominal, [c for c in filtered_nominal.columns if c != "asset_class"]),
        ("Annualised returns (real)", filtered_real, [c for c in filtered_real.columns if c != "asset_class"]),
    ]

    for slide_title, df, percent_cols in slide_specs:
        slide = prs.slides.add_slide(blank)
        add_ppt_title(slide, slide_title)
        img = render_table_to_image(df, percent_cols=percent_cols, first_col_width=0.24)
        slide.shapes.add_picture(img, Inches(0.45), Inches(1.2), width=Inches(12.35))

    slide = prs.slides.add_slide(blank)
    add_ppt_title(slide, "Growth of wealth (10 years, nominal)")
    slide.shapes.add_picture(render_growth_chart_to_image(growth_df), Inches(0.5), Inches(1.2), width=Inches(12.2))

    slide = prs.slides.add_slide(blank)
    add_ppt_title(slide, f"UK yield curve (as at {report_end_date_text})")
    slide.shapes.add_picture(render_yield_chart_to_image(yield_df), Inches(0.5), Inches(1.2), width=Inches(12.2))

    slide = prs.slides.add_slide(blank)
    add_ppt_title(slide, "Important notes")
    sections = [
        ("Methodology", methodology_notes, (0.8, 1.3, 5.8, 2.0)),
        ("Indices used", index_notes, (0.8, 3.45, 5.8, 2.7)),
        ("Further information", further_notes, (6.75, 1.3, 5.75, 4.85)),
    ]
    for heading, lines, (left, top, width, height) in sections:
        heading_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.28))
        hp = heading_box.text_frame.paragraphs[0]
        hrun = hp.add_run()
        hrun.text = heading
        hrun.font.name = report_font
        hrun.font.size = Pt(14)
        hrun.font.bold = True
        hrun.font.color.rgb = ppt_rgb(TEXT_GREY)

        body_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.32), Inches(width), Inches(height))
        tf_notes = body_box.text_frame
        tf_notes.word_wrap = True
        for idx, note in enumerate(lines):
            p = tf_notes.paragraphs[0] if idx == 0 else tf_notes.add_paragraph()
            p.text = note
            p.level = 0
            for run in p.runs:
                run.font.name = report_font
            p.font.size = Pt(11.5)
            p.font.color.rgb = ppt_rgb(TEXT_GREY)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# =====================================
# PAGE SETUP
# =====================================
st.set_page_config(page_title=APP_TITLE, layout="wide")

st.markdown(
    f"""
    <style>
    html, body, .stApp,
    [data-testid="stAppViewContainer"],
    [data-testid="stMarkdownContainer"],
    p, li, label, button, input, textarea, table {{
        font-family: "Calibri Light", Calibri, "Segoe UI", Arial, sans-serif !important;
    }}

    .stApp {{ background-color: white; }}

    header[data-testid="stHeader"],
    div[data-testid="stToolbar"],
    div[data-testid="stDecoration"],
    div[data-testid="stStatusWidget"] {{
        display: none !important;
        visibility: hidden !important;
        height: 0 !important;
    }}

    #MainMenu {{ visibility: hidden !important; }}

    [data-testid="stAppViewContainer"] > .main {{ padding-top: 0 !important; }}

    .block-container {{
        padding-top: 0.25rem !important;
        padding-bottom: 1rem !important;
        max-width: 1500px;
    }}

    .top-header-grid {{
        display: grid;
        grid-template-columns: 1fr auto;
        align-items: center;
        gap: 18px;
        margin-bottom: 2px;
    }}

    .top-title-wrap {{
        display: flex;
        align-items: center;
    }}

    .dashboard-title {{
        font-size: 40px;
        font-weight: 500;
        margin-bottom: 2px;
        color: black;
        line-height: 1.15;
        padding-top: 0;
    }}

    .header-logo {{
        display: flex;
        justify-content: flex-end;
        align-items: center;
        padding-top: 0;
    }}

    .header-logo img {{
        max-height: 50px;
        width: auto;
        object-fit: contain;
        display: block;
    }}

    .toolbar-label {{
        font-size: 13px;
        font-weight: 700;
        color: #444;
        margin-bottom: 6px;
    }}

    .toolbar-label-muted {{ color: #888 !important; }}

    .toolbar-meta {{
        text-align: right;
        font-size: 13px;
        color: {TEXT_GREY};
        padding-top: 22px;
        line-height: 1.2;
        white-space: nowrap;
    }}

    div[data-testid="stCheckbox"] label,
    div[data-testid="stCheckbox"] label p,
    div[data-testid="stCheckbox"] span {{
        color: #333 !important;
    }}

    .snapshot-toolbar-note {{
        font-size: 13px;
        color: {TEXT_GREY};
        margin: 0 0 14px 0;
        line-height: 1.25;
    }}

    .factor-style-shell {{
        width: 100%;
        border: 0.2px solid {MID_GREY};
        border-radius: 12px;
        padding: 12px;
        margin-bottom: 10px;
        background: {WHITE};
        box-sizing: border-box;
    }}

    .factor-style-grid {{
        display: grid;
        grid-template-columns: minmax(92px, 0.95fr) repeat(3, minmax(0, 1fr));
        gap: 8px;
        align-items: stretch;
    }}

    .factor-style-header {{
        font-size: 15px;
        font-weight: 700;
        color: {TEXT_GREY};
        padding: 4px 2px 8px 2px;
        text-align: center;
    }}

    .factor-style-corner {{
        text-align: left;
    }}

    .factor-style-row-label {{
        display: flex;
        align-items: center;
        font-size: 15px;
        font-weight: 700;
        color: {TEXT_GREY};
        padding: 0 4px;
    }}

    .factor-style-cell {{
        min-height: 72px;
        border-radius: 10px;
        display: flex;
        align-items: center;
        justify-content: center;
        font-size: 24px;
        font-weight: 700;
        line-height: 1;
        box-sizing: border-box;
    }}

    .factor-style-cell-empty {{
        background: {LIGHT_GREY};
        color: #9a9a9a;
    }}

    .country-card-grid {{
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(126px, 1fr));
        gap: 8px;
        margin-bottom: 12px;
    }}

    .country-card {{
        border-radius: 10px;
        padding: 10px 10px 9px 10px;
        min-height: 86px;
        color: #ffffff;
    }}

    .country-card-title {{
        font-size: 14px;
        font-weight: 700;
        line-height: 1.15;
        margin-bottom: 3px;
    }}

    .country-card-meta {{
        font-size: 10px;
        font-weight: 600;
        opacity: 0.92;
        margin-bottom: 10px;
        letter-spacing: 0.03em;
    }}

    .country-card-value {{
        font-size: 24px;
        font-weight: 700;
        line-height: 1.0;
        margin-bottom: 0;
    }}

    .region-card-shell {{
        display: flex;
        flex-direction: column;
        gap: 10px;
        margin-bottom: 12px;
    }}

    .region-card-row {{
        display: flex;
        gap: 10px;
        width: 100%;
        align-items: stretch;
    }}

    .region-global-card,
    .region-card {{
        border-radius: 12px;
        padding: 12px 14px;
        min-height: 118px;
        color: #ffffff;
        box-sizing: border-box;
    }}

    .region-card-title {{
        font-size: 20px;
        font-weight: 800;
        line-height: 1.1;
        margin-bottom: 4px;
    }}

    .region-card-meta {{
        font-size: 12px;
        font-weight: 600;
        opacity: 0.86;
        margin-bottom: 18px;
        letter-spacing: 0.02em;
    }}

    .region-card-value {{
        font-size: 30px;
        font-weight: 800;
        line-height: 1.0;
        margin-bottom: 0;
    }}

    .patchwork-shell {{
        margin-bottom: 14px;
    }}

    .patchwork-grid {{
        display: grid;
        gap: 8px 8px;
        align-items: start;
        margin-bottom: 12px;
    }}

    .patchwork-year-header {{
        font-size: 16px;
        font-weight: 800;
        color: #444;
        text-align: center;
        padding-bottom: 3px;
    }}

    .patchwork-metric-header {{
        min-width: 100px;
    }}

    .patchwork-year-col {{
        display: flex;
        flex-direction: column;
        gap: 6px;
    }}

    .patchwork-metric-col {{
        min-width: 100px;
    }}

    .patchwork-cell {{
        border-radius: 8px;
        min-height: 52px;
        padding: 7px 8px 6px 8px;
        color: #ffffff;
        display: flex;
        flex-direction: column;
        justify-content: space-between;
        box-sizing: border-box;
    }}

    .patchwork-cell-country {{
        font-size: 11px;
        font-weight: 800;
        line-height: 1.1;
    }}

    .patchwork-cell-return {{
        font-size: 18px;
        font-weight: 800;
        line-height: 1.0;
    }}

    .patchwork-summary-cell {{
        min-height: 52px;
    }}

    .period-shell {{
        background: transparent;
        padding: 8px 8px 12px 8px;
        min-height: 100%;
    }}

    .period-title {{
        text-align: center;
        font-size: 28px;
        font-weight: 500;
        margin-bottom: 10px;
        color: black;
    }}

    .group-card {{
        padding: 10px 8px 8px 8px;
        margin-bottom: 10px;
        background: #e7e7e7;
    }}

    .section-title {{
        text-align: center;
        font-size: 18px;
        font-weight: 500;
        margin-top: 0;
        margin-bottom: 4px;
        color: black;
    }}

    .section-title-empty {{
        height: 4px;
        margin-bottom: 2px;
    }}

    .section-subtitle {{
        text-align: center;
        font-size: 13px;
        font-style: italic;
        color: #444;
        margin-bottom: 6px;
    }}

    .big-tile, .small-tile {{
        color: white;
        text-align: center;
        font-weight: 700;
        border-radius: 0;
        line-height: 1.2;
    }}

    .big-tile {{
        padding: 10px 8px;
        margin-bottom: 8px;
        font-size: 22px;
    }}

    .small-tile {{
        padding: 8px 6px;
        font-size: 18px;
        margin-bottom: 8px;
    }}

    .tile-label, .tile-label-on-colour {{
        display: block;
        color: black !important;
        font-size: 13px;
        font-style: italic;
        font-weight: 500;
        margin-bottom: 4px;
    }}

    .tile-label-plain {{
        font-style: normal !important;
        font-weight: 700 !important;
    }}

    .spacer {{ height: 2px; }}

    .page-section-title {{
        font-size: 20px;
        font-weight: 600;
        color: {TEXT_GREY};
        margin: 6px 0 10px 0;
    }}

    .methodology-text {{
        margin-top: 18px;
        margin-bottom: 10px;
        font-size: 13px;
        color: {TEXT_GREY} !important;
        text-align: center;
        line-height: 1.45;
    }}

    .methodology-text a {{
        color: {BRAND_ORANGE_DARK};
        text-decoration: none;
    }}

    .footer-bar {{
        margin-top: 20px;
        display: flex;
        justify-content: center;
        align-items: center;
        width: 100%;
    }}

    .footer-logo img {{
        max-height: 48px;
        width: auto;
        object-fit: contain;
        display: block;
    }}

    .diag-title {{
        font-size: 20px;
        font-weight: 500;
        color: black !important;
        margin-top: 4px;
        margin-bottom: 8px;
    }}

    .diag-note {{
        font-size: 13px;
        color: {TEXT_GREY} !important;
        margin-bottom: 14px;
        line-height: 1.4;
    }}

    .stButton button {{
        border-radius: 6px !important;
        font-weight: 600 !important;
    }}

    [data-testid="collapsedControl"] {{
        display: none !important;
    }}

    section[data-testid="stSidebar"] {{
        min-width: 58px !important;
        max-width: 58px !important;
        transition: min-width 0.18s ease, max-width 0.18s ease !important;
        overflow-x: hidden !important;
        background: #ffffff !important;
        border-right: 1px solid {MID_GREY} !important;
    }}

    section[data-testid="stSidebar"]:hover {{
        min-width: 228px !important;
        max-width: 228px !important;
    }}

    section[data-testid="stSidebar"] > div:first-child {{
        width: 58px !important;
        transition: width 0.18s ease !important;
        overflow-x: hidden !important;
    }}

    section[data-testid="stSidebar"]:hover > div:first-child {{
        width: 228px !important;
    }}

    .sidebar-hamburger {{
        font-size: 24px;
        font-weight: 700;
        color: #111827;
        margin: 2px 0 14px 2px;
        line-height: 1;
    }}

    .sidebar-nav-title {{
        font-size: 11px;
        font-weight: 700;
        letter-spacing: 0.08em;
        text-transform: uppercase;
        color: #6b7280;
        margin: 10px 0 8px 2px;
        white-space: nowrap;
    }}

    section[data-testid="stSidebar"] .sidebar-nav-title,
    section[data-testid="stSidebar"] div[data-testid="stButton"] {{
        opacity: 0;
        visibility: hidden;
        max-height: 0;
        overflow: hidden;
        pointer-events: none;
        margin: 0 !important;
        transition: opacity 0.12s ease, max-height 0.12s ease, visibility 0.12s ease;
    }}

    section[data-testid="stSidebar"]:hover .sidebar-nav-title,
    section[data-testid="stSidebar"]:hover div[data-testid="stButton"] {{
        opacity: 1;
        visibility: visible;
        max-height: 80px;
        pointer-events: auto;
    }}

    .stSidebar .stButton button {{
        min-height: 40px !important;
        width: 100% !important;
        justify-content: flex-start !important;
        padding-left: 12px !important;
    }}

    .stSidebar .stButton button[kind="primary"] {{
        background: {BRAND_ORANGE} !important;
        border: 1px solid {BRAND_ORANGE} !important;
        color: #ffffff !important;
    }}

    .stSidebar .stButton button[kind="secondary"] {{
        background: #f7f7f7 !important;
        border: 1px solid #d8d8d8 !important;
        color: #111111 !important;
    }}

    .stDownloadButton button {{
        background: #f2f2f2 !important;
        color: black !important;
        border: 1px solid #cfcfcf !important;
        border-radius: 4px !important;
        font-weight: 500 !important;
    }}

    .stDownloadButton button * {{
        color: black !important;
    }}

    .stDownloadButton button:hover,
    .stButton button:hover {{
        border-color: {BRAND_ORANGE} !important;
    }}

    .stTabs [data-baseweb="tab-list"] {{
        gap: 6px;
    }}

    .stTabs [data-baseweb="tab"] {{
        color: {TEXT_GREY} !important;
        background: #f3f3f3 !important;
        border: 0.2px solid {MID_GREY} !important;
        border-radius: 8px 8px 0 0 !important;
        padding: 0.45rem 0.8rem !important;
    }}

    .stTabs [aria-selected="true"] {{
        color: #111111 !important;
        background: #ffffff !important;
        border-bottom-color: #ffffff !important;
        font-weight: 600 !important;
    }}

    .stTabs [data-baseweb="tab"] * {{
        color: inherit !important;
        -webkit-text-fill-color: inherit !important;
    }}

    .stMetric, .stMetric * {{
        color: black !important;
    }}

    h1, h2, h3, h4, h5, h6,
    [data-testid="stHeading"],
    [data-testid="stHeading"] * {{
        color: #000000 !important;
        -webkit-text-fill-color: #000000 !important;
    }}

    .stMultiSelect label, .stSelectbox label {{
        color: {TEXT_GREY} !important;
    }}

    .st-key-diagnostics_toggle_button button {{
        background: #111827 !important;
        color: #ffffff !important;
        border: 1px solid #111827 !important;
        border-radius: 4px !important;
        font-weight: 600 !important;
        text-align: left !important;
        justify-content: flex-start !important;
        padding: 0.55rem 0.9rem !important;
        margin: 0 !important;
    }}

    .st-key-diagnostics_toggle_button button *,
    .st-key-diagnostics_toggle_button button span {{
        color: #ffffff !important;
        -webkit-text-fill-color: #ffffff !important;
    }}

    div[data-testid="stVerticalBlock"] > div:empty {{
        display: none !important;
    }}

    div[data-testid="stButton"],
    div[data-testid="stButton"] > div {{
        margin: 0 !important;
        padding: 0 !important;
    }}

    div[data-testid="stButton"] button {{
        margin-bottom: 0 !important;
    }}

    div[data-testid="stVegaLiteChart"] {{
        background: {CHART_BG_GREY} !important;
        border: 0.2px solid {MID_GREY} !important;
        border-radius: 12px !important;
        padding: 10px 18px 8px 12px !important;
        overflow: hidden !important;
    }}

    div[data-testid="stVegaLiteChart"] canvas,
    div[data-testid="stVegaLiteChart"] svg {{
        border-radius: 10px !important;
    }}

    .table-shell {{
        width: 100%;
        border: 0.2px solid {MID_GREY};
        border-radius: 12px;
        overflow: hidden;
        margin-bottom: 8px;
    }}

    .custom-data-table {{
        width: 100%;
        border-collapse: collapse;
        table-layout: fixed;
    }}

    .custom-data-table thead th {{
        background: {LIGHT_GREY};
        color: {TEXT_GREY};
        border: 0.2px solid {MID_GREY};
        padding: 10px 12px;
        text-align: center;
        font-weight: 600;
        white-space: normal;
        line-height: 1.2;
        overflow: hidden;
        text-overflow: ellipsis;
    }}

    .custom-data-table tbody td {{
        color: {TEXT_GREY};
        border: 0.2px solid {MID_GREY};
        padding: 10px 12px;
        text-align: center;
        white-space: nowrap;
        overflow: hidden;
        text-overflow: ellipsis;
    }}

    .custom-data-table tbody tr:nth-child(odd) td {{
        background: {LIGHT_GREY};
    }}

    .custom-data-table tbody tr:nth-child(even) td {{
        background: {WHITE};
    }}

    .custom-data-table th:first-child,
    .custom-data-table td:first-child {{
        text-align: left;
    }}

    .table-empty {{
        padding: 16px;
        color: {TEXT_GREY};
        text-align: center;
        background: {WHITE};
    }}
    </style>
    """,
    unsafe_allow_html=True,
)


# =====================================
# LOAD + PREP
# =====================================
if not DATA_FILE.exists():
    st.error(f"File not found: {DATA_FILE}")
    st.stop()

current_file_mtime = DATA_FILE.stat().st_mtime

try:
    ts, mapping = load_data(str(DATA_FILE), current_file_mtime)
except Exception as exc:
    st.exception(exc)
    st.stop()

regions_df = load_regions_data(str(DATA_FILE), current_file_mtime)
sectors_df = load_sectors_data(str(DATA_FILE), current_file_mtime)
factors_df = load_factors_data(str(DATA_FILE), current_file_mtime)

monthly_levels = build_monthly_index_levels(ts, mapping)
if not monthly_levels:
    st.error(
        "No mapped asset classes were found. Check that the mapping sheet asset_class names "
        "match the expected names in DISPLAY_GROUPS / ASSET_CLASS_ALIASES."
    )
    st.stop()

inflation_levels = None
inflation_source_note = "Workbook time_series"
inflation_debug_message = None

try:
    inflation_levels, inflation_source_note, inflation_debug_message = build_best_available_inflation_levels(ts)
except Exception as exc:
    st.warning(f"Inflation series could not be built. Real mode may be unavailable. Details: {exc}")

inflation_monthly_returns = (
    build_monthly_returns_from_levels(inflation_levels)
    if inflation_levels is not None and not inflation_levels.dropna().empty
    else pd.Series(dtype=float)
)

primary_tickers = mapping.get("live_fund_primary", pd.Series(dtype=str)).dropna().tolist()
secondary_tickers = mapping.get("live_fund_secondary", pd.Series(dtype=str)).dropna().tolist()
all_live_tickers = tuple(sorted({normalise_ticker(x) for x in primary_tickers + secondary_tickers if normalise_ticker(x)}))

try:
    live_prices = fetch_yf_prices(all_live_tickers, ETF_DOWNLOAD_START)
except Exception as exc:
    st.warning(f"Live price download failed. Falling back to index-only results. Details: {exc}")
    live_prices = pd.DataFrame()

stitched_series_map, live_diag = build_stitched_asset_series(
    monthly_levels=monthly_levels,
    mapping=mapping,
    prices_df=live_prices,
)

chart_series_map, chart_diag = build_chart_preferred_series(
    monthly_levels=monthly_levels,
    mapping=mapping,
    prices_df=live_prices,
)

if not stitched_series_map:
    st.error("No stitched or index-only asset series could be built.")
    st.stop()

available_assets = [a for a in DEFAULT_ASSET_ORDER if a in stitched_series_map or a in chart_series_map]
default_chart_assets = [a for a in DEFAULT_CHART_ASSETS if a in available_assets]
whole_period_start = pd.Timestamp("1989-07-31")
common_inception_text = whole_period_start.strftime("%d/%m/%Y")

asset_coverage_diag = build_asset_coverage_table(
    mapping=mapping,
    monthly_levels=monthly_levels,
    stitched_series_map=stitched_series_map,
    chart_series_map=chart_series_map,
    live_diag=live_diag,
    chart_diag=chart_diag,
)
mapping_diag = build_mapping_diagnostics_table(mapping, ts)
live_price_diag = build_live_price_diagnostics(live_prices)
ons_fetch_summary, ons_fetch_preview = build_ons_fetch_diagnostics(ONS_CPI_INDEX_CSV_URL)
yield_curve_df, boe_yield_summary, boe_yield_preview = build_boe_yield_curve_diagnostics(
    BOE_YIELD_CURVE_ZIP_URL,
    DIVIDENDDATA_INDEX_LINKED_GILTS_URL,
)
yield_curve_history_df, boe_month_end_yield_summary, boe_month_end_yield_preview = build_boe_month_end_yield_curve_history(
    BOE_NOMINAL_MONTH_END_ZIP_URL,
    BOE_REAL_MONTH_END_ZIP_URL,
)
global_yield_curve_df = pd.DataFrame(columns=["country", "maturity_label", "maturity_years", "yield_percent", "curve_date", "source_url"])
global_yield_summary = pd.DataFrame(columns=["metric", "value"])
global_yield_preview = pd.DataFrame(columns=["country", "maturity_label", "maturity_years", "yield_percent", "source_url"])
asset_style_map = (
    mapping.drop_duplicates(subset=["asset_class"], keep="last")
    .set_index("asset_class")
    .get("growth_defensive", pd.Series(dtype=float))
    .to_dict()
)


# =====================================
# TOP HEADER
# =====================================
logo_html = ""
if ALBION_LOGO_FILE.exists():
    logo_html = (
        f'<div class="header-logo"><img src="data:image/png;base64,'
        f'{img_to_base64(str(ALBION_LOGO_FILE), ALBION_LOGO_FILE.stat().st_mtime)}"></div>'
    )

if "top_page_selector" not in st.session_state:
    st.session_state["top_page_selector"] = "Dashboard"

with st.sidebar:
    st.markdown('<div class="sidebar-hamburger">&#9776;</div>', unsafe_allow_html=True)
    st.markdown('<div class="sidebar-nav-title">Navigation</div>', unsafe_allow_html=True)
    if st.button(
        PAGE_LABELS["Dashboard"],
        key="sidebar_page_dashboard_btn",
        use_container_width=True,
        type="primary" if st.session_state["top_page_selector"] == "Dashboard" else "secondary",
    ):
        st.session_state["top_page_selector"] = "Dashboard"
        st.rerun()
    if st.button(
        PAGE_LABELS["Charts"],
        key="sidebar_page_charts_btn",
        use_container_width=True,
        type="primary" if st.session_state["top_page_selector"] == "Charts" else "secondary",
    ):
        st.session_state["top_page_selector"] = "Charts"
        st.rerun()
    if st.button(
        PAGE_LABELS["Risk"],
        key="sidebar_page_risk_btn",
        use_container_width=True,
        type="primary" if st.session_state["top_page_selector"] == "Risk" else "secondary",
    ):
        st.session_state["top_page_selector"] = "Risk"
        st.rerun()
    if st.button(
        PAGE_LABELS["Factors"],
        key="sidebar_page_factors_btn",
        use_container_width=True,
        type="primary" if st.session_state["top_page_selector"] == "Factors" else "secondary",
    ):
        st.session_state["top_page_selector"] = "Factors"
        st.rerun()
    if st.button(
        PAGE_LABELS["Geo"],
        key="sidebar_page_geo_btn",
        use_container_width=True,
        type="primary" if st.session_state["top_page_selector"] == "Geo" else "secondary",
    ):
        st.session_state["top_page_selector"] = "Geo"
        st.rerun()
    if st.button(
        PAGE_LABELS["Sector"],
        key="sidebar_page_sector_btn",
        use_container_width=True,
        type="primary" if st.session_state["top_page_selector"] == "Sector" else "secondary",
    ):
        st.session_state["top_page_selector"] = "Sector"
        st.rerun()
    if st.button(
        PAGE_LABELS["Yield"],
        key="sidebar_page_yield_btn",
        use_container_width=True,
        type="primary" if st.session_state["top_page_selector"] == "Yield" else "secondary",
    ):
        st.session_state["top_page_selector"] = "Yield"
        st.rerun()

st.markdown('<div class="top-header-grid">', unsafe_allow_html=True)
left_col, right_col = st.columns([2.2, 1.0])

with left_col:
    st.markdown(f'<div class="top-title-wrap"><div class="dashboard-title">{APP_TITLE}</div></div>', unsafe_allow_html=True)

with right_col:
    st.markdown(logo_html, unsafe_allow_html=True)

st.markdown('</div>', unsafe_allow_html=True)

page_name = st.session_state["top_page_selector"]
if "show_diagnostics" not in st.session_state:
    st.session_state["show_diagnostics"] = False


# =====================================
# CONTENT
# =====================================
if page_name == "Dashboard":
    display_mode = st.session_state.get("display_mode_toolbar", "Absolute")
    return_basis = st.session_state.get("return_basis_toolbar", "Nominal")
    relative_detail_mode = st.session_state.get("relative_basis_toolbar", "Major")

    is_relative_mode = display_mode == "Relative"
    is_real_mode = return_basis == "Real"
    effective_real_mode = is_real_mode and inflation_levels is not None and not inflation_levels.dropna().empty

    end_date_dashboard = get_dashboard_end_date(
        stitched_series_map=stitched_series_map,
        live_diag=live_diag,
        inflation_levels=inflation_levels,
        is_real_mode=effective_real_mode,
    )
    dashboard_saved_end = pd.Timestamp(
        st.session_state.get("dashboard_end_date_filter", end_date_dashboard).strftime("%Y-%m-%d")
    )
    dashboard_saved_end = min(dashboard_saved_end, end_date_dashboard)

    toolbar_wrap_cols = st.columns([4.15, 2.45])
    with toolbar_wrap_cols[0]:
        toolbar_cols = st.columns([0.91, 0.91, 0.91, 0.60])

        with toolbar_cols[0]:
            st.markdown('<div class="toolbar-label">Display mode:</div>', unsafe_allow_html=True)
            display_mode = st.segmented_control(
                label="Display mode",
                options=["Absolute", "Relative"],
                default=st.session_state.get("display_mode_toolbar", "Absolute"),
                key="display_mode_toolbar",
                label_visibility="collapsed",
            ) or "Absolute"

        is_relative_mode = display_mode == "Relative"

        with toolbar_cols[1]:
            st.markdown('<div class="toolbar-label">Relative basis:</div>', unsafe_allow_html=True)
            relative_detail_mode = st.segmented_control(
                label="Relative basis",
                options=["Major", "Minor"],
                default=st.session_state.get("relative_basis_toolbar", "Major"),
                key="relative_basis_toolbar",
                label_visibility="collapsed",
            ) or "Major"

        with toolbar_cols[2]:
            st.markdown('<div class="toolbar-label">Return basis:</div>', unsafe_allow_html=True)
            return_basis = st.segmented_control(
                label="Return basis",
                options=["Nominal", "Real"],
                default=st.session_state.get("return_basis_toolbar", "Nominal"),
                key="return_basis_toolbar",
                label_visibility="collapsed",
            ) or "Nominal"

        is_real_mode = return_basis == "Real"
        effective_real_mode = is_real_mode and inflation_levels is not None and not inflation_levels.dropna().empty
        end_date_dashboard = get_dashboard_end_date(
            stitched_series_map=stitched_series_map,
            live_diag=live_diag,
            inflation_levels=inflation_levels,
            is_real_mode=effective_real_mode,
        )
        dashboard_saved_end = min(dashboard_saved_end, end_date_dashboard)

        with toolbar_cols[3]:
            st.markdown('<div class="toolbar-label">End date:</div>', unsafe_allow_html=True)
            dashboard_end_input = st.date_input(
                "End date",
                value=dashboard_saved_end.date(),
                min_value=whole_period_start.date(),
                max_value=end_date_dashboard.date(),
                key="dashboard_end_date_filter",
                label_visibility="collapsed",
                format="DD/MM/YYYY",
            )

    with toolbar_wrap_cols[1]:
        right_cols = st.columns([4.0, 1.0])
        with right_cols[1]:
            st.markdown('<div class="toolbar-label">Report builder:</div>', unsafe_allow_html=True)
            export_report_placeholder = st.empty()

    dashboard_end_date_selected = pd.Timestamp(dashboard_end_input)
    dashboard_series_window = filter_series_map_to_window(
        stitched_series_map,
        None,
        dashboard_end_date_selected,
    )
    dashboard_inflation_window = (
        filter_series_to_window(inflation_levels, None, dashboard_end_date_selected)
        if inflation_levels is not None
        else None
    )

    if is_real_mode and not effective_real_mode:
        st.warning("Real mode selected but no usable UK inflation series was found. Falling back to nominal results.")

    display_groups = get_display_groups(is_relative_mode, relative_detail_mode)

    absolute_returns_df = order_asset_rows(
        calc_horizon_returns_from_levels(dashboard_series_window, dashboard_end_date_selected, list(DASHBOARD_HORIZONS.keys()))
    )

    if is_relative_mode:
        nominal_display_returns_df = order_asset_rows(
            convert_to_relative_returns(absolute_returns_df, relative_detail_mode=relative_detail_mode)
        )
    else:
        nominal_display_returns_df = absolute_returns_df.copy()

    inflation_returns_dashboard_df = (
        calc_horizon_returns_from_levels({"UK inflation": dashboard_inflation_window}, dashboard_end_date_selected, list(DASHBOARD_HORIZONS.keys()))
        if dashboard_inflation_window is not None and not dashboard_inflation_window.dropna().empty
        else pd.DataFrame()
    )

    displayed_returns_dashboard_df = (
        order_asset_rows(convert_to_real_returns(nominal_display_returns_df, inflation_returns_dashboard_df))
        if effective_real_mode
        else nominal_display_returns_df.copy()
    )

    report_nominal_abs_df = order_asset_rows(
        calc_horizon_returns_from_levels(dashboard_series_window, dashboard_end_date_selected, list(DASHBOARD_HORIZONS.keys()))
    )
    report_real_abs_df = (
        order_asset_rows(convert_to_real_returns(report_nominal_abs_df, inflation_returns_dashboard_df))
        if not inflation_returns_dashboard_df.empty
        else report_nominal_abs_df.copy()
    )
    report_relative_minor_df = order_asset_rows(
        convert_to_relative_returns(report_nominal_abs_df, relative_detail_mode="Minor")
    )
    report_nominal_returns_df = order_asset_rows(
        merge_return_tables(
            calc_horizon_returns_from_levels(
                dashboard_series_window,
                dashboard_end_date_selected,
                ["YTD", "1Y", "3Y", "5Y", "7Y", "10Y", "15Y", "20Y", "25Y"],
            ),
            calc_whole_period_returns(dashboard_series_window, dashboard_end_date_selected, whole_period_start),
        )
    )
    report_real_returns_df = (
        order_asset_rows(
            convert_to_real_returns(
                report_nominal_returns_df,
                order_asset_rows(
                    merge_return_tables(
                        calc_horizon_returns_from_levels(
                            {"UK inflation": dashboard_inflation_window},
                            dashboard_end_date_selected,
                            ["YTD", "1Y", "3Y", "5Y", "7Y", "10Y", "15Y", "20Y", "25Y"],
                        ),
                        calc_whole_period_returns({"UK inflation": dashboard_inflation_window}, dashboard_end_date_selected, whole_period_start),
                    )
                ),
            )
        )
        if dashboard_inflation_window is not None and not dashboard_inflation_window.dropna().empty
        else report_nominal_returns_df.copy()
    )
    report_growth_df = build_growth_of_wealth_df(
        chart_series_map=filter_series_map_to_window(chart_series_map, None, dashboard_end_date_selected),
        selected_assets=default_chart_assets,
        end_date=dashboard_end_date_selected,
        period_key="10Y",
        is_real_mode=False,
        inflation_levels=None,
    )
    report_yield_df = build_yield_curve_display_df(yield_curve_df, ["Nominal", "Real", "Breakeven inflation"])
    snapshot_assets = get_assets_from_display_groups(DISPLAY_GROUPS_ABSOLUTE)
    report_index_mapping = (
        mapping.drop_duplicates(subset=["asset_class"], keep="last")
        .loc[lambda df: df["asset_class"].isin(snapshot_assets), ["asset_class", "index_name"]]
        .copy()
    )
    report_index_notes = [
        f"{display_name(row.asset_class)}: {row.index_name}"
        for row in report_index_mapping.itertuples()
        if str(row.index_name).strip()
    ]
    methodology_notes = [
        f"Report end date: {dashboard_end_date_selected.strftime('%d/%m/%Y')}.",
        "Snapshot slides use absolute nominal, absolute real, and relative minor annualised returns across 20Y, 10Y, 5Y and YTD.",
        "Real returns are calculated using UK inflation via (1+asset return)/(1+inflation return)-1.",
        "Growth of wealth uses nominal GBP series over 10 years and prioritises live daily history where available before stitched monthly index history.",
        "UK nominal and real spot curves are sourced from the Bank of England. Short-end real yields may be extended using DividendData where available.",
    ]
    further_notes = [
        "The app diagnostics section provides the underlying mapping, live ticker coverage, yield-curve source checks and inflation notes.",
        "Relative minor mode shows growth assets relative to their configured local benchmark and defensive assets relative to cash.",
        "More information on the Albion indices is available at smartersuccess.net/indices.",
    ]
    report_bytes = build_quarterly_market_metrics_report(
        report_end_date_text=dashboard_end_date_selected.strftime("%d/%m/%Y"),
        report_end_date_long_text=dashboard_end_date_selected.strftime("%d %B %Y"),
        report_month_text=dashboard_end_date_selected.strftime("%B %Y"),
        nominal_abs_df=report_nominal_abs_df[["asset_class", "20Y", "10Y", "5Y", "YTD"]],
        real_abs_df=report_real_abs_df[["asset_class", "20Y", "10Y", "5Y", "YTD"]],
        relative_minor_df=report_relative_minor_df[["asset_class", "20Y", "10Y", "5Y", "YTD"]],
        returns_nominal_df=report_nominal_returns_df[["asset_class"] + RETURNS_TABLE_PERIODS],
        returns_real_df=report_real_returns_df[["asset_class"] + RETURNS_TABLE_PERIODS],
        growth_df=report_growth_df,
        yield_df=report_yield_df,
        bullet_points=build_market_commentary_bullets(report_nominal_abs_df),
        methodology_notes=methodology_notes,
        index_notes=report_index_notes,
        further_notes=further_notes,
        snapshot_assets=snapshot_assets,
    )
    export_report_placeholder.download_button(
        "Export",
        data=report_bytes,
        file_name=f"quarterly_market_metrics_{dashboard_end_date_selected.strftime('%Y_%m')}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="export_quarterly_market_metrics_report",
        use_container_width=True,
    )

    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Annualised {"real" if effective_real_mode else "nominal"} returns in GBP to <b>{dashboard_end_date_selected.strftime("%d/%m/%Y")}</b>
        </div>
        """,
        unsafe_allow_html=True,
    )

    lookup = build_lookup_table(displayed_returns_dashboard_df)

    cols = st.columns(4)
    period_order = ["20Y", "10Y", "5Y", "YTD"]

    for col, period in zip(cols, period_order):
        period_vals = displayed_returns_dashboard_df[period].dropna()
        vmin = period_vals.min() if len(period_vals) else -0.05
        vmax = period_vals.max() if len(period_vals) else 0.15

        with col:
            st.markdown('<div class="period-shell">', unsafe_allow_html=True)
            st.markdown(f'<div class="period-title">{DASHBOARD_HORIZONS[period]}</div>', unsafe_allow_html=True)

            for group in display_groups:
                title = group["title"]
                items = group["items"]
                labels = group["labels"]

                st.markdown('<div class="group-card">', unsafe_allow_html=True)

                if title:
                    st.markdown(f'<div class="section-title">{title}</div>', unsafe_allow_html=True)
                else:
                    st.markdown('<div class="section-title-empty"></div>', unsafe_allow_html=True)

                if len(items) == 1:
                    item = items[0]
                    val = lookup.get(item, {}).get(period, np.nan)
                    colour = heat_colour(val, vmin, vmax)
                    label = labels.get(item, "")
                    subtitle_html = f'<div class="section-subtitle">{label}</div>' if label else ""

                    st.markdown(
                        f"""
                        {subtitle_html}
                        <div class="big-tile" style="background:{colour}; width:50%; margin-left:auto; margin-right:auto;">
                            {format_pct(val)}
                        </div>
                        """,
                        unsafe_allow_html=True,
                    )

                elif len(items) == 2:
                    two_cols = st.columns(2)
                    for idx, item in enumerate(items):
                        val = lookup.get(item, {}).get(period, np.nan)
                        colour = heat_colour(val, vmin, vmax)
                        label = labels.get(item, item)
                        with two_cols[idx]:
                            st.markdown(
                                f"""
                                <div class="small-tile" style="background:{colour}; min-height:74px;">
                                    <span class="tile-label tile-label-plain">{label}</span>
                                    {format_pct(val)}
                                </div>
                                """,
                                unsafe_allow_html=True,
                            )

                elif len(items) == 3:
                    broad = items[0]
                    broad_val = lookup.get(broad, {}).get(period, np.nan)
                    broad_colour = heat_colour(broad_val, vmin, vmax)
                    broad_label = labels.get(broad, "")
                    broad_label_html = f'<span class="tile-label-on-colour">{broad_label}</span>' if broad_label else ""

                    st.markdown(
                        f"""
                        <div class="big-tile" style="background:{broad_colour};">
                            {broad_label_html}
                            {format_pct(broad_val)}
                        </div>
                        """,
                        unsafe_allow_html=True,
                    )

                    small_cols = st.columns(2)
                    for i, item in enumerate(items[1:]):
                        val = lookup.get(item, {}).get(period, np.nan)
                        colour = heat_colour(val, vmin, vmax)
                        label = labels.get(item, "")
                        with small_cols[i]:
                            st.markdown(
                                f"""
                                <div class="small-tile" style="background:{colour};">
                                    <span class="tile-label">{label}</span>
                                    {format_pct(val)}
                                </div>
                                """,
                                unsafe_allow_html=True,
                            )

                elif len(items) == 4:
                    row1 = st.columns(2)
                    row2 = st.columns(2)
                    for idx, item in enumerate(items):
                        val = lookup.get(item, {}).get(period, np.nan)
                        colour = heat_colour(val, vmin, vmax)
                        label = labels.get(item, item)
                        target_cols = row1 if idx < 2 else row2
                        with target_cols[idx % 2]:
                            st.markdown(
                                f"""
                                <div class="small-tile" style="background:{colour}; min-height:74px;">
                                    <span class="tile-label tile-label-plain">{label}</span>
                                    {format_pct(val)}
                                </div>
                                """,
                                unsafe_allow_html=True,
                            )
                
                elif len(items) == 5:
                    row1 = st.columns(3)
                    row2 = st.columns(2)

                    for idx, item in enumerate(items):
                        val = lookup.get(item, {}).get(period, np.nan)
                        colour = heat_colour(val, vmin, vmax)
                        label = labels.get(item, item)

                        if idx < 3:
                            target_cols = row1
                            target_idx = idx
                        else:
                            target_cols = row2
                            target_idx = idx - 3

                        with target_cols[target_idx]:
                            st.markdown(
                                f"""
                                <div class="small-tile" style="background:{colour}; min-height:74px;">
                                    <span class="tile-label tile-label-plain">{label}</span>
                                    {format_pct(val)}
                                </div>
                                """,
                                unsafe_allow_html=True,
                            )

                st.markdown("</div>", unsafe_allow_html=True)
                st.markdown('<div class="spacer"></div>', unsafe_allow_html=True)

            st.markdown("</div>", unsafe_allow_html=True)

    st.markdown(
        f'<div class="methodology-text">{get_methodology_paragraph("Dashboard", is_relative_mode, relative_detail_mode, effective_real_mode, inflation_source_note)}</div>',
        unsafe_allow_html=True,
    )

elif page_name == "Charts":
    detail_return_basis = st.session_state.get("detail_return_basis_toolbar", "Nominal")
    detail_period = st.session_state.get("detail_period_toolbar", "YTD")

    is_real_mode = detail_return_basis == "Real"
    effective_real_mode = is_real_mode and inflation_levels is not None and not inflation_levels.dropna().empty

    end_date_charts = get_dashboard_end_date(
        stitched_series_map=stitched_series_map,
        live_diag=live_diag,
        inflation_levels=inflation_levels,
        is_real_mode=effective_real_mode,
    )
    charts_min_date, _ = get_series_map_date_bounds(chart_series_map)
    charts_default_start = whole_period_start if charts_min_date is None else max(whole_period_start, charts_min_date)
    charts_saved_start = pd.Timestamp(
        st.session_state.get("charts_start_date_filter", charts_default_start).strftime("%Y-%m-%d")
    )
    charts_saved_end = pd.Timestamp(
        st.session_state.get("charts_end_date_filter", end_date_charts).strftime("%Y-%m-%d")
    )
    charts_saved_start = max(charts_default_start, min(charts_saved_start, end_date_charts))
    charts_saved_end = max(charts_saved_start, min(charts_saved_end, end_date_charts))
    chart_period_options = list(CHART_PERIODS.keys()) + ["Custom"]

    toolbar_wrap_cols = st.columns([5.35, 1.25])
    with toolbar_wrap_cols[0]:
        toolbar_cols = st.columns([0.92, 1.95, 0.66, 0.66])

        with toolbar_cols[0]:
            st.markdown('<div class="toolbar-label">Return basis:</div>', unsafe_allow_html=True)
            detail_return_basis = st.segmented_control(
                label="Return basis",
                options=["Nominal", "Real"],
                default=st.session_state.get("detail_return_basis_toolbar", "Nominal"),
                key="detail_return_basis_toolbar",
                label_visibility="collapsed",
            ) or "Nominal"

    is_real_mode = detail_return_basis == "Real"
    effective_real_mode = is_real_mode and inflation_levels is not None and not inflation_levels.dropna().empty
    end_date_charts = get_dashboard_end_date(
        stitched_series_map=stitched_series_map,
        live_diag=live_diag,
        inflation_levels=inflation_levels,
        is_real_mode=effective_real_mode,
    )
    charts_saved_start = max(charts_default_start, min(charts_saved_start, end_date_charts))
    charts_saved_end = max(charts_saved_start, min(charts_saved_end, end_date_charts))
    detail_period = detail_period if detail_period in chart_period_options else "YTD"

    prev_chart_period = st.session_state.get("_charts_prev_period")
    prev_chart_start = st.session_state.get("_charts_prev_start_date")
    prev_chart_end = st.session_state.get("_charts_prev_end_date")

    charts_active_start = charts_saved_start.normalize()
    charts_active_end = charts_saved_end.normalize()

    start_changed = prev_chart_start is not None and charts_active_start != pd.Timestamp(prev_chart_start).normalize()
    end_changed = prev_chart_end is not None and charts_active_end != pd.Timestamp(prev_chart_end).normalize()
    period_changed = prev_chart_period is not None and detail_period != prev_chart_period

    if prev_chart_period is None:
        if detail_period == "Custom":
            detail_period = match_chart_period_from_dates(charts_active_start, charts_active_end, charts_default_start)
        if detail_period != "Custom":
            charts_active_start = get_chart_period_start_date(detail_period, charts_active_end, charts_default_start)
    elif period_changed and detail_period != "Custom":
        charts_active_start = get_chart_period_start_date(detail_period, charts_active_end, charts_default_start)
    elif end_changed and not start_changed and detail_period != "Custom":
        charts_active_start = get_chart_period_start_date(detail_period, charts_active_end, charts_default_start)
    elif start_changed or end_changed:
        detail_period = match_chart_period_from_dates(charts_active_start, charts_active_end, charts_default_start)

    st.session_state["detail_period_toolbar"] = detail_period
    st.session_state["charts_start_date_filter"] = charts_active_start.date()
    st.session_state["charts_end_date_filter"] = charts_active_end.date()

    with toolbar_cols[1]:
        st.markdown('<div class="toolbar-label">Chart period:</div>', unsafe_allow_html=True)
        detail_period = st.segmented_control(
            label="Chart period",
            options=chart_period_options,
            key="detail_period_toolbar",
            label_visibility="collapsed",
        ) or "YTD"

    with toolbar_cols[2]:
        st.markdown('<div class="toolbar-label">Start date:</div>', unsafe_allow_html=True)
        charts_start_input = st.date_input(
            "Start date",
            min_value=charts_default_start.date(),
            max_value=end_date_charts.date(),
            key="charts_start_date_filter",
            label_visibility="collapsed",
            format="DD/MM/YYYY",
            disabled=detail_period != "Custom",
        )

    charts_start_date = pd.Timestamp(charts_start_input)

    with toolbar_cols[3]:
        st.markdown('<div class="toolbar-label">End date:</div>', unsafe_allow_html=True)
        charts_end_input = st.date_input(
            "End date",
            min_value=charts_start_date.date(),
            max_value=end_date_charts.date(),
            key="charts_end_date_filter",
            label_visibility="collapsed",
            format="DD/MM/YYYY",
            disabled=detail_period != "Custom",
        )

    charts_end_date_selected = pd.Timestamp(charts_end_input)
    st.session_state["_charts_prev_period"] = detail_period
    st.session_state["_charts_prev_start_date"] = charts_start_date.date()
    st.session_state["_charts_prev_end_date"] = charts_end_date_selected.date()
    chart_series_window = filter_series_map_to_window(
        chart_series_map,
        charts_start_date,
        charts_end_date_selected,
    )
    stitched_series_window = filter_series_map_to_window(
        stitched_series_map,
        charts_start_date,
        charts_end_date_selected,
    )
    charts_inflation_window = (
        filter_series_to_window(inflation_levels, charts_start_date, charts_end_date_selected)
        if inflation_levels is not None
        else None
    )
    charts_whole_period_start = charts_start_date

    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Return in GBP from <b>{charts_whole_period_start.strftime("%d/%m/%Y")}</b> to <b>{charts_end_date_selected.strftime("%d/%m/%Y")}</b>
        </div>
        """,
        unsafe_allow_html=True,
    )

    if is_real_mode and not effective_real_mode:
        st.warning("Real mode selected but no usable UK inflation series was found. Falling back to nominal results.")

    saved_default_assets = st.session_state.get("detail_selected_assets", default_chart_assets)
    saved_default_assets = [a for a in saved_default_assets if a in available_assets] or default_chart_assets

    selected_assets = st.multiselect(
        "Asset classes",
        options=available_assets,
        default=saved_default_assets,
        key="detail_selected_assets",
    )

    if not selected_assets:
        st.info("Select at least one asset class to populate the chart and tables.")
        selected_assets = []

    nominal_returns_charts_df = order_asset_rows(
        merge_return_tables(
            calc_horizon_returns_from_levels(
                stitched_series_map,
                end_date_charts,
                ["YTD", "1Y", "3Y", "5Y", "7Y", "10Y", "15Y", "20Y", "25Y"],
            ),
            calc_whole_period_returns(stitched_series_map, end_date_charts, whole_period_start),
        )
    )

    inflation_returns_charts_df = (
        order_asset_rows(
            merge_return_tables(
                calc_horizon_returns_from_levels(
                    {"UK inflation": inflation_levels},
                    end_date_charts,
                    ["YTD", "1Y", "3Y", "5Y", "7Y", "10Y", "15Y", "20Y", "25Y"],
                ),
                calc_whole_period_returns({"UK inflation": inflation_levels}, end_date_charts, whole_period_start),
            )
        )
        if inflation_levels is not None and not inflation_levels.dropna().empty
        else pd.DataFrame()
    )

    displayed_returns_charts_df = (
        order_asset_rows(convert_to_real_returns(nominal_returns_charts_df, inflation_returns_charts_df))
        if effective_real_mode
        else nominal_returns_charts_df.copy()
    )

    calendar_year_df = build_calendar_year_returns(stitched_series_map, end_date_charts, years_back=10)

    if effective_real_mode and inflation_levels is not None and not inflation_levels.dropna().empty:
        inflation_calendar_year_df = build_calendar_year_returns({"UK inflation": inflation_levels}, end_date_charts, years_back=10)
        infl_lookup = inflation_calendar_year_df.set_index("asset_class").to_dict(orient="index").get("UK inflation", {})
        calendar_year_real_df = calendar_year_df.copy()
        for year_col in [c for c in calendar_year_real_df.columns if c != "asset_class"]:
            infl_val = infl_lookup.get(year_col, np.nan)
            calendar_year_real_df[year_col] = calendar_year_real_df[year_col].map(lambda x: safe_relative_return(x, infl_val))
        calendar_year_df = calendar_year_real_df

    growth_df = build_growth_of_wealth_df(
        chart_series_map=chart_series_window,
        selected_assets=selected_assets,
        end_date=charts_end_date_selected,
        period_key=detail_period,
        is_real_mode=effective_real_mode,
        inflation_levels=charts_inflation_window,
    )

    st.markdown('<div class="page-section-title">Growth of wealth</div>', unsafe_allow_html=True)
    if growth_df.empty:
        st.info("No chart data available for the current selection.")
    else:
        st.altair_chart(build_chart(growth_df, selected_assets, detail_period), width="stretch")

    st.markdown('<div class="page-section-title">Annualised returns</div>', unsafe_allow_html=True)
    returns_display_df = (
        displayed_returns_charts_df[displayed_returns_charts_df["asset_class"].isin(selected_assets)][["asset_class"] + RETURNS_TABLE_PERIODS]
        if selected_assets else displayed_returns_charts_df[["asset_class"] + RETURNS_TABLE_PERIODS]
    )
    st.markdown(
        build_html_table(
            returns_display_df,
            percent_cols=RETURNS_TABLE_PERIODS,
            conditional_cols=RETURNS_TABLE_PERIODS,
        ),
        unsafe_allow_html=True,
    )

    st.markdown('<div class="page-section-title">Calendar year returns</div>', unsafe_allow_html=True)
    calendar_display_df = (
        calendar_year_df[calendar_year_df["asset_class"].isin(selected_assets)]
        if selected_assets else calendar_year_df
    )
    year_cols = [c for c in calendar_display_df.columns if c != "asset_class"]
    st.markdown(
        build_html_table(
            calendar_display_df[["asset_class"] + year_cols],
            percent_cols=year_cols,
            conditional_cols=year_cols,
        ),
        unsafe_allow_html=True,
    )

    st.markdown(
        f'<div class="methodology-text">{get_methodology_paragraph("Charts", False, "Major", effective_real_mode, inflation_source_note)}</div>',
        unsafe_allow_html=True,
    )

elif page_name == "Risk":
    risk_return_basis = st.session_state.get("risk_return_basis_toolbar", "Nominal")
    risk_period = st.session_state.get("risk_period_toolbar", "10Y")

    is_real_mode = risk_return_basis == "Real"
    effective_real_mode = is_real_mode and inflation_levels is not None and not inflation_levels.dropna().empty

    end_date_risk = get_dashboard_end_date(
        stitched_series_map=stitched_series_map,
        live_diag=live_diag,
        inflation_levels=inflation_levels,
        is_real_mode=effective_real_mode,
    )

    toolbar_wrap_cols = st.columns([4.9, 1.7])
    with toolbar_wrap_cols[0]:
        toolbar_cols = st.columns([0.92, 2.55])

        with toolbar_cols[0]:
            st.markdown('<div class="toolbar-label">Return basis:</div>', unsafe_allow_html=True)
            risk_return_basis = st.segmented_control(
                label="Return basis",
                options=["Nominal", "Real"],
                default=st.session_state.get("risk_return_basis_toolbar", "Nominal"),
                key="risk_return_basis_toolbar",
                label_visibility="collapsed",
            ) or "Nominal"

        is_real_mode = risk_return_basis == "Real"
        effective_real_mode = is_real_mode and inflation_levels is not None and not inflation_levels.dropna().empty

        with toolbar_cols[1]:
            st.markdown('<div class="toolbar-label">Risk period:</div>', unsafe_allow_html=True)
            risk_period = st.segmented_control(
                label="Risk period",
                options=list(RISK_PERIODS.keys()),
                default=st.session_state.get("risk_period_toolbar", "10Y"),
                key="risk_period_toolbar",
                label_visibility="collapsed",
            ) or "10Y"

    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Annualised risk / return in GBP to <b>{end_date_risk.strftime("%d/%m/%Y")}</b>
        </div>
        """,
        unsafe_allow_html=True,
    )

    if is_real_mode and not effective_real_mode:
        st.warning("Real mode selected but no usable UK inflation series was found. Falling back to nominal results.")

    saved_risk_assets = st.session_state.get("risk_selected_assets", default_chart_assets)
    saved_risk_assets = [a for a in saved_risk_assets if a in available_assets] or default_chart_assets

    selected_assets = st.multiselect(
        "Asset classes",
        options=available_assets,
        default=saved_risk_assets,
        key="risk_selected_assets",
    )

    if not selected_assets:
        st.info("Select at least one asset class to populate the chart and table.")
        selected_assets = []

    risk_series_map = {}
    for asset_class, series in stitched_series_map.items():
        risk_series = series
        if effective_real_mode:
            if inflation_levels is None or inflation_levels.dropna().empty:
                continue
            risk_series = build_real_level_series(series, inflation_levels)
        risk_series_map[asset_class] = risk_series

    risk_scatter_df = build_risk_scatter_df(
        series_map=risk_series_map,
        selected_assets=selected_assets,
        end_date=end_date_risk,
        period_key=st.session_state.get("risk_period_toolbar", risk_period),
        whole_period_start=whole_period_start,
    )
    risk_table_df = build_risk_metrics_table(
        series_map=risk_series_map,
        selected_assets=selected_assets,
        end_date=end_date_risk,
        period_keys=list(RISK_PERIODS.keys()),
        whole_period_start=whole_period_start,
    )
    risk_summary_df = build_risk_summary_table(
        series_map=risk_series_map,
        asset_style_map=asset_style_map,
        selected_assets=selected_assets,
        end_date=end_date_risk,
        start_date=whole_period_start,
    )
    correlation_matrix_df = build_correlation_matrix_table(
        series_map=risk_series_map,
        selected_assets=selected_assets,
        end_date=end_date_risk,
        start_date=whole_period_start,
    )

    st.markdown('<div class="page-section-title">Volatility/return chart</div>', unsafe_allow_html=True)
    if risk_scatter_df.empty:
        st.info("No risk data available for the current selection.")
    else:
        st.altair_chart(build_risk_scatter_chart(risk_scatter_df, selected_assets), width="stretch")

    st.markdown('<div class="page-section-title">Volatility/return table</div>', unsafe_allow_html=True)
    risk_percent_cols = [c for c in risk_table_df.columns if c != "asset_class"]
    st.markdown(
        build_html_table(
            risk_table_df,
            percent_cols=risk_percent_cols,
            conditional_cols=risk_percent_cols,
            header_wrap_cols=risk_percent_cols,
            invert_conditional_cols=[c for c in risk_percent_cols if c.endswith("Vol")],
            rank_conditional_cols=[c for c in risk_percent_cols if c.endswith("Vol")],
        ),
        unsafe_allow_html=True,
    )

    st.markdown('<div class="page-section-title">Since inception risk metrics</div>', unsafe_allow_html=True)
    risk_summary_display_df = risk_summary_df.copy()
    ratio_cols = [
        "Return/vol ratio",
        "Sharpe ratio",
        "Information ratio",
        "Sortino ratio",
        "Calmar ratio",
    ]
    for col in ratio_cols:
        if col in risk_summary_display_df.columns:
            risk_summary_display_df[col] = risk_summary_display_df[col].map(
                lambda x: np.nan if pd.isna(x) else round(float(x), 2)
            )

    risk_summary_percent_cols = [
        c for c in ["Period return", "Period vol", "Worst drawdown", "Tracking error"] if c in risk_summary_display_df.columns
    ]
    risk_summary_conditional_cols = [c for c in risk_summary_display_df.columns if c != "asset_class"]
    st.markdown(
        build_html_table(
            risk_summary_display_df,
            percent_cols=risk_summary_percent_cols,
            conditional_cols=risk_summary_conditional_cols,
            header_wrap_cols=risk_summary_conditional_cols,
            invert_conditional_cols=["Period vol", "Tracking error"],
            rank_conditional_cols=["Period vol", "Worst drawdown", "Tracking error"],
            decimal_cols=ratio_cols,
        ),
        unsafe_allow_html=True,
    )

    st.markdown('<div class="page-section-title">Since inception correlation matrix</div>', unsafe_allow_html=True)
    correlation_cols = [c for c in correlation_matrix_df.columns if c != "asset_class"]
    st.markdown(
        build_html_table(
            correlation_matrix_df,
            conditional_cols=correlation_cols,
            header_wrap_cols=correlation_cols,
            decimal_cols=correlation_cols,
            correlation_conditional_cols=correlation_cols,
        ),
        unsafe_allow_html=True,
    )

    risk_methodology = (
        f"This tab shows annualised {'real' if effective_real_mode else 'nominal'} return and annualised volatility in GBP "
        f"to <b>{end_date_risk.strftime('%d/%m/%Y')}</b>. Volatility is calculated from monthly returns and annualised using the square root of 12."
    )
    if effective_real_mode:
        risk_methodology += f" Current inflation source: {inflation_source_note}."
    risk_methodology += (
        " Glossary: Return/vol ratio is annualised return divided by annualised volatility. "
        "Sharpe ratio is excess return over Cash (GBP) divided by volatility. "
        "Information ratio is excess return versus the assigned benchmark divided by tracking error. "
        "Sortino ratio is return divided by downside deviation. "
        "Worst drawdown is the largest peak-to-trough fall since common inception. "
        "Calmar ratio is annualised return divided by the absolute worst drawdown. "
        "Tracking error is the annualised standard deviation of monthly excess returns versus the assigned benchmark. "
        "The assigned benchmark is Global market for growth assets and Cash (GBP) for defensive assets."
    )
    st.markdown(f'<div class="methodology-text">{risk_methodology}</div>', unsafe_allow_html=True)

elif page_name == "Factors":
    factors_period = st.session_state.get("factors_period_toolbar", "YTD")
    factors_currency = st.session_state.get("factors_currency_toolbar", "USD")
    fx_values_df = fetch_fx_value_series()
    _, factor_fx_end_date = build_currency_performance_matrix(fx_values_df, factors_period)
    factor_period_options = get_geo_period_options(factor_fx_end_date, fx_values_df, factors_currency)
    if factors_period not in factor_period_options:
        factors_period = "MAX" if "MAX" in factor_period_options else factor_period_options[0]
        st.session_state["factors_period_toolbar"] = factors_period

    toolbar_wrap_cols = st.columns([4.9, 1.5])
    with toolbar_wrap_cols[0]:
        toolbar_cols = st.columns([1.38, 0.28, 1.14])

        with toolbar_cols[0]:
            st.markdown('<div class="toolbar-label">Period:</div>', unsafe_allow_html=True)
            factors_period = st.segmented_control(
                label="Period",
                options=factor_period_options,
                default=st.session_state.get("factors_period_toolbar", factors_period),
                key="factors_period_toolbar",
                label_visibility="collapsed",
            ) or "YTD"

        with toolbar_cols[1]:
            st.markdown('<div class="toolbar-label">Currency:</div>', unsafe_allow_html=True)
            factors_currency = st.selectbox(
                "Currency",
                options=GEO_NEUTRAL_CURRENCIES,
                index=GEO_NEUTRAL_CURRENCIES.index(st.session_state.get("factors_currency_toolbar", "USD")),
                key="factors_currency_toolbar",
                label_visibility="collapsed",
            )

    factor_style_box_df, factor_start_anchor, factor_end_date, factor_common_inception = build_factor_style_box_df(
        factors_df=factors_df,
        fx_values_df=fx_values_df,
        period_key=factors_period,
        neutral_currency=factors_currency,
        region_key="US",
    )

    st.markdown('<div class="page-section-title">US equity style box</div>', unsafe_allow_html=True)
    factor_start_text = pd.Timestamp(factor_start_anchor).strftime("%d/%m/%Y") if factor_start_anchor is not None else "-"
    factor_end_text = pd.Timestamp(factor_end_date).strftime("%d/%m/%Y") if factor_end_date is not None else "-"
    common_inception_text = (
        pd.Timestamp(factor_common_inception).strftime("%d/%m/%Y") if factor_common_inception is not None else "-"
    )
    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            US size and style factor returns in {factors_currency} from <b>{factor_start_text}</b> to <b>{factor_end_text}</b>. The <b>MAX</b> option uses the common inception across the included US funds, currently <b>{common_inception_text}</b>.
        </div>
        """,
        unsafe_allow_html=True,
    )

    if factor_style_box_df.empty:
        st.info("No US factor data were available from the factors sheet tickers for the selected period.")
    else:
        style_box_cols = st.columns([0.9, 1.2, 0.9])
        with style_box_cols[1]:
            st.markdown(build_factor_style_box_html(factor_style_box_df), unsafe_allow_html=True)

    st.markdown(
        '<div class="methodology-text">This tab uses the workbook <b>factors</b> sheet and currently filters it to <b>Region = US</b>. '
        'Each tile shows the selected-period return for the mapped ETF or fund in that size/style bucket, converted into the chosen currency where quote-currency data are available. '
        'Tile colours are relative within the 3x3 box, so stronger buckets are greener and weaker buckets are redder.</div>',
        unsafe_allow_html=True,
    )

elif page_name == "Geo":
    geo_period = st.session_state.get("geo_period_toolbar", "YTD")
    geo_neutral_currency = st.session_state.get("geo_neutral_currency_toolbar", "USD")
    fx_values_df = fetch_fx_value_series()
    _, currency_end_date = build_currency_performance_matrix(fx_values_df, geo_period)
    geo_period_options = get_geo_period_options(currency_end_date, fx_values_df, geo_neutral_currency)
    if geo_period not in geo_period_options:
        geo_period = "MAX" if "MAX" in geo_period_options else geo_period_options[0]
        st.session_state["geo_period_toolbar"] = geo_period
    country_performance_df, country_end_date = build_country_performance_df(
        regions_df=regions_df,
        fx_values_df=fx_values_df,
        period_key=geo_period,
        neutral_currency=geo_neutral_currency,
        preferred_series_map=chart_series_map,
    )
    country_series_map, country_series_end_date = build_country_series_map(
        regions_df=regions_df,
        fx_values_df=fx_values_df,
        neutral_currency=geo_neutral_currency,
        preferred_series_map=chart_series_map,
    )
    region_series_map, region_series_end_date = build_region_series_map(
        regions_df=regions_df,
        fx_values_df=fx_values_df,
        neutral_currency=geo_neutral_currency,
    )
    region_performance_df, region_end_date = build_region_performance_df(
        regions_df=regions_df,
        fx_values_df=fx_values_df,
        period_key=geo_period,
        neutral_currency=geo_neutral_currency,
        preferred_series_map=chart_series_map,
    )

    toolbar_wrap_cols = st.columns([4.9, 1.5])
    with toolbar_wrap_cols[0]:
        toolbar_cols = st.columns([1.38, 0.28, 1.14])

        with toolbar_cols[0]:
            st.markdown('<div class="toolbar-label">Period:</div>', unsafe_allow_html=True)
            geo_period = st.segmented_control(
                label="Period",
                options=geo_period_options,
                default=st.session_state.get("geo_period_toolbar", geo_period),
                key="geo_period_toolbar",
                label_visibility="collapsed",
            ) or "YTD"

        with toolbar_cols[1]:
            st.markdown('<div class="toolbar-label">Currency:</div>', unsafe_allow_html=True)
            geo_neutral_currency = st.selectbox(
                "Currency",
                options=GEO_NEUTRAL_CURRENCIES,
                index=GEO_NEUTRAL_CURRENCIES.index(st.session_state.get("geo_neutral_currency_toolbar", "USD")),
                key="geo_neutral_currency_toolbar",
                label_visibility="collapsed",
            )

    if geo_period != st.session_state.get("geo_period_toolbar", "YTD"):
        st.session_state["geo_period_toolbar"] = geo_period

    currency_matrix_df, currency_end_date = build_currency_performance_matrix(fx_values_df, geo_period)
    country_performance_df, country_end_date = build_country_performance_df(
        regions_df=regions_df,
        fx_values_df=fx_values_df,
        period_key=geo_period,
        neutral_currency=geo_neutral_currency,
        preferred_series_map=chart_series_map,
    )
    country_series_map, country_series_end_date = build_country_series_map(
        regions_df=regions_df,
        fx_values_df=fx_values_df,
        neutral_currency=geo_neutral_currency,
        preferred_series_map=chart_series_map,
    )
    region_series_map, region_series_end_date = build_region_series_map(
        regions_df=regions_df,
        fx_values_df=fx_values_df,
        neutral_currency=geo_neutral_currency,
    )
    region_performance_df, region_end_date = build_region_performance_df(
        regions_df=regions_df,
        fx_values_df=fx_values_df,
        period_key=geo_period,
        neutral_currency=geo_neutral_currency,
        preferred_series_map=chart_series_map,
    )

    st.markdown('<div class="page-section-title">Country performance</div>', unsafe_allow_html=True)
    country_start_text = (
        get_geo_period_start_anchor(country_end_date, geo_period, fx_values_df, geo_neutral_currency).strftime("%d/%m/%Y")
        if country_end_date is not None
        else "-"
    )
    country_end_text = country_end_date.strftime("%d/%m/%Y") if country_end_date is not None else "-"
    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Country and global market performance in {geo_neutral_currency} from <b>{country_start_text}</b> to <b>{country_end_text}</b>
        </div>
        """,
        unsafe_allow_html=True,
    )

    if country_performance_df.empty:
        st.info("No country data were available from the regions sheet tickers for the selected period.")
    else:
        st.markdown(build_country_tiles_html(country_performance_df), unsafe_allow_html=True)

    st.markdown('<div class="page-section-title">Regional performance</div>', unsafe_allow_html=True)
    region_start_text = (
        get_geo_period_start_anchor(region_end_date, geo_period, fx_values_df, geo_neutral_currency).strftime("%d/%m/%Y")
        if region_end_date is not None
        else "-"
    )
    region_end_text = region_end_date.strftime("%d/%m/%Y") if region_end_date is not None else "-"
    msci_source_date_text = MSCI_ACWI_IMI_REGION_SOURCE_DATE.strftime("%d/%m/%Y")
    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Regional and global market performance in {geo_neutral_currency} from <b>{region_start_text}</b> to <b>{region_end_text}</b>. Tile widths use MSCI ACWI IMI regional weights as at <b>{msci_source_date_text}</b>.
        </div>
        """,
        unsafe_allow_html=True,
    )
    if region_performance_df.empty:
        st.info("No regional data were available from the regions sheet tickers for the selected period.")
    else:
        st.markdown(build_region_tiles_html(region_performance_df), unsafe_allow_html=True)

    st.markdown('<div class="page-section-title">Patchwork quilt</div>', unsafe_allow_html=True)
    patchwork_view = st.segmented_control(
        label="Patchwork view",
        options=list(PATCHWORK_COUNTRY_SETS.keys()),
        default=st.session_state.get("geo_patchwork_view", "Largest 10"),
        key="geo_patchwork_view",
        label_visibility="collapsed",
    ) or "Largest 10"
    country_rankings_df = (
        fetch_companiesmarketcap_country_rankings(COMPANIESMARKETCAP_COUNTRIES_URL)
        if patchwork_view in {"Largest 10", "Largest 20"}
        else pd.DataFrame()
    )
    patchwork_series_map = region_series_map if patchwork_view == "Regional" else country_series_map
    patchwork_end_date = region_series_end_date if patchwork_view == "Regional" else country_series_end_date
    patchwork_labels = resolve_patchwork_labels(
        patchwork_view,
        list(patchwork_series_map.keys()),
        country_rankings_df,
    )
    patchwork_df, patchwork_legend_df, patchwork_years = build_country_patchwork_quilt(
        country_series_map=patchwork_series_map,
        end_date=patchwork_end_date,
        years_back=10,
        include_labels=patchwork_labels,
    )
    patchwork_start_text = f"31/12/{patchwork_years[0]-1}" if patchwork_years else "-"
    patchwork_end_text = f"31/12/{patchwork_years[-1]}" if patchwork_years else "-"
    patchwork_latest_text = pd.Timestamp(patchwork_end_date).strftime("%d/%m/%Y") if patchwork_end_date is not None else "-"
    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Calendar-year {'regional' if patchwork_view == 'Regional' else 'country'} rankings in {geo_neutral_currency} from <b>{patchwork_start_text}</b> to <b>{patchwork_end_text}</b>. The <b>Whole period</b> column shows annualised returns across that full window, and the <b>YTD</b> column runs from <b>31/12/{pd.Timestamp(patchwork_end_date).year - 1 if patchwork_end_date is not None else '-'}</b> to <b>{patchwork_latest_text}</b>. Only entries with a full 10-year calendar history are included.
        </div>
        """,
        unsafe_allow_html=True,
    )
    if patchwork_df.empty:
        st.info("No entries in the selected patchwork set had a full 10-year calendar history in the selected currency.")
    else:
        st.markdown(build_country_patchwork_html(patchwork_df, patchwork_legend_df, patchwork_years), unsafe_allow_html=True)

    st.markdown(
        '<div class="methodology-text">This tab shows country performance using the tickers on the workbook <b>regions</b> sheet where <b>Country = 1</b> and <b>Available = Yes</b>, '
        'regional performance using the rows where <b>Region = 1</b>, plus a Global market comparator from the preferred chart series used on the growth chart. '
        'The regional block uses MSCI ACWI IMI regional weights sourced from '
        f'<a href="{MSCI_ACWI_IMI_REGION_SOURCE_URL}" target="_blank">MSCI</a> as at {msci_source_date_text}. '
        'Yahoo Finance prices are converted into the selected neutral currency where quote-currency data are available. '
        f'The <b>MAX</b> option uses a fixed start date anchored to <b>{GEO_MAX_START.strftime("%d/%m/%Y")}</b> or the selected FX history where later. '
        'The patchwork quilt ranks each included country or region by calendar-year return, top to bottom within each year, and shows only entries with a full 10-year history across the displayed quilt window. '
        'Using smaller curated sets such as Largest 10, Largest 20, or Regional also allows a more distinct colour palette than All countries.</div>',
        unsafe_allow_html=True,
    )

elif page_name == "Sector":
    sector_period = st.session_state.get("sector_period_toolbar", "YTD")
    sector_neutral_currency = st.session_state.get("sector_neutral_currency_toolbar", "USD")
    fx_values_df = fetch_fx_value_series()
    _, sector_fx_end_date = build_currency_performance_matrix(fx_values_df, sector_period)
    sector_period_options = get_geo_period_options(sector_fx_end_date, fx_values_df, sector_neutral_currency)
    if sector_period not in sector_period_options:
        sector_period = "MAX" if "MAX" in sector_period_options else sector_period_options[0]
        st.session_state["sector_period_toolbar"] = sector_period

    sector_series_map, sector_series_end_date = build_labelled_series_map(
        source_df=sectors_df,
        label_col="sector",
        fx_values_df=fx_values_df,
        neutral_currency=sector_neutral_currency,
        preferred_series_map=chart_series_map,
    )
    sector_labels = {str(label).strip() for label in sectors_df["sector"].astype(str).tolist()} if not sectors_df.empty else set()
    sector_max_start_anchor = get_common_series_inception_anchor(sector_series_map, include_labels=sector_labels)

    sector_performance_df, sector_end_date = build_labelled_performance_df(
        source_df=sectors_df,
        label_col="sector",
        fx_values_df=fx_values_df,
        period_key=sector_period,
        neutral_currency=sector_neutral_currency,
        preferred_series_map=chart_series_map,
        max_start_anchor=sector_max_start_anchor,
    )

    toolbar_wrap_cols = st.columns([4.9, 1.5])
    with toolbar_wrap_cols[0]:
        toolbar_cols = st.columns([1.38, 0.28, 1.14])

        with toolbar_cols[0]:
            st.markdown('<div class="toolbar-label">Period:</div>', unsafe_allow_html=True)
            sector_period = st.segmented_control(
                label="Period",
                options=sector_period_options,
                default=st.session_state.get("sector_period_toolbar", sector_period),
                key="sector_period_toolbar",
                label_visibility="collapsed",
            ) or "YTD"

        with toolbar_cols[1]:
            st.markdown('<div class="toolbar-label">Currency:</div>', unsafe_allow_html=True)
            sector_neutral_currency = st.selectbox(
                "Currency",
                options=GEO_NEUTRAL_CURRENCIES,
                index=GEO_NEUTRAL_CURRENCIES.index(st.session_state.get("sector_neutral_currency_toolbar", "USD")),
                key="sector_neutral_currency_toolbar",
                label_visibility="collapsed",
            )

    sector_series_map, sector_series_end_date = build_labelled_series_map(
        source_df=sectors_df,
        label_col="sector",
        fx_values_df=fx_values_df,
        neutral_currency=sector_neutral_currency,
        preferred_series_map=chart_series_map,
    )
    sector_labels = {str(label).strip() for label in sectors_df["sector"].astype(str).tolist()} if not sectors_df.empty else set()
    sector_max_start_anchor = get_common_series_inception_anchor(sector_series_map, include_labels=sector_labels)
    sector_performance_df, sector_end_date = build_labelled_performance_df(
        source_df=sectors_df,
        label_col="sector",
        fx_values_df=fx_values_df,
        period_key=sector_period,
        neutral_currency=sector_neutral_currency,
        preferred_series_map=chart_series_map,
        max_start_anchor=sector_max_start_anchor,
    )

    st.markdown('<div class="page-section-title">Sector performance</div>', unsafe_allow_html=True)
    sector_start_text = (
        (
            pd.Timestamp(sector_max_start_anchor).strftime("%d/%m/%Y")
            if sector_period == "MAX" and sector_max_start_anchor is not None
            else get_geo_period_start_anchor(sector_end_date, sector_period, fx_values_df, sector_neutral_currency).strftime("%d/%m/%Y")
        )
        if sector_end_date is not None
        else "-"
    )
    sector_end_text = sector_end_date.strftime("%d/%m/%Y") if sector_end_date is not None else "-"
    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Sector performance in {sector_neutral_currency} from <b>{sector_start_text}</b> to <b>{sector_end_text}</b>
        </div>
        """,
        unsafe_allow_html=True,
    )
    if sector_performance_df.empty:
        st.info("No sector data were available from the sectors sheet tickers for the selected period.")
    else:
        st.markdown(build_label_tiles_html(sector_performance_df), unsafe_allow_html=True)

    st.markdown('<div class="page-section-title">Patchwork quilt</div>', unsafe_allow_html=True)
    sector_patchwork_df, sector_patchwork_legend_df, sector_patchwork_years = build_country_patchwork_quilt(
        country_series_map=sector_series_map,
        end_date=sector_series_end_date,
        years_back=10,
    )
    sector_patchwork_start_text = f"31/12/{sector_patchwork_years[0]-1}" if sector_patchwork_years else "-"
    sector_patchwork_end_text = f"31/12/{sector_patchwork_years[-1]}" if sector_patchwork_years else "-"
    sector_patchwork_latest_text = pd.Timestamp(sector_end_date).strftime("%d/%m/%Y") if sector_end_date is not None else "-"
    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Calendar-year sector rankings in {sector_neutral_currency} from <b>{sector_patchwork_start_text}</b> to <b>{sector_patchwork_end_text}</b>. The <b>Whole period</b> column shows annualised returns across that full window, and the <b>YTD</b> column runs from <b>31/12/{pd.Timestamp(sector_end_date).year - 1 if sector_end_date is not None else '-'}</b> to <b>{sector_patchwork_latest_text}</b>. Only sectors with a full 10-year calendar history are included.
        </div>
        """,
        unsafe_allow_html=True,
    )
    if sector_patchwork_df.empty:
        st.info("No sectors had a full 10-year calendar history in the selected currency.")
    else:
        st.markdown(build_country_patchwork_html(sector_patchwork_df, sector_patchwork_legend_df, sector_patchwork_years), unsafe_allow_html=True)

    st.markdown(
        '<div class="methodology-text">This tab shows sector performance using the tickers on the workbook <b>sectors</b> sheet where <b>Available = Yes</b>. '
        'Yahoo Finance prices are converted into the selected neutral currency where quote-currency data are available. '
        f'The <b>MAX</b> option uses the common inception date across the available sector ETFs, currently <b>{pd.Timestamp(sector_max_start_anchor).strftime("%d/%m/%Y") if sector_max_start_anchor is not None else "-"}</b>. '
        'The patchwork quilt ranks each included sector by calendar-year return, top to bottom within each year, and shows only sectors with a full 10-year history across the displayed quilt window.</div>',
        unsafe_allow_html=True,
    )

else:
    global_yield_curve_df, global_yield_summary, global_yield_preview = build_global_yield_curve_diagnostics(
        WORLD_GOVERNMENT_BONDS_BASE_URL,
        WORLD_GOVERNMENT_BONDS_COUNTRIES,
    )
    st.markdown('<div class="page-section-title">UK yield curve</div>', unsafe_allow_html=True)
    selected_yield_series = st.multiselect(
        "Curve series",
        options=["Nominal", "Real", "Breakeven inflation"],
        default=st.session_state.get("yield_selected_series", ["Nominal", "Real", "Breakeven inflation"]),
        key="yield_selected_series",
    )
    boe_fetch_ok = (
        not boe_yield_summary.empty
        and "Fetch status" in boe_yield_summary["metric"].values
        and boe_yield_summary.loc[boe_yield_summary["metric"] == "Fetch status", "value"].astype(str).iloc[0] == "OK"
    )
    boe_month_end_fetch_ok = (
        not boe_month_end_yield_summary.empty
        and "Fetch status" in boe_month_end_yield_summary["metric"].values
        and boe_month_end_yield_summary.loc[boe_month_end_yield_summary["metric"] == "Fetch status", "value"].astype(str).iloc[0] == "OK"
    )
    overlay_options = build_uk_yield_curve_overlay_options(yield_curve_history_df) if boe_month_end_fetch_ok else []
    selected_overlay_defs: list[dict[str, object]] = []
    if overlay_options:
        st.markdown('<div class="toolbar-label" style="margin-top:4px;">Historical month-end overlays:</div>', unsafe_allow_html=True)
        overlay_cols = st.columns(len(overlay_options))
        for idx, option in enumerate(overlay_options):
            with overlay_cols[idx]:
                is_selected = st.checkbox(str(option["label"]), key=f"yield_overlay_{option['key']}")
                if is_selected:
                    selected_overlay_defs.append(option)
    elif not boe_month_end_fetch_ok:
        st.info("Historical BOE month-end curves could not be loaded, so only the latest UK curves are shown.")

    historical_curve_frames = []
    for sort_idx, option in enumerate(selected_overlay_defs, start=1):
        curve_date = pd.Timestamp(option["curve_date"]).normalize()
        selected_points = yield_curve_history_df[
            pd.to_datetime(yield_curve_history_df["curve_date"], errors="coerce").dt.normalize() == curve_date
        ][["maturity_years", "yield_percent", "curve_type", "curve_date"]].copy()
        if selected_points.empty:
            continue
        selected_points["snapshot_label"] = f"{option['label']} ({curve_date.strftime('%d/%m/%Y')})"
        selected_points["snapshot_sort"] = sort_idx
        historical_curve_frames.append(selected_points)

    historical_curve_display_df = (
        pd.concat(historical_curve_frames, ignore_index=True)
        if historical_curve_frames
        else pd.DataFrame(columns=["maturity_years", "yield_percent", "curve_type", "curve_date", "snapshot_label", "snapshot_sort"])
    )
    yield_curve_display_df = build_yield_curve_display_df(yield_curve_df, selected_yield_series, historical_curve_display_df)
    has_curve_points = not yield_curve_display_df.empty

    if not boe_fetch_ok and not has_curve_points:
        st.warning("Bank of England yield-curve data could not be loaded. See diagnostics for details.")
    else:
        latest_dates = (
            yield_curve_df.groupby("curve_type")["curve_date"]
            .max()
            .dropna()
            .sort_index()
        ) if not yield_curve_df.empty else pd.Series(dtype="datetime64[ns]")
        latest_text = " | ".join(
            [f"{curve}: {pd.Timestamp(curve_date).strftime('%d/%m/%Y')}" for curve, curve_date in latest_dates.items()]
        ) if not latest_dates.empty else ""
        overlay_text = " | ".join(
            f"{option['label']}: {pd.Timestamp(option['curve_date']).strftime('%d/%m/%Y')}" for option in selected_overlay_defs
        ) if selected_overlay_defs else ""
        meta_parts = []
        if latest_text:
            meta_parts.append(f"Latest curve dates: <b>{latest_text}</b>")
        if overlay_text:
            meta_parts.append(f"Historical overlays: <b>{overlay_text}</b>")
        if meta_parts:
            st.markdown(
                f'<div class="toolbar-meta" style="text-align:left; padding-top:0; margin-bottom:10px;">{" | ".join(meta_parts)}</div>',
                unsafe_allow_html=True,
            )
        if has_curve_points:
            st.altair_chart(build_yield_curve_chart(yield_curve_display_df), width="stretch")
        else:
            st.info("BOE fetch succeeded but no curve points were available to plot. See diagnostics for details.")

    st.markdown('<div class="page-section-title">UK historical yield</div>', unsafe_allow_html=True)
    uk_historical_yield_df = build_uk_historical_yield_df(yield_curve_history_df) if boe_month_end_fetch_ok else pd.DataFrame()
    default_uk_historical_series = [
        "Nominal 10Y",
        "Real 10Y",
        "BEI 10Y",
    ]
    available_uk_historical_series = uk_historical_yield_df["series_label"].dropna().drop_duplicates().tolist() if not uk_historical_yield_df.empty else []
    hist_cols = st.columns([1.6, 4.4])
    with hist_cols[0]:
        st.markdown('<div class="toolbar-label">Period:</div>', unsafe_allow_html=True)
        uk_hist_period = st.segmented_control(
            label="UK historical yield period",
            options=UK_HISTORICAL_YIELD_PERIODS,
            default=st.session_state.get("uk_hist_yield_period_toolbar", "MAX"),
            key="uk_hist_yield_period_toolbar",
            label_visibility="collapsed",
        ) or "MAX"
    with hist_cols[1]:
        st.markdown('<div class="toolbar-label">Series:</div>', unsafe_allow_html=True)
        uk_hist_series = st.multiselect(
            "UK historical yield series",
            options=available_uk_historical_series,
            default=[label for label in default_uk_historical_series if label in available_uk_historical_series],
            key="uk_hist_yield_series_toolbar",
            label_visibility="collapsed",
        )

    if uk_historical_yield_df.empty:
        st.info("No BOE month-end history was available for the UK historical yield chart.")
    else:
        uk_hist_end_date = pd.Timestamp(uk_historical_yield_df["curve_date"].max()).normalize()
        uk_hist_min_date = pd.Timestamp(uk_historical_yield_df["curve_date"].min()).normalize()
        uk_hist_start_date = get_uk_historical_yield_start_date(uk_hist_period, uk_hist_end_date, uk_hist_min_date)
        uk_historical_window_df = uk_historical_yield_df[
            (pd.to_datetime(uk_historical_yield_df["curve_date"], errors="coerce").dt.normalize() >= uk_hist_start_date)
            & (pd.to_datetime(uk_historical_yield_df["curve_date"], errors="coerce").dt.normalize() <= uk_hist_end_date)
        ].copy()
        chart_series = uk_hist_series or available_uk_historical_series
        note = (
            f"Month-end UK nominal and real yields from <b>{uk_hist_start_date.strftime('%d/%m/%Y')}</b> to <b>{uk_hist_end_date.strftime('%d/%m/%Y')}</b>. "
            "The shortest real and BEI line uses the earliest available BOE real-curve maturity, which is <b>2.5Y</b> in the month-end archive."
        )
        st.markdown(f'<div class="snapshot-toolbar-note">{note}</div>', unsafe_allow_html=True)
        if uk_historical_window_df.empty or not chart_series:
            st.info("No UK historical yield series were available for the selected period and filters.")
        else:
            st.altair_chart(build_uk_historical_yield_chart(uk_historical_window_df, chart_series, uk_hist_period), width="stretch")

    st.markdown('<div class="page-section-title">UK historical term spreads</div>', unsafe_allow_html=True)
    uk_term_spread_df = build_uk_term_spread_df(yield_curve_history_df) if boe_month_end_fetch_ok else pd.DataFrame()
    default_uk_term_spreads = ["Nominal 10-2Y", "Nominal 2-0.5Y"]
    available_uk_term_spreads = uk_term_spread_df["series_label"].dropna().drop_duplicates().tolist() if not uk_term_spread_df.empty else []
    spread_cols = st.columns([1.6, 4.4])
    with spread_cols[0]:
        st.markdown('<div class="toolbar-label">Period:</div>', unsafe_allow_html=True)
        uk_term_spread_period = st.segmented_control(
            label="UK term spread period",
            options=UK_HISTORICAL_YIELD_PERIODS,
            default=st.session_state.get("uk_term_spread_period_toolbar", "MAX"),
            key="uk_term_spread_period_toolbar",
            label_visibility="collapsed",
        ) or "MAX"
    with spread_cols[1]:
        st.markdown('<div class="toolbar-label">Spreads:</div>', unsafe_allow_html=True)
        uk_term_spread_selection = st.multiselect(
            "UK term spread series",
            options=available_uk_term_spreads,
            default=[label for label in default_uk_term_spreads if label in available_uk_term_spreads],
            key="uk_term_spread_series_toolbar",
            label_visibility="collapsed",
        )

    if uk_term_spread_df.empty:
        st.info("No BOE month-end history was available for the UK term-spread chart.")
    else:
        uk_term_end_date = pd.Timestamp(uk_term_spread_df["curve_date"].max()).normalize()
        uk_term_min_date = pd.Timestamp(uk_term_spread_df["curve_date"].min()).normalize()
        uk_term_start_date = get_uk_historical_yield_start_date(uk_term_spread_period, uk_term_end_date, uk_term_min_date)
        uk_term_window_df = uk_term_spread_df[
            (pd.to_datetime(uk_term_spread_df["curve_date"], errors="coerce").dt.normalize() >= uk_term_start_date)
            & (pd.to_datetime(uk_term_spread_df["curve_date"], errors="coerce").dt.normalize() <= uk_term_end_date)
        ].copy()
        spread_series = uk_term_spread_selection or available_uk_term_spreads
        st.markdown(
            f'<div class="snapshot-toolbar-note">Month-end UK term spreads from <b>{uk_term_start_date.strftime("%d/%m/%Y")}</b> to <b>{uk_term_end_date.strftime("%d/%m/%Y")}</b>.</div>',
            unsafe_allow_html=True,
        )
        if uk_term_window_df.empty or not spread_series:
            st.info("No UK term-spread series were available for the selected period and filters.")
        else:
            st.altair_chart(build_uk_term_spread_chart(uk_term_window_df, spread_series, uk_term_spread_period), width="stretch")

    st.markdown('<div class="page-section-title">Global yield curves</div>', unsafe_allow_html=True)
    global_latest_curve_date = (
        pd.to_datetime(global_yield_curve_df["curve_date"], errors="coerce").max()
        if not global_yield_curve_df.empty and "curve_date" in global_yield_curve_df.columns
        else pd.NaT
    )
    global_latest_curve_text = (
        pd.Timestamp(global_latest_curve_date).strftime("%d/%m/%Y")
        if pd.notna(global_latest_curve_date)
        else "unknown date"
    )
    st.markdown(
        '<div class="diag-note" style="margin-top:-4px; margin-bottom:10px;">'
        f'WorldGovernmentBonds nominal curves, latest dated <b>{global_latest_curve_text}</b>; methodology may differ from the UK chart above.'
        '</div>',
        unsafe_allow_html=True,
    )
    global_fetch_ok = (
        not global_yield_summary.empty
        and "Fetch status" in global_yield_summary["metric"].values
        and global_yield_summary.loc[global_yield_summary["metric"] == "Fetch status", "value"].astype(str).iloc[0] == "OK"
    )
    if not global_fetch_ok and global_yield_curve_df.empty:
        st.warning("WorldGovernmentBonds yield-curve data could not be loaded. See diagnostics for details.")
    elif global_yield_curve_df.empty:
        st.info("WorldGovernmentBonds fetch succeeded but no country curve points were available to plot.")
    else:
        st.altair_chart(build_global_yield_curve_chart(global_yield_curve_df), width="stretch")

    fx_values_df = fetch_fx_value_series()
    fx_period = st.session_state.get("rates_fx_period_toolbar", "YTD")
    currency_matrix_df, currency_end_date = build_currency_performance_matrix(fx_values_df, fx_period)
    fx_period_options = get_fx_period_options(currency_end_date, fx_values_df)

    st.markdown('<div class="page-section-title">Currency performance</div>', unsafe_allow_html=True)
    toolbar_wrap_cols = st.columns([5.1, 1.3])
    with toolbar_wrap_cols[0]:
        toolbar_cols = st.columns([1.55, 1.15])
        with toolbar_cols[0]:
            st.markdown('<div class="toolbar-label">Period:</div>', unsafe_allow_html=True)
            fx_period = st.segmented_control(
                label="FX period",
                options=fx_period_options,
                default=st.session_state.get("rates_fx_period_toolbar", fx_period),
                key="rates_fx_period_toolbar",
                label_visibility="collapsed",
            ) or "YTD"

    currency_matrix_df, currency_end_date = build_currency_performance_matrix(fx_values_df, fx_period)
    currency_start_text = (
        get_fx_period_start_anchor(currency_end_date, fx_period, fx_values_df).strftime("%d/%m/%Y")
        if currency_end_date is not None
        else "-"
    )
    end_date_text = currency_end_date.strftime("%d/%m/%Y") if currency_end_date is not None else "-"
    st.markdown(
        f"""
        <div class="snapshot-toolbar-note">
            Daily FX performance matrix from <b>{currency_start_text}</b> to <b>{end_date_text}</b>
        </div>
        """,
        unsafe_allow_html=True,
    )
    if currency_matrix_df.empty:
        st.info("No FX data were available from yfinance for the selected matrix.")
    else:
        st.markdown(build_currency_matrix_html(currency_matrix_df), unsafe_allow_html=True)

    yield_note = (
        "Latest nominal and real UK spot curves are fetched from the Bank of England yield-curve archive. "
        "The app reads the latest non-blank dated row from the '4. spot curve' sheet in the current-month nominal and real daily workbooks. "
        "Historical UK month-end overlays are fetched from the Bank of England nominal and real month-end archives and can add last month end, last quarter end, 6M ago, 9M ago, and 1Y ago snapshots. "
        "Where DividendData provides shorter-maturity index-linked gilts than the Bank of England real curve, those short-end gilt real yields are prepended to extend the real curve only below the first BOE real maturity. "
        "Breakeven inflation is then calculated point-by-point as ((1+nominal yield)/(1+real yield))-1 on the maturities where both nominal and real yields are available, both for the latest curve and any selected historical month-end overlays. "
        "The global chart fetches the latest nominal yield-curve tables from WorldGovernmentBonds country pages and plots the 'Last' annualised yield column for the ten largest government bond markets. "
        "The FX matrix shows daily row-versus-column currency performance from yfinance over the selected period."
    )
    st.markdown(f'<div class="methodology-text">{yield_note}</div>', unsafe_allow_html=True)


# =====================================
# FOOTER
# =====================================
if POWERED_BY_FILE.exists():
    st.markdown(
        f"""
        <div class="footer-bar">
            <div class="footer-logo">
                <img src="data:image/png;base64,{img_to_base64(str(POWERED_BY_FILE), POWERED_BY_FILE.stat().st_mtime)}">
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


# =====================================
# DIAGNOSTICS
# =====================================
diag_label = (
    "▼ Diagnostics and underlying data"
    if st.session_state["show_diagnostics"]
    else "► Diagnostics and underlying data"
)

if st.button(diag_label, key="diagnostics_toggle_button", use_container_width=True):
    st.session_state["show_diagnostics"] = not st.session_state["show_diagnostics"]

if st.session_state["show_diagnostics"]:
    if global_yield_summary.empty:
        global_yield_curve_df, global_yield_summary, global_yield_preview = build_global_yield_curve_diagnostics(
            WORLD_GOVERNMENT_BONDS_BASE_URL,
            WORLD_GOVERNMENT_BONDS_COUNTRIES,
        )
    generic_end_date = get_dashboard_end_date(
        stitched_series_map=stitched_series_map,
        live_diag=live_diag,
        inflation_levels=inflation_levels,
        is_real_mode=False,
    )
    dashboard_anchor_diag = build_return_anchor_table(
        stitched_series_map,
        generic_end_date,
        ["YTD", "1Y", "3Y", "5Y", "10Y", "20Y", "Period"],
        whole_period_start,
    )
    chart_anchor_diag = build_return_anchor_table(
        chart_series_map,
        generic_end_date,
        ["YTD", "1Y", "3Y", "5Y", "10Y", "20Y", "Period"],
        whole_period_start,
    )
    inflation_anchor_diag = (
        build_return_anchor_table(
            {"UK inflation": inflation_levels},
            generic_end_date,
            ["YTD", "1Y", "3Y", "5Y", "10Y", "20Y", "Period"],
            whole_period_start,
        )
        if inflation_levels is not None and not inflation_levels.dropna().empty
        else pd.DataFrame()
    )

    with st.container(border=True):
        st.markdown('<div class="diag-title">Live fund stitching diagnostics</div>', unsafe_allow_html=True)
        st.markdown(
            f"""
            <div class="diag-note">
                Index history is always preferred where available. Live yfinance history is only used after index history ends for the dashboard and tables.
                The growth of wealth chart uses a separate cached series that prioritises daily live-fund history where available and stitches monthly index history before that to extend the chart backwards.
                Inflation source: <b>{inflation_source_note}</b>. Generic dashboard end date: <b>{generic_end_date.strftime("%d/%m/%Y")}</b>.
                Common inception date used for period table: <b>{common_inception_text}</b>.
            </div>
            """,
            unsafe_allow_html=True,
        )

        if inflation_debug_message:
            st.info(inflation_debug_message)

        tabs = st.tabs(["Overview", "Dashboard stitching", "Chart series", "Returns", "Inflation", "Yield curves", "Mapping & prices"])

        with tabs[0]:
            c1, c2, c3, c4, c5, c6 = st.columns(6)
            c1.metric("Mapped rows", int(len(mapping)))
            c2.metric("Unique asset classes", int(mapping["asset_class"].nunique()))
            c3.metric("Monthly index series", int(len(monthly_levels)))
            c4.metric("Dashboard series", int(len(stitched_series_map)))
            c5.metric("Chart series", int(len(chart_series_map)))
            c6.metric("Live tickers returned", int(len(live_prices.columns)) if not live_prices.empty else 0)

            st.subheader("Asset coverage summary")
            st.dataframe(prepare_dataframe_for_display(asset_coverage_diag), width="stretch", hide_index=True)
            st.download_button(
                label="Download asset coverage (CSV)",
                data=dataframe_to_csv_download(asset_coverage_diag),
                file_name="asset_coverage_diagnostics.csv",
                mime="text/csv",
                key="download_asset_coverage_csv",
            )

            missing_assets = asset_coverage_diag[
                (asset_coverage_diag["index_points"] == 0)
                | (asset_coverage_diag["dashboard_points"] == 0)
                | (asset_coverage_diag["chart_points"] == 0)
            ].copy()
            st.subheader("Coverage gaps")
            if missing_assets.empty:
                st.write("No major coverage gaps detected in mapped assets.")
            else:
                st.dataframe(prepare_dataframe_for_display(missing_assets), width="stretch", hide_index=True)

        with tabs[1]:
            if not live_diag.empty:
                summary = (
                    live_diag["series_type"]
                    .fillna("unknown")
                    .value_counts(dropna=False)
                    .rename_axis("series_type")
                    .reset_index(name="count")
                )
                counts = summary.set_index("series_type")["count"].to_dict()

                c1, c2, c3, c4 = st.columns(4)
                c1.metric("Stitched", int(counts.get("stitched", 0)))
                c2.metric("Index only", int(counts.get("index_only", 0)))
                c3.metric("Live only", int(counts.get("live_only", 0)))
                c4.metric("Missing", int(counts.get("missing", 0)))

                st.subheader("Summary")
                st.dataframe(prepare_dataframe_for_display(summary), width="stretch", hide_index=True)

                diag_show = format_diagnostic_table(live_diag)
                preferred_cols = [
                    "asset_class",
                    "live_fund_primary",
                    "live_fund_secondary",
                    "selected_ticker",
                    "selected_source",
                    "series_type",
                    "index_last_date",
                    "live_first_date",
                    "live_last_date",
                    "stitch_anchor_date",
                    "note",
                ]
                display_cols = [c for c in preferred_cols if c in diag_show.columns]
                st.subheader("Detailed diagnostics")
                st.dataframe(prepare_dataframe_for_display(diag_show[display_cols]), width="stretch", hide_index=True)

                st.download_button(
                    label="Download live diagnostics (CSV)",
                    data=dataframe_to_csv_download(live_diag),
                    file_name="live_fund_stitch_diagnostics.csv",
                    mime="text/csv",
                    use_container_width=False,
                    key="download_live_diagnostics_csv",
                )
            else:
                st.write("No live diagnostics available.")

        with tabs[2]:
            if not chart_diag.empty:
                summary = (
                    chart_diag["series_type"]
                    .fillna("unknown")
                    .value_counts(dropna=False)
                    .rename_axis("series_type")
                    .reset_index(name="count")
                )
                st.subheader("Summary")
                st.dataframe(prepare_dataframe_for_display(summary), width="stretch", hide_index=True)

                chart_show = format_diagnostic_table(chart_diag)
                st.subheader("Chart series diagnostics")
                st.dataframe(prepare_dataframe_for_display(chart_show), width="stretch", hide_index=True)
                st.download_button(
                    label="Download chart diagnostics (CSV)",
                    data=dataframe_to_csv_download(chart_diag),
                    file_name="chart_series_diagnostics.csv",
                    mime="text/csv",
                    key="download_chart_diagnostics_csv",
                )
            else:
                st.write("No chart series diagnostics available.")

        with tabs[3]:
            returns_tabs = st.tabs(["Dashboard returns", "Charts returns", "Calendar year returns", "Return anchors"])

            with returns_tabs[0]:
                dashboard_nominal = order_asset_rows(
                    calc_horizon_returns_from_levels(stitched_series_map, generic_end_date, list(DASHBOARD_HORIZONS.keys()))
                )
                st.dataframe(convert_pct_table_for_display(dashboard_nominal), width="stretch", hide_index=True)
                st.download_button(
                    label="Download dashboard returns (CSV)",
                    data=dataframe_to_csv_download(dashboard_nominal),
                    file_name="dashboard_returns_diagnostics.csv",
                    mime="text/csv",
                    key="download_dashboard_returns_diag_csv",
                )

            with returns_tabs[1]:
                charts_nominal = order_asset_rows(
                    merge_return_tables(
                        calc_horizon_returns_from_levels(
                            stitched_series_map,
                            generic_end_date,
                            ["YTD", "1Y", "3Y", "5Y", "7Y", "10Y", "15Y", "20Y", "25Y"],
                        ),
                        calc_whole_period_returns(stitched_series_map, generic_end_date, whole_period_start),
                    )
                )
                st.dataframe(convert_pct_table_for_display(charts_nominal), width="stretch", hide_index=True)
                st.download_button(
                    label="Download charts returns (CSV)",
                    data=dataframe_to_csv_download(charts_nominal),
                    file_name="charts_returns_diagnostics.csv",
                    mime="text/csv",
                    key="download_charts_returns_diag_csv",
                )

            with returns_tabs[2]:
                calendar_diag = build_calendar_year_returns(stitched_series_map, generic_end_date, 10)
                st.dataframe(convert_pct_table_for_display(calendar_diag), width="stretch", hide_index=True)
                st.download_button(
                    label="Download calendar returns (CSV)",
                    data=dataframe_to_csv_download(calendar_diag),
                    file_name="calendar_year_returns_diagnostics.csv",
                    mime="text/csv",
                    key="download_calendar_returns_diag_csv",
                )

            with returns_tabs[3]:
                anchor_tabs = st.tabs(["Dashboard anchor dates", "Chart anchor dates"])

                with anchor_tabs[0]:
                    st.dataframe(prepare_dataframe_for_display(dashboard_anchor_diag), width="stretch", hide_index=True)
                    st.download_button(
                        label="Download dashboard anchors (CSV)",
                        data=dataframe_to_csv_download(dashboard_anchor_diag),
                        file_name="dashboard_return_anchor_diagnostics.csv",
                        mime="text/csv",
                        key="download_dashboard_anchor_diag_csv",
                    )

                with anchor_tabs[1]:
                    st.dataframe(prepare_dataframe_for_display(chart_anchor_diag), width="stretch", hide_index=True)
                    st.download_button(
                        label="Download chart anchors (CSV)",
                        data=dataframe_to_csv_download(chart_anchor_diag),
                        file_name="chart_return_anchor_diagnostics.csv",
                        mime="text/csv",
                        key="download_chart_anchor_diag_csv",
                    )

        with tabs[4]:
            c1, c2 = st.columns(2)

            with c1:
                if inflation_levels is not None and not inflation_levels.dropna().empty:
                    infl_diag_table = order_asset_rows(
                        merge_return_tables(
                            calc_horizon_returns_from_levels(
                                {"UK inflation": inflation_levels},
                                generic_end_date,
                                ["YTD", "1Y", "3Y", "5Y", "7Y", "10Y", "15Y", "20Y", "25Y"],
                            ),
                            calc_whole_period_returns({"UK inflation": inflation_levels}, generic_end_date, whole_period_start),
                        )
                    )
                    st.subheader("Inflation return table")
                    st.dataframe(convert_pct_table_for_display(infl_diag_table), width="stretch", hide_index=True)
                else:
                    st.write("No inflation return table available.")

            with c2:
                if inflation_levels is not None and not inflation_levels.empty:
                    infl_last = pd.to_datetime(inflation_levels.index.max()).strftime("%d/%m/%Y")
                    jan26_val = inflation_monthly_returns.get(pd.Timestamp("2026-01-31"), np.nan)
                    st.metric("Inflation source", inflation_source_note)
                    st.metric("Inflation last date", infl_last)
                    st.metric("Jan-26 monthly return", "-" if pd.isna(jan26_val) else f"{jan26_val:.4%}")
                else:
                    st.write("No inflation series available.")

            st.subheader("ONS fetch diagnostics")
            st.dataframe(prepare_dataframe_for_display(ons_fetch_summary), width="stretch", hide_index=True)

            if not ons_fetch_preview.empty:
                ons_preview_show = ons_fetch_preview.copy()
                ons_preview_show["value"] = ons_preview_show["value"].map(lambda x: np.nan if pd.isna(x) else round(float(x), 4))
                st.dataframe(prepare_dataframe_for_display(ons_preview_show.tail(24)), width="stretch", hide_index=True)
                st.download_button(
                    label="Download ONS parse preview (CSV)",
                    data=dataframe_to_csv_download(ons_fetch_preview),
                    file_name="ons_cpi_parse_preview.csv",
                    mime="text/csv",
                    key="download_ons_parse_preview_csv",
                )

            if inflation_levels is not None and not inflation_levels.empty:
                levels_df = inflation_levels.reset_index()
                levels_df.columns = ["Date", "Inflation level"]
                levels_df["Date"] = pd.to_datetime(levels_df["Date"]).dt.strftime("%d/%m/%Y")

                monthly_df = inflation_monthly_returns.reset_index()
                monthly_df.columns = ["Date", "UK inflation monthly return"]
                monthly_df["Date"] = pd.to_datetime(monthly_df["Date"]).dt.strftime("%d/%m/%Y")
                monthly_df["UK inflation monthly return"] = monthly_df["UK inflation monthly return"].map(
                    lambda x: np.nan if pd.isna(x) else round(x * 100, 4)
                )

                lower_tabs = st.tabs(["Inflation levels tail", "Inflation monthly tail", "Inflation anchors"])

                with lower_tabs[0]:
                    st.dataframe(prepare_dataframe_for_display(levels_df.tail(24)), width="stretch", hide_index=True)

                with lower_tabs[1]:
                    st.dataframe(prepare_dataframe_for_display(monthly_df.tail(24)), width="stretch", hide_index=True)

                with lower_tabs[2]:
                    st.dataframe(prepare_dataframe_for_display(inflation_anchor_diag), width="stretch", hide_index=True)
            else:
                st.write("No inflation series available.")

        with tabs[5]:
            st.subheader("Bank of England yield-curve diagnostics")
            st.dataframe(prepare_dataframe_for_display(boe_yield_summary), width="stretch", hide_index=True)
            if not boe_yield_preview.empty:
                st.dataframe(prepare_dataframe_for_display(boe_yield_preview), width="stretch", hide_index=True)
                st.download_button(
                    label="Download BOE yield diagnostics (CSV)",
                    data=dataframe_to_csv_download(boe_yield_preview),
                    file_name="boe_yield_curve_diagnostics.csv",
                    mime="text/csv",
                    key="download_boe_yield_diag_csv",
                )
            if not yield_curve_df.empty:
                curve_preview = yield_curve_df.copy()
                curve_preview["curve_date"] = pd.to_datetime(curve_preview["curve_date"], errors="coerce").dt.strftime("%d/%m/%Y").fillna("")
                st.subheader("Latest BOE spot-curve points")
                st.dataframe(prepare_dataframe_for_display(curve_preview), width="stretch", hide_index=True)

            st.subheader("BOE month-end history diagnostics")
            st.dataframe(prepare_dataframe_for_display(boe_month_end_yield_summary), width="stretch", hide_index=True)
            if not boe_month_end_yield_preview.empty:
                st.dataframe(prepare_dataframe_for_display(boe_month_end_yield_preview), width="stretch", hide_index=True)
                st.download_button(
                    label="Download BOE month-end diagnostics (CSV)",
                    data=dataframe_to_csv_download(boe_month_end_yield_preview),
                    file_name="boe_month_end_yield_diagnostics.csv",
                    mime="text/csv",
                    key="download_boe_month_end_yield_diag_csv",
                )
            if not yield_curve_history_df.empty:
                history_preview = yield_curve_history_df.copy()
                history_preview["curve_date"] = pd.to_datetime(history_preview["curve_date"], errors="coerce").dt.strftime("%d/%m/%Y").fillna("")
                st.subheader("BOE month-end spot-curve points")
                st.dataframe(prepare_dataframe_for_display(history_preview.tail(300)), width="stretch", hide_index=True)

            st.subheader("WorldGovernmentBonds diagnostics")
            st.dataframe(prepare_dataframe_for_display(global_yield_summary), width="stretch", hide_index=True)
            if not global_yield_preview.empty:
                st.dataframe(prepare_dataframe_for_display(global_yield_preview), width="stretch", hide_index=True)
                st.download_button(
                    label="Download global yield diagnostics (CSV)",
                    data=dataframe_to_csv_download(global_yield_preview),
                    file_name="global_yield_curve_diagnostics.csv",
                    mime="text/csv",
                    key="download_global_yield_diag_csv",
                )
            if not global_yield_curve_df.empty:
                global_curve_preview = global_yield_curve_df.copy()
                global_curve_preview["curve_date"] = pd.to_datetime(global_curve_preview["curve_date"], errors="coerce").dt.strftime("%d/%m/%Y").fillna("")
                st.subheader("Latest global nominal curve points")
                st.dataframe(prepare_dataframe_for_display(global_curve_preview), width="stretch", hide_index=True)

        with tabs[6]:
            mapping_tabs = st.tabs(["Validated mapping", "Raw mapping", "Live prices"])

            with mapping_tabs[0]:
                st.subheader("Validated mapping health")
                st.dataframe(prepare_dataframe_for_display(mapping_diag), width="stretch", hide_index=True)
                st.download_button(
                    label="Download mapping diagnostics (CSV)",
                    data=dataframe_to_csv_download(mapping_diag),
                    file_name="mapping_diagnostics.csv",
                    mime="text/csv",
                    key="download_mapping_diag_csv",
                )

            with mapping_tabs[1]:
                st.subheader("Mapping table")
                st.dataframe(prepare_dataframe_for_display(mapping), width="stretch", hide_index=True)

            with mapping_tabs[2]:
                st.subheader("Live price coverage")
                if live_price_diag.empty:
                    st.write("No live price history available.")
                else:
                    st.dataframe(prepare_dataframe_for_display(live_price_diag), width="stretch", hide_index=True)
                    st.download_button(
                        label="Download live price coverage (CSV)",
                        data=dataframe_to_csv_download(live_price_diag),
                        file_name="live_price_coverage.csv",
                        mime="text/csv",
                        key="download_live_price_coverage_csv",
                    )
